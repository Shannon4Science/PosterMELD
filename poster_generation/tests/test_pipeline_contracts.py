from copy import deepcopy
import json
from pathlib import Path
import zipfile

import pytest
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE

from src.agents.affiliation_logo_agent import AffiliationLogoAgent
from src.agents.adaptive_column_relayout import AdaptiveColumnRelayoutAgent
from src.agents.background_image_agent import BackgroundImageAgent
from src.agents.block_content_refiner import BlockContentRefiner
from src.agents.block_occupancy_analyzer import BlockOccupancyAnalyzer
from src.agents.block_vlm_reviewer import BlockVLMReviewer
from src.agents.color_agent import ColorAgent
from src.agents.curator import StoryBoardCurator
from src.agents.font_agent import FontAgent
from src.agents.generated_teaser_agent import GeneratedTeaserAgent
from src.agents.header_block_reviewer import HeaderBlockReviewer
from src.agents.header_planner import HeaderPlanner
from src.agents.layout_agent import LayoutAgent
from src.agents.layout_with_balancer import LayoutWithBalancerAgent
from src.agents.micro_layout_refiner import MicroLayoutRefiner
from src.agents.parser import Parser
from src.agents.poster_keypoint_selector import PosterKeypointSelector
from src.agents.renderer import Renderer
from src.agents.section_title_designer import SectionTitleDesigner
from src.agents.standard_template_preselector import StandardTemplatePreselector
from src.agents.template_capacity_planner import TemplateCapacityPlanner
from src.agents.template_block_planner import TemplateBlockPlanner
from src.agents.vlm_layout_reviewer import VLMLayoutReviewer
from src.agents.visual_asset_agent import VisualAssetAgent
from src.agents.visual_legibility_reviewer import VisualLegibilityReviewer
from src.config.poster_config import load_config
from src.layout.template_selector import TemplateSelector
from src.state.poster_state import create_state, _get_model_config
from src.template_extraction.block_template_registry import (
    build_runtime_template,
    get_block_template_info,
    list_block_template_ids,
    load_block_template_layout,
)
from src.template_extraction.extract_templates import build_template
from src.template_extraction.registry import list_extracted_template_ids, load_extracted_template
from src.tools.image_api import ImageTools
from src.tools.layout_api import LayoutTemplates
from src.tools.pptx_api import PPTXDirector
from src.utils.text_cleanup import normalize_text_for_poster, normalize_title_for_poster
from utils.langgraph_utils import LangGraphAgent
from src.workflow.pipeline import (
    _build_final_gate_refinement_occupancy,
    _final_artifact_failures,
    _run_final_quality_gate,
    _section_geometry_issues,
    resolve_poster_dimensions,
)
from src.tools.mineru_api import MinerUClient, MinerUExtraction


def test_parser_visual_assets_registry_matches_images_tables():
    parser = Parser.__new__(Parser)
    figures = {
        "1": {
            "caption": "Figure 1",
            "path": "/tmp/figure-1.png",
            "aspect": 1.5,
        }
    }
    tables = {
        "2": {
            "caption": "Table 2",
            "path": "/tmp/table-2.png",
            "aspect": 2.0,
        }
    }

    visual_assets = parser._build_visual_registry(figures, tables)

    assert visual_assets["figure_1"]["source_path"] == figures["1"]["path"]
    assert visual_assets["figure_1"]["asset_type"] == "figure"
    assert visual_assets["table_2"]["source_path"] == tables["2"]["path"]
    assert visual_assets["table_2"]["asset_type"] == "table"


def test_parser_extracts_mineru_assets_from_content_list(tmp_path):
    extract_dir = tmp_path / "content" / "mineru_raw"
    images_dir = extract_dir / "images"
    images_dir.mkdir(parents=True)
    Image.new("RGB", (120, 60), "white").save(images_dir / "figure-a.jpg")
    Image.new("RGB", (200, 80), "white").save(images_dir / "table-a.jpg")
    content_list_path = extract_dir / "paper_content_list.json"
    items = [
        {
            "type": "image",
            "img_path": "images/figure-a.jpg",
            "image_caption": ["Figure 1: Architecture overview."],
            "page_idx": 0,
            "bbox": [1, 2, 3, 4],
        },
        {
            "type": "table",
            "img_path": "images/table-a.jpg",
            "table_caption": ["Table 1: Main results."],
            "page_idx": 1,
        },
    ]
    content_list_path.write_text(json.dumps(items), encoding="utf-8")
    extraction = MinerUExtraction(
        raw_text="# Paper",
        extract_dir=extract_dir,
        zip_path=extract_dir / "mineru_result.zip",
        content_list_path=content_list_path,
        content_items=items,
        report={"backend": "mineru"},
    )
    parser = Parser.__new__(Parser)
    parser.name = "parser"

    figures, tables = parser._extract_mineru_assets(extraction, tmp_path / "assets")

    assert figures["1"]["caption"] == "Figure 1: Architecture overview."
    assert figures["1"]["aspect"] == 2.0
    assert Path(figures["1"]["path"]).exists()
    assert tables["1"]["caption"] == "Table 1: Main results."
    assert tables["1"]["aspect"] == 2.5
    assert Path(tables["1"]["path"]).exists()


def test_parser_renders_mineru_asset_from_source_pdf_bbox_at_high_resolution(tmp_path):
    import fitz

    pdf_path = tmp_path / "paper.pdf"
    document = fitz.open()
    page = document.new_page(width=600, height=800)
    page.draw_rect(fitz.Rect(60, 80, 300, 240), color=(0, 0, 0), width=1)
    page.insert_text((80, 130), "Readable vector figure label", fontsize=18)
    document.save(pdf_path)
    document.close()

    extract_dir = tmp_path / "content" / "mineru_raw"
    images_dir = extract_dir / "images"
    images_dir.mkdir(parents=True)
    Image.new("RGB", (80, 40), "white").save(images_dir / "low-res.jpg")
    items = [{
        "type": "image",
        "img_path": "images/low-res.jpg",
        "image_caption": ["Figure 1: Vector source."],
        "page_idx": 0,
        "bbox": [100, 100, 500, 300],
    }]
    extraction = MinerUExtraction(
        raw_text="# Paper",
        extract_dir=extract_dir,
        zip_path=extract_dir / "mineru_result.zip",
        content_list_path=None,
        content_items=items,
        pdf_path=pdf_path,
        report={"backend": "mineru"},
    )
    parser = Parser.__new__(Parser)
    parser.name = "parser"

    figures, _ = parser._extract_mineru_assets(extraction, tmp_path / "assets")

    assert figures["1"]["extraction_method"] == "pdf_bbox_render"
    assert figures["1"]["width"] >= 1500
    assert figures["1"]["width"] > 80
    assert figures["1"]["height"] > 40


def test_mineru_client_uploads_polls_and_unpacks_zip(tmp_path, monkeypatch):
    pdf_path = tmp_path / "paper.pdf"
    pdf_path.write_bytes(b"%PDF-1.4\n")
    fixture_dir = tmp_path / "fixture"
    (fixture_dir / "images").mkdir(parents=True)
    Image.new("RGB", (80, 40), "white").save(fixture_dir / "images" / "figure-a.png")
    (fixture_dir / "full.md").write_text("# Parsed Paper\n\nBody.", encoding="utf-8")
    (fixture_dir / "paper_content_list.json").write_text(
        json.dumps([{"type": "image", "img_path": "images/figure-a.png", "image_caption": ["Figure"]}]),
        encoding="utf-8",
    )
    zip_path = tmp_path / "mineru_fixture.zip"
    with zipfile.ZipFile(zip_path, "w") as zf:
        zf.write(fixture_dir / "full.md", "full.md")
        zf.write(fixture_dir / "paper_content_list.json", "paper_content_list.json")
        zf.write(fixture_dir / "images" / "figure-a.png", "images/figure-a.png")

    calls = {"post": 0, "put": 0, "get": 0}

    class FakeResponse:
        def __init__(self, payload=None, content=b"", status_code=200):
            self._payload = payload or {}
            self.content = content
            self.status_code = status_code
            self.text = json.dumps(self._payload)

        def json(self):
            return self._payload

        def raise_for_status(self):
            if self.status_code >= 400:
                raise RuntimeError(self.text)

    def fake_post(url, headers=None, json=None, timeout=None):
        calls["post"] += 1
        assert headers["Authorization"] == "Bearer test-key"
        assert json["model_version"] == "vlm"
        return FakeResponse({"code": 0, "data": {"batch_id": "batch-1", "file_urls": [{"url": "https://upload"}]}})

    def fake_put(url, data=None, headers=None, timeout=None):
        calls["put"] += 1
        assert url == "https://upload"
        assert headers is None
        return FakeResponse({})

    def fake_get(url, headers=None, timeout=None):
        calls["get"] += 1
        if "extract-results" in url:
            return FakeResponse({
                "code": 0,
                "data": {
                    "extract_result": [
                        {"file_name": "paper.pdf", "state": "done", "full_zip_url": "https://download/result.zip"}
                    ]
                },
            })
        return FakeResponse(content=zip_path.read_bytes())

    monkeypatch.setattr("src.tools.mineru_api.requests.post", fake_post)
    monkeypatch.setattr("src.tools.mineru_api.requests.put", fake_put)
    monkeypatch.setattr("src.tools.mineru_api.requests.get", fake_get)

    client = MinerUClient("test-key", base_url="https://mineru.net", poll_interval=0.01)
    extraction = client.parse_pdf(pdf_path, tmp_path / "content")

    assert extraction.raw_text.startswith("# Parsed Paper")
    assert extraction.content_list_path is not None
    assert extraction.content_items[0]["type"] == "image"
    assert calls == {"post": 1, "put": 1, "get": 2}


def test_parser_falls_back_to_marker_when_mineru_fails(tmp_path, monkeypatch):
    parser = Parser.__new__(Parser)
    parser.name = "parser"
    parser.config_data = {
        "pdf_processing": {
            "backend": "mineru",
            "fallback_backend": "marker",
            "mineru": {"model_version": "vlm"},
        }
    }
    content_dir = tmp_path / "content"
    content_dir.mkdir()

    monkeypatch.setattr("src.agents.parser.MinerUClient.from_env", lambda config: (_ for _ in ()).throw(RuntimeError("mineru down")))
    monkeypatch.setattr(parser, "_extract_raw_text_with_marker", lambda pdf_path, content_dir: ("marker raw", ("marker",)))

    raw_text, raw_result = parser._extract_raw_text(str(tmp_path / "paper.pdf"), content_dir)
    report = json.loads((content_dir / "mineru_report.json").read_text(encoding="utf-8"))

    assert raw_text == "marker raw"
    assert raw_result == ("marker",)
    assert report["fallback_used"] is True
    assert report["fallback_backend"] == "marker"
    assert "api_key" not in report


def test_parser_extracts_affiliations_from_paper_header():
    parser = Parser.__new__(Parser)
    raw_text = """
    # Example Paper
    A. One, B. Two
    Department of Computer Science and Engineering, Washington University in St. Louis, USA
    Department of Computer Science and Engineering, George Mason University
    Brown School at Washington University in St. Louis, USA

    #### Abstract
    Body starts here.
    """

    affiliations = parser._extract_affiliations(raw_text)

    assert "Washington University in St. Louis" in affiliations
    assert "George Mason University" in affiliations
    assert "Brown School at Washington University in St. Louis" in affiliations


def test_parser_does_not_extract_reference_doi_as_paper_doi():
    parser = Parser.__new__(Parser)
    raw_text = """
    # Can Watermarked LLMs Be Identified By Users Via Crafted Prompts?
    Tsinghua University

    A BSTRACT
    Body starts here.

    References
    Some unrelated paper. DOI: 10.1145/3626772.3661377
    """

    assert parser._extract_doi(raw_text) is None


def test_parser_extracts_watermark_paper_affiliations_from_header():
    parser = Parser.__new__(Parser)
    raw_text = """
    # Can Watermarked LLMs Be Identified By Users Via Crafted Prompts?
    Tsinghua University
    Beijing University of Posts and Telecommunications
    The Chinese University of Hongkong
    University of Illinois at Chicago
    Hongkong University of Science and Technology (Guangzhou)

    A BSTRACT
    Body starts here.
    """

    affiliations = parser._extract_affiliations(raw_text)

    assert "Tsinghua University" in affiliations
    assert "Beijing University of Posts and Telecommunications" in affiliations
    assert "The Chinese University of Hong Kong" in affiliations
    assert "University of Illinois at Chicago" in affiliations
    assert "Hong Kong University of Science and Technology (Guangzhou)" in affiliations


def test_parser_uses_cached_llm_outputs_when_model_calls_fail(tmp_path, monkeypatch):
    output_dir = tmp_path / "output"
    content_dir = output_dir / "content"
    content_dir.mkdir(parents=True)
    cached_narrative = {
        "and": "Problem context",
        "but": "Core tension",
        "therefore": "Poster takeaway",
        "meta": {
            "poster_title": "Cached Poster Title",
            "authors": "A. Author and B. Author",
        },
    }
    cached_classified_visuals = {
        "key_visual": "figure_1",
        "problem_illustration": [],
        "method_workflow": ["figure_1"],
        "main_results": [],
        "comparative_results": [],
        "supporting": [],
    }
    cached_structured_sections = {
        "paper_sections": [
            {"section_name": "Intro", "section_type": "foundation", "content": "Intro content"},
            {"section_name": "Method", "section_type": "method", "content": "Method content"},
            {"section_name": "Results", "section_type": "evaluation", "content": "Results content"},
        ]
    }
    (content_dir / "narrative_content.json").write_text(json.dumps(cached_narrative), encoding="utf-8")
    (content_dir / "classified_visuals.json").write_text(json.dumps(cached_classified_visuals), encoding="utf-8")
    (content_dir / "structured_sections.json").write_text(json.dumps(cached_structured_sections), encoding="utf-8")

    parser = Parser.__new__(Parser)
    parser.name = "parser"
    state = create_state(str(tmp_path / "paper.pdf"))
    state["output_dir"] = str(output_dir)

    monkeypatch.setattr(parser, "_extract_raw_text", lambda pdf_path, content_dir: ("raw text", object()))
    monkeypatch.setattr(parser, "_extract_assets", lambda raw_result, poster_name, assets_dir: ({}, {}))
    monkeypatch.setattr(parser, "_extract_title_authors", lambda raw_text, config, state: (_ for _ in ()).throw(RuntimeError("title model down")))
    monkeypatch.setattr(parser, "_generate_narrative_content", lambda raw_text, config, state: (_ for _ in ()).throw(RuntimeError("narrative model down")))
    monkeypatch.setattr(parser, "_classify_visual_assets", lambda figures, tables, raw_text, config, state: (_ for _ in ()).throw(RuntimeError("visual model down")))
    monkeypatch.setattr(parser, "_extract_structured_sections", lambda raw_text, config, state: (_ for _ in ()).throw(RuntimeError("section model down")))
    monkeypatch.setattr(parser, "_extract_affiliations", lambda raw_text: [])
    monkeypatch.setattr(parser, "_extract_doi", lambda raw_text: None)

    result = parser(state)

    assert result["errors"] == []
    assert result["narrative_content"]["meta"]["poster_title"] == "Cached Poster Title"
    assert result["classified_visuals"] == cached_classified_visuals
    assert result["structured_sections"] == cached_structured_sections


def test_poster_keypoint_selector_caps_to_ten_by_reading_order(tmp_path, monkeypatch):
    payload = {
        "paper_poster_keypoints": [
            {"id": index, "key_point": f"Poster-worthy claim {index}", "section": "Introduction"}
            for index in range(1, 13)
        ],
        "reading_order": list(range(12, 0, -1)),
    }

    class FakeResponse:
        content = json.dumps(payload)
        input_tokens = 10
        output_tokens = 20

    class FakeAgent:
        def __init__(self, *args, **kwargs):
            pass

        def step(self, message):
            return FakeResponse()

    monkeypatch.setattr("src.agents.poster_keypoint_selector.LangGraphAgent", FakeAgent)
    state = create_state(str(tmp_path / "paper.pdf"))
    state["output_dir"] = str(tmp_path / "output")
    state["raw_text"] = "Full paper text with enough content for keypoint selection."

    result = PosterKeypointSelector()(state)

    assert len(result["paper_poster_keypoints"]) == 10
    assert result["poster_reading_order"] == list(range(1, 11))
    assert result["paper_poster_keypoints"][0]["original_id"] == 12
    assert result["poster_keypoint_selection_report"]["dropped_original_ids"] == [2, 1]
    assert Path(state["output_dir"], "content", "poster_keypoint_selection.json").exists()


def test_poster_keypoint_selector_falls_back_to_structured_sections(tmp_path, monkeypatch):
    class FailingAgent:
        def __init__(self, *args, **kwargs):
            pass

        def step(self, message):
            raise RuntimeError("offline")

    monkeypatch.setattr("src.agents.poster_keypoint_selector.LangGraphAgent", FailingAgent)
    state = create_state(str(tmp_path / "paper.pdf"))
    state["output_dir"] = str(tmp_path / "output")
    state["raw_text"] = "Paper text."
    state["structured_sections"] = {
        "paper_sections": [
            {"section_name": "Introduction", "key_points": ["Problem statement is poster relevant."]},
            {"section_name": "Method", "key_points": ["The framework has a stable arena-based update rule."]},
        ]
    }

    result = PosterKeypointSelector()(state)

    assert len(result["paper_poster_keypoints"]) == 2
    assert result["poster_keypoint_selection_report"]["source"] == "structured_sections_fallback"
    assert result["paper_poster_keypoints"][1]["section"] == "Method"


def test_curator_aligns_story_board_to_keypoints():
    curator = StoryBoardCurator()
    state = create_state("/tmp/paper.pdf")
    state["paper_poster_keypoints"] = [
        {"id": index, "key_point": f"Key contribution or result {index}", "section": "Method" if index <= 6 else "Experiments"}
        for index in range(1, 11)
    ]
    state["poster_reading_order"] = list(range(1, 11))
    story_board = {
        "spatial_content_plan": {
            "sections": [
                {
                    "section_id": "method_summary",
                    "section_title": "Long Method Summary Title",
                    "column_assignment": "middle",
                    "vertical_priority": "top",
                    "text_content": ["Existing method detail."],
                    "visual_assets": [],
                }
            ]
        }
    }
    visual_context = {"valid_visual_ids": [], "keypoint_target_count": 10}

    curator._align_sections_to_keypoints(story_board, state, visual_context, classified_visuals={})

    sections = story_board["spatial_content_plan"]["sections"]
    assert len(sections) == 10
    assert [section["keypoint_id"] for section in sections] == list(range(1, 11))
    assert all(section["source_section"] in {"Method", "Experiments"} for section in sections)
    assert all(len(section["section_title"].split()) <= curator.validation_config["max_title_words"] for section in sections)


def test_curator_normalizes_poster_text_items():
    curator = StoryBoardCurator()

    cleaned = curator._clean_poster_text_items(
        [
            "• **Problem:** Dense paper text should become a clean poster item.",
            "    ◦ Nested details should not keep sub-bullet markers.",
            "1) Ordered list markers should be stripped.",
            "",
            "Step 2: Workflow prefixes should be removed.",
            "The results are presented in Table 2.",
            "Baseline Methods We compare the proposed approach to the following baselines: - 1. *Random:* Each.",
            "This suggests the inductive bias introduced in the hierarchical.",
        ],
        max_items=6,
    )

    assert cleaned == [
        "**Problem:** Dense paper text should become a clean poster item.",
        "Nested details should not keep sub-bullet markers.",
        "Ordered list markers should be stripped.",
        "Workflow prefixes should be removed.",
    ]
    assert curator._clean_section_title("Main results with table") == "Main Results"
    assert curator._clean_section_title("Paper S Main") == "Paper's Main"
    assert normalize_text_for_poster("Paper's Main") == "Paper's Main"
    assert normalize_text_for_poster("Main Results") == "Main Results"
    assert normalize_text_for_poster(
        "Baseline Methods We compare the proposed approach to the following baselines: - 1. Random: Each."
    ) == ""
    assert normalize_text_for_poster("This suggests the inductive bias introduced in the hierarchical.") == ""
    assert normalize_text_for_poster(
        "Overall empirical conclusion: HAGS is the strongest method across cost."
    ) == "Overall empirical conclusion: HAGS is the strongest method."
    # unrendered LaTeX math and sub/sup markup are stripped so they neither show as
    # garbage nor bloat a block past its panel
    assert "$" not in normalize_text_for_poster("For anchor $x_{i}$ the loss $$\\mathcal{L} = -\\log p$$ holds.")
    assert "\\" not in normalize_text_for_poster("For anchor $x_{i}$ the loss $$\\mathcal{L} = -\\log p$$ holds.")
    assert normalize_text_for_poster("The score S<sub>i</sub> uses the k<sup>th</sup> layer.") == (
        "The score Si uses the kth layer."
    )
    assert LayoutAgent()._section_title_label({}, "Paper S Main", create_state("/tmp/paper.pdf"))["title"] == "Paper's Main"
    assert normalize_title_for_poster("Active Geospatial Search For Effcient Tenant Eviction Outreach") == (
        "Active Geospatial Search for Efficient Tenant Eviction Outreach"
    )


def test_curator_groups_keypoints_for_dense_landscape_template(tmp_path):
    curator = StoryBoardCurator()
    capacity_state = create_state(
        str(tmp_path / "paper.pdf"),
        layout_template="cluster_104_landscape",
        width=54,
        height=27,
    )
    capacity_state["output_dir"] = str(tmp_path / "output_dense")
    capacity_state = TemplateCapacityPlanner()(capacity_state)
    state = create_state("/tmp/paper.pdf")
    state["paper_poster_keypoints"] = [
        {"id": 1, "key_point": "Arena comparisons can produce unstable Elo rankings when the same battles are processed in different orders.", "section": "Introduction"},
        {"id": 2, "key_point": "The paper frames stable LLM evaluation as likelihood-based estimation over the complete battle set.", "section": "Introduction"},
        {"id": 3, "key_point": "m-ELO reformulates Elo scoring as maximum likelihood estimation instead of sequential updates.", "section": "Method"},
        {"id": 4, "key_point": "am-ELO extends m-ELO by modeling annotator ability in the pairwise probability function.", "section": "Method"},
        {"id": 5, "key_point": "The framework estimates model scores and annotator abilities jointly from arena records.", "section": "Method"},
        {"id": 6, "key_point": "Experiments compare Elo, m-ELO, and am-ELO on arena-style LLM evaluation data.", "section": "Experiments"},
        {"id": 7, "key_point": "Results show am-ELO achieves lower loss than the baseline Elo-style estimators.", "section": "Results"},
        {"id": 8, "key_point": "Prediction experiments indicate better generalization for am-ELO.", "section": "Results"},
        {"id": 9, "key_point": "Robustness tests analyze perturbations to arena outcomes and annotator behavior.", "section": "Robustness"},
        {"id": 10, "key_point": "The final takeaway is that annotator-aware MLE gives more stable arena-based LLM evaluation.", "section": "Conclusion"},
    ]
    state["poster_reading_order"] = list(range(1, 11))
    story_board = {"spatial_content_plan": {"sections": []}}
    visual_context = {
        "valid_visual_ids": ["figure_2", "figure_3", "table_3"],
        "keypoint_target_count": 10,
        "keypoint_section_target_count": 7,
        "keypoint_grouping_mode": True,
        "requested_layout_template": "cluster_104_landscape",
        "template_fast_mode": True,
        "fast_block_contract": capacity_state["fast_block_contract"],
        "fast_visual_policy": capacity_state["fast_visual_policy"],
        "template_layout": capacity_state["layout_template_metadata"],
        "visual_assets_heights": {
            "figure_2": {"aspect_ratio": 1.7},
            "figure_3": {"aspect_ratio": 2.3},
            "table_3": {"aspect_ratio": 0.94},
        },
    }
    classified_visuals = {
        "key_visual": "figure_2",
        "method_workflow": [],
        "main_results": ["figure_3"],
        "comparative_results": ["table_3"],
        "supporting": [],
    }

    curator._align_sections_to_keypoints(story_board, state, visual_context, classified_visuals)

    sections = story_board["spatial_content_plan"]["sections"]
    source_ids = [keypoint_id for section in sections for keypoint_id in section["source_keypoint_ids"]]
    visual_ids = [
        visual["visual_id"]
        for section in sections
        for visual in section.get("visual_assets", [])
    ]
    assert len(sections) == 7
    assert source_ids == list(range(1, 11))
    assert all(section.get("source_keypoint_ids") for section in sections)
    assert "figure_2" in visual_ids
    assert "figure_3" in visual_ids
    assert "table_3" not in visual_ids


def test_curator_groups_keypoints_for_six_slot_landscape_template(tmp_path):
    curator = StoryBoardCurator()
    capacity_state = create_state(
        str(tmp_path / "paper.pdf"),
        layout_template="cluster_43_landscape",
        width=54,
        height=27,
    )
    capacity_state["output_dir"] = str(tmp_path / "output_six")
    capacity_state = TemplateCapacityPlanner()(capacity_state)
    state = create_state("/tmp/paper.pdf")
    state["paper_poster_keypoints"] = [
        {"id": index, "key_point": f"am-ELO poster point {index} with stable arena evaluation evidence.", "section": "Method" if index <= 6 else "Results"}
        for index in range(1, 11)
    ]
    state["poster_reading_order"] = list(range(1, 11))
    story_board = {"spatial_content_plan": {"sections": []}}
    visual_context = {
        "valid_visual_ids": ["figure_2", "figure_3", "table_2", "table_3"],
        "keypoint_target_count": 10,
        "keypoint_section_target_count": 6,
        "keypoint_grouping_mode": True,
        "requested_layout_template": "cluster_43_landscape",
        "template_fast_mode": True,
        "fast_block_contract": capacity_state["fast_block_contract"],
        "fast_visual_policy": capacity_state["fast_visual_policy"],
        "template_layout": capacity_state["layout_template_metadata"],
        "visual_assets_heights": {
            "figure_2": {"aspect_ratio": 1.6},
            "figure_3": {"aspect_ratio": 2.0},
            "table_2": {"aspect_ratio": 2.4},
            "table_3": {"aspect_ratio": 1.2},
        },
    }
    classified_visuals = {
        "key_visual": "figure_2",
        "method_workflow": [],
        "main_results": ["figure_3", "table_2"],
        "comparative_results": ["table_3"],
        "supporting": [],
    }

    curator._align_sections_to_keypoints(story_board, state, visual_context, classified_visuals)

    sections = story_board["spatial_content_plan"]["sections"]
    assert len(sections) == 6
    assert [keypoint_id for section in sections for keypoint_id in section["source_keypoint_ids"]] == list(range(1, 11))
    assert [section["preferred_slot_id"] for section in sections] == [
        "slot_1",
        "slot_2",
        "slot_3",
        "slot_4",
        "slot_5",
        "slot_6",
    ]
    visual_ids = [
        visual["visual_id"]
        for section in sections
        for visual in section.get("visual_assets", [])
    ]
    assert {"figure_2", "figure_3", "table_2"}.issubset(set(visual_ids))


def test_curator_standard_landscape_uses_distinct_professional_section_titles(tmp_path):
    capacity_state = create_state(
        str(tmp_path / "paper.pdf"),
        layout_template="cluster_43_landscape",
        width=54,
        height=27,
    )
    capacity_state["output_dir"] = str(tmp_path / "output_titles")
    capacity_state = TemplateCapacityPlanner()(capacity_state)
    curator = StoryBoardCurator()
    state = create_state(str(tmp_path / "paper.pdf"))
    state["paper_poster_keypoints"] = [
        {"id": 1, "key_point": "Eviction outreach requires sequential, budget-constrained search under uncertainty.", "section": "Introduction"},
        {"id": 2, "key_point": "AGS formalizes the exploration-exploitation tradeoff with geospatial and travel-cost considerations.", "section": "Introduction"},
        {"id": 3, "key_point": "HAGS is introduced to make AGS scalable to large urban search spaces.", "section": "Introduction"},
        {"id": 4, "key_point": "The paper's main application is finding properties with tenants at risk of imminent eviction to support preventive outreach.", "section": "Introduction"},
        {"id": 5, "key_point": "Prior active search methods do not adequately model geospatial structure and travel costs.", "section": "Related Work"},
        {"id": 6, "key_point": "Visual active search is related but not directly applicable to parcel-level eviction outreach.", "section": "Related Work"},
        {"id": 7, "key_point": "Previous eviction prediction work lacks adaptive sequential search and multimodal decision policies.", "section": "Related Work"},
        {"id": 8, "key_point": "AGS formulates outreach as budget-constrained sequential geospatial search.", "section": "Methodology"},
        {"id": 9, "key_point": "The policy combines a prediction module with a search module under remaining budget.", "section": "Methodology"},
        {"id": 10, "key_point": "Overall empirical conclusion: HAGS is the strongest method across cost models, budgets, and target rates.", "section": "Results"},
    ]
    state["poster_reading_order"] = list(range(1, 11))
    stale_result = "Overall empirical conclusion: HAGS is the strongest method across cost models, budgets, and target rates."
    story_board = {
        "spatial_content_plan": {
            "sections": [
                {"section_id": "old_motivation", "text_content": []},
                {"section_id": "old_hags", "text_content": [stale_result]},
                {"section_id": "old_application", "text_content": []},
                {
                    "section_id": "old_prior",
                    "text_content": [
                        "Hierarchy is the main reason the method succeeds at urban scale.",
                        stale_result,
                    ],
                },
                {
                    "section_id": "old_prediction",
                    "text_content": [
                        "AGS fills the gap between predictive eviction analytics and operational outreach planning."
                    ],
                },
                {"section_id": "old_policy", "text_content": ["Hierarchy is the main reason the method succeeds at urban scale."]},
            ]
        }
    }
    visual_context = {
        "valid_visual_ids": ["figure_1", "figure_2", "table_1"],
        "keypoint_target_count": 10,
        "keypoint_section_target_count": 6,
        "keypoint_grouping_mode": True,
        "requested_layout_template": "cluster_43_landscape",
        "template_fast_mode": True,
        "fast_block_contract": capacity_state["fast_block_contract"],
        "fast_visual_policy": capacity_state["fast_visual_policy"],
        "template_layout": capacity_state["layout_template_metadata"],
        "visual_assets_heights": {
            "figure_1": {"aspect_ratio": 1.6},
            "figure_2": {"aspect_ratio": 2.0},
            "table_1": {"aspect_ratio": 2.4},
        },
    }
    classified_visuals = {
        "key_visual": "figure_1",
        "method_workflow": [],
        "main_results": ["figure_2", "table_1"],
        "comparative_results": [],
        "supporting": [],
    }

    curator._align_sections_to_keypoints(story_board, state, visual_context, classified_visuals)

    titles = [section["section_title"] for section in story_board["spatial_content_plan"]["sections"]]
    assert len(titles) == len(set(titles))
    assert "Paper S Main" not in titles
    assert "Paper's Main" not in titles
    assert "HAGS Overview" in titles
    assert "Main Application" in titles
    assert "Prior Methods" in titles
    assert "Search Policy" in titles or "Main Results" in titles
    by_title = {section["section_title"]: section for section in story_board["spatial_content_plan"]["sections"]}
    assert by_title["Prior Prediction"]["source_keypoint_ids"] == [7, 8]
    assert by_title["Search Policy"]["source_keypoint_ids"] == [9, 10]
    assert stale_result not in "\n".join(by_title["HAGS Overview"]["text_content"])
    assert stale_result not in "\n".join(by_title["Prior Methods"]["text_content"])
    assert stale_result in "\n".join(by_title["Search Policy"]["text_content"])


def test_curator_standard_visual_assignment_preserves_generated_teaser(tmp_path):
    capacity_state = create_state(
        str(tmp_path / "paper.pdf"),
        layout_template="cluster_43_landscape",
        width=54,
        height=27,
    )
    capacity_state["output_dir"] = str(tmp_path / "output_teaser_preserve")
    capacity_state = TemplateCapacityPlanner()(capacity_state)
    curator = StoryBoardCurator()
    state = create_state(str(tmp_path / "paper.pdf"))
    state["paper_poster_keypoints"] = [
        {"id": index, "key_point": f"Poster keypoint {index} with method or result evidence.", "section": "Method" if index <= 6 else "Results"}
        for index in range(1, 11)
    ]
    state["poster_reading_order"] = list(range(1, 11))
    story_board = {
        "spatial_content_plan": {
            "sections": [
                {
                    "section_id": "motivation",
                    "section_title": "Motivation",
                    "text_content": ["Motivation text."],
                    "visual_assets": [{"visual_id": "generated_teaser_1"}],
                    "generated_teaser_summary": True,
                }
            ]
        }
    }
    visual_context = {
        "valid_visual_ids": ["generated_teaser_1", "figure_1", "figure_2", "table_1", "table_2"],
        "keypoint_target_count": 10,
        "keypoint_section_target_count": 6,
        "keypoint_grouping_mode": True,
        "requested_layout_template": "cluster_43_landscape",
        "template_fast_mode": True,
        "fast_block_contract": capacity_state["fast_block_contract"],
        "fast_visual_policy": capacity_state["fast_visual_policy"],
        "template_layout": capacity_state["layout_template_metadata"],
        "visual_assets": {
            "generated_teaser_1": {"asset_type": "figure", "aspect": 2.2},
            "figure_1": {"asset_type": "figure", "aspect": 1.6},
            "figure_2": {"asset_type": "figure", "aspect": 2.0},
            "table_1": {"asset_type": "table", "aspect": 2.4},
            "table_2": {"asset_type": "table", "aspect": 2.2},
        },
    }
    classified_visuals = {
        "key_visual": "figure_1",
        "method_workflow": [],
        "main_results": ["figure_2", "table_1"],
        "comparative_results": ["table_2"],
        "supporting": [],
    }

    curator._align_sections_to_keypoints(story_board, state, visual_context, classified_visuals)

    first_section = story_board["spatial_content_plan"]["sections"][0]
    visual_ids = [
        visual["visual_id"]
        for section in story_board["spatial_content_plan"]["sections"]
        for visual in section.get("visual_assets", [])
    ]
    assert first_section["section_title"] == "Motivation"
    assert first_section["visual_assets"][0]["visual_id"] == "generated_teaser_1"
    assert len([visual_id for visual_id in visual_ids if not visual_id.startswith("generated_teaser")]) == 4
    figure_slots = set(visual_context["fast_visual_policy"]["figure_slots"])
    table_slots = set(visual_context["fast_visual_policy"]["table_slots"])
    for section in story_board["spatial_content_plan"]["sections"]:
        slot_id = section["preferred_slot_id"]
        for visual in section.get("visual_assets", []):
            visual_id = visual["visual_id"]
            if visual_id.startswith("figure_"):
                assert slot_id in figure_slots
            if visual_id.startswith("table_"):
                assert slot_id in table_slots


def test_template_capacity_planner_builds_landscape_fast_contract(tmp_path):
    state = create_state(
        str(tmp_path / "paper.pdf"),
        layout_template="cluster_43_landscape",
        width=54,
        height=27,
    )
    state["output_dir"] = str(tmp_path / "output")

    result = TemplateCapacityPlanner()(state)

    contract = result["fast_block_contract"]
    assert result["template_fast_mode"] is True
    assert contract["template_id"] == "cluster_43_landscape"
    assert contract["slot_order"] == ["slot_1", "slot_2", "slot_3", "slot_4", "slot_5", "slot_6"]
    assert len(contract["blocks"]) == 6
    assert contract["by_slot"]["slot_1"]["min_chars"] > 0
    assert contract["by_slot"]["slot_6"]["visual_policy"] == "table_with_callouts"
    assert result["fast_visual_policy"]["figure_count"] == 2
    assert result["fast_visual_policy"]["table_count"] == 2
    assert Path(state["output_dir"], "content", "fast_block_contract.json").exists()


def test_template_capacity_planner_applies_rich_visual_density(tmp_path):
    state = create_state(
        str(tmp_path / "paper.pdf"),
        layout_template="cluster_43_landscape",
        width=54,
        height=27,
        visual_density="rich",
    )
    state["output_dir"] = str(tmp_path / "output_rich")

    result = TemplateCapacityPlanner()(state)

    assert result["fast_visual_policy"]["visual_density"] == "rich"
    assert result["fast_visual_policy"]["figure_count"] == 3
    assert result["fast_visual_policy"]["table_count"] == 2
    assert result["fast_visual_policy"]["max_visuals_total"] == 5
    assert result["fast_visual_policy"]["visual_footprint"]["enabled"] is True
    assert result["fast_block_contract"]["by_slot"]["slot_2"]["visual_footprint"]["min_width"] > 0


def test_template_capacity_planner_orders_visual_slots_by_area(tmp_path):
    state = create_state(
        str(tmp_path / "paper.pdf"),
        layout_template="cluster_96_landscape",
        width=54,
        height=27,
        visual_density="rich",
    )
    state["output_dir"] = str(tmp_path / "output_area_order")

    result = TemplateCapacityPlanner()(state)

    policy = result["fast_visual_policy"]
    assert policy["figure_slots"] == ["slot_2", "slot_1"]
    assert policy["table_slots"] == ["slot_5"]
    rejected = policy["rejected_visual_slots"]
    assert any(item["slot_id"] == "slot_3" for item in rejected["figure"])
    assert any(item["slot_id"] == "slot_6" for item in rejected["table"])


def test_template_capacity_planner_excludes_narrow_portrait_figure_slots(tmp_path):
    state = create_state(
        str(tmp_path / "paper.pdf"),
        layout_template="cluster_22_portrait",
        width=36,
        height=50.88,
        visual_density="rich",
    )
    state["output_dir"] = str(tmp_path / "output_portrait_slots")

    result = TemplateCapacityPlanner()(state)

    policy = result["fast_visual_policy"]
    assert "slot_3" in policy["figure_slots"]
    assert "slot_2" not in policy["figure_slots"]
    assert policy["figure_count"] == len(policy["figure_slots"])
    rejected = policy["rejected_visual_slots"]["figure"]
    assert any(item["slot_id"] == "slot_2" and "width" in item["failed"] for item in rejected)
    rejected_slot_ids = {
        item["slot_id"]
        for items in policy["rejected_visual_slots"].values()
        for item in items
    }
    assert rejected_slot_ids <= set(result["fast_block_contract"]["by_slot"])
    slot_2_contract = result["fast_block_contract"]["by_slot"]["slot_2"]
    assert slot_2_contract["visual_policy"] == "text_summary"
    assert 350 <= slot_2_contract["target_chars"] <= 450
    assert slot_2_contract["capacity_warning"] == "visual_slot_too_narrow_text_fallback"


def test_template_capacity_planner_uses_conservative_portrait_text_capacity():
    planner = TemplateCapacityPlanner()
    landscape = planner._estimate_capacity({"w": 12.41, "h": 11.6256, "poster_orientation": "landscape"}, "text_summary")
    portrait = planner._estimate_capacity({"w": 12.41, "h": 11.6256, "poster_orientation": "portrait"}, "text_summary")

    assert portrait["target_chars"] < landscape["target_chars"]
    assert portrait["target_chars"] == 390


def test_standard_template_preselector_auto_selects_default_standard_landscape_template(tmp_path):
    state = create_state(str(tmp_path / "paper.pdf"), layout_template="auto", width=54, height=36)
    state["output_dir"] = str(tmp_path / "output")
    state["structured_sections"] = {
        "paper_sections": [
            {"section_name": f"Section {index}", "section_type": "method", "key_points": ["A"]}
            for index in range(6)
        ]
    }
    state["classified_visuals"] = {"main_results": ["table_1"]}
    state["visual_assets"] = {
        "figure_1": {"asset_type": "figure", "aspect": 2.0},
        "figure_2": {"asset_type": "figure", "aspect": 1.5},
        "table_1": {"asset_type": "table", "aspect": 2.4},
    }

    result = StandardTemplatePreselector()(state)

    assert result["resolved_layout_template"] == "cluster_43_landscape"
    assert result["poster_width"] == 54.0
    assert result["poster_height"] == 27.0
    assert result["enable_block_vlm_review"] is True
    assert Path(state["output_dir"], "content", "standard_template_selection_report.json").exists()


def test_template_capacity_planner_builds_standard_template_contract(tmp_path):
    state = create_state(
        str(tmp_path / "paper.pdf"),
        layout_template="cluster_43_landscape",
        width=54,
        height=27,
    )
    state["output_dir"] = str(tmp_path / "output")

    result = TemplateCapacityPlanner()(state)

    contract = result["fast_block_contract"]
    source_ids = [keypoint_id for block in contract["blocks"] for keypoint_id in block["source_keypoint_ids"]]
    assert result["template_fast_mode"] is True
    assert contract["template_id"] == "cluster_43_landscape"
    assert len(contract["blocks"]) == 6
    assert source_ids == list(range(1, 11))
    assert result["fast_visual_policy"]["figure_slots"]
    assert result["fast_visual_policy"]["table_slots"]


def test_poster_keypoint_selector_prompt_includes_fast_contract():
    selector = PosterKeypointSelector()
    fast_contract = {
        "template_id": "cluster_43_landscape",
        "blocks": [
            {
                "slot_id": "slot_1",
                "slot_role": "Motivation",
                "visual_policy": "text_only",
                "target_chars": 500,
                "source_keypoint_ids": [1, 2],
            }
        ],
    }

    prompt = selector._build_prompt("Paper text.", fast_contract)

    assert "Fast template-first capacity context" in prompt
    assert "cluster_43_landscape" in prompt
    assert "motivation, method/architecture" in prompt


def test_curator_groups_keypoints_for_standard_landscape_template(tmp_path):
    capacity_state = create_state(
        str(tmp_path / "paper.pdf"),
        layout_template="cluster_43_landscape",
        width=54,
        height=27,
    )
    capacity_state["output_dir"] = str(tmp_path / "output")
    capacity_state = TemplateCapacityPlanner()(capacity_state)

    curator = StoryBoardCurator()
    state = create_state(str(tmp_path / "paper.pdf"))
    state["paper_poster_keypoints"] = [
        {"id": index, "key_point": f"HAGS poster point {index} with method or result evidence.", "section": "Method" if index <= 6 else "Results"}
        for index in range(1, 11)
    ]
    state["poster_reading_order"] = list(range(1, 11))
    story_board = {"spatial_content_plan": {"sections": []}}
    visual_context = {
        "valid_visual_ids": ["figure_1", "figure_2", "table_1"],
        "keypoint_target_count": 10,
        "keypoint_section_target_count": 6,
        "keypoint_grouping_mode": True,
        "requested_layout_template": "cluster_43_landscape",
        "template_fast_mode": True,
        "fast_block_contract": capacity_state["fast_block_contract"],
        "fast_visual_policy": capacity_state["fast_visual_policy"],
        "template_layout": capacity_state["layout_template_metadata"],
        "visual_assets_heights": {
            "figure_1": {"aspect_ratio": 1.6},
            "figure_2": {"aspect_ratio": 2.0},
            "table_1": {"aspect_ratio": 2.4},
        },
    }
    classified_visuals = {
        "key_visual": "figure_1",
        "method_workflow": [],
        "main_results": ["figure_2", "table_1"],
        "comparative_results": [],
        "supporting": [],
    }

    curator._align_sections_to_keypoints(story_board, state, visual_context, classified_visuals)

    sections = story_board["spatial_content_plan"]["sections"]
    source_ids = [keypoint_id for section in sections for keypoint_id in section["source_keypoint_ids"]]
    visual_ids = [
        visual["visual_id"]
        for section in sections
        for visual in section.get("visual_assets", [])
    ]
    assert len(sections) == 6
    assert source_ids == list(range(1, 11))
    assert all(section.get("preferred_slot_id") for section in sections)
    assert all(section.get("capacity_budget") for section in sections)
    assert {"figure_1", "figure_2", "table_1"}.issubset(set(visual_ids))
    key_holder = next(section for section in sections if section.get("visual_assets") and section["visual_assets"][0]["visual_id"] == "figure_1")
    assert key_holder["column_assignment"] == "middle"
    assert key_holder["vertical_priority"] == "top"


def test_curator_backfills_standard_template_visuals_by_large_blocks(tmp_path):
    capacity_state = create_state(
        str(tmp_path / "paper.pdf"),
        layout_template="cluster_96_landscape",
        width=54,
        height=27,
        visual_density="rich",
    )
    capacity_state["output_dir"] = str(tmp_path / "output")
    capacity_state = TemplateCapacityPlanner()(capacity_state)

    curator = StoryBoardCurator()
    state = create_state(str(tmp_path / "paper.pdf"))
    state["paper_poster_keypoints"] = [
        {"id": index, "key_point": f"Poster point {index} with method or result evidence.", "section": "Method" if index <= 6 else "Results"}
        for index in range(1, 11)
    ]
    state["poster_reading_order"] = list(range(1, 11))
    story_board = {"spatial_content_plan": {"sections": []}}
    visual_context = {
        "valid_visual_ids": ["figure_1", "figure_2", "table_1", "table_2"],
        "keypoint_target_count": 10,
        "keypoint_section_target_count": 6,
        "keypoint_grouping_mode": True,
        "requested_layout_template": "cluster_96_landscape",
        "template_fast_mode": True,
        "fast_block_contract": capacity_state["fast_block_contract"],
        "fast_visual_policy": capacity_state["fast_visual_policy"],
        "template_layout": capacity_state["layout_template_metadata"],
        "visual_assets_heights": {
            "figure_1": {"aspect_ratio": 2.52},
            "figure_2": {"aspect_ratio": 2.84},
            "table_1": {"aspect_ratio": 4.3},
            "table_2": {"aspect_ratio": 4.2},
        },
    }
    classified_visuals = {
        "key_visual": "figure_2",
        "method_workflow": ["figure_1"],
        "main_results": ["table_1", "table_2"],
        "comparative_results": [],
        "supporting": [],
    }

    curator._align_sections_to_keypoints(story_board, state, visual_context, classified_visuals)

    sections = story_board["spatial_content_plan"]["sections"]
    by_slot = {section["preferred_slot_id"]: section for section in sections}
    visual_ids = [
        visual["visual_id"]
        for section in sections
        for visual in section.get("visual_assets", [])
    ]
    slot_2_visual_ids = [visual["visual_id"] for visual in by_slot["slot_2"]["visual_assets"]]
    slot_3_visual_ids = [visual["visual_id"] for visual in by_slot["slot_3"]["visual_assets"]]
    assert slot_2_visual_ids[0] == "figure_2"
    assert not any(visual_id.startswith("table_") for visual_id in slot_2_visual_ids)
    assert not any(visual_id.startswith("table_") for visual_id in slot_3_visual_ids)
    table_slot_visual_ids = [
        visual["visual_id"]
        for slot_id in capacity_state["fast_visual_policy"]["table_slots"]
        for visual in by_slot[slot_id].get("visual_assets", [])
    ]
    assert sum(visual_id.startswith("table_") for visual_id in table_slot_visual_ids) == capacity_state["fast_visual_policy"]["table_count"]
    assert sum(visual_id.startswith("figure_") for visual_id in visual_ids) == 2
    assert sum(visual_id.startswith("table_") for visual_id in visual_ids) == capacity_state["fast_visual_policy"]["table_count"]


def test_template_block_planner_preserves_feasible_fast_assigned_visuals(tmp_path):
    capacity_state = create_state(
        str(tmp_path / "paper.pdf"),
        layout_template="cluster_96_landscape",
        width=54,
        height=27,
        visual_density="rich",
    )
    capacity_state["output_dir"] = str(tmp_path / "output_capacity")
    capacity_state = TemplateCapacityPlanner()(capacity_state)

    curator = StoryBoardCurator()
    state = create_state(str(tmp_path / "paper.pdf"), layout_template="cluster_96_landscape", width=54, height=27)
    state["paper_poster_keypoints"] = [
        {"id": index, "key_point": f"Poster point {index} with method or result evidence.", "section": "Method" if index <= 6 else "Results"}
        for index in range(1, 11)
    ]
    state["poster_reading_order"] = list(range(1, 11))
    state["visual_assets"] = {
        "figure_1": {"asset_type": "figure", "aspect": 2.52},
        "figure_2": {"asset_type": "figure", "aspect": 2.84},
        "table_2": {"asset_type": "table", "aspect": 4.216867469879518},
        "table_4": {"asset_type": "table", "aspect": 4.307377049180328},
    }
    visual_context = {
        "valid_visual_ids": list(state["visual_assets"].keys()),
        "keypoint_target_count": 10,
        "keypoint_section_target_count": 6,
        "keypoint_grouping_mode": True,
        "requested_layout_template": "cluster_96_landscape",
        "template_fast_mode": True,
        "fast_block_contract": capacity_state["fast_block_contract"],
        "fast_visual_policy": capacity_state["fast_visual_policy"],
        "template_layout": capacity_state["layout_template_metadata"],
        "visual_assets": state["visual_assets"],
        "visual_assets_heights": {
            visual_id: {"aspect_ratio": asset["aspect"]}
            for visual_id, asset in state["visual_assets"].items()
        },
    }
    classified_visuals = {
        "key_visual": "figure_2",
        "method_workflow": ["figure_1"],
        "main_results": ["table_2", "table_4"],
        "comparative_results": [],
        "supporting": [],
    }
    story_board = {"spatial_content_plan": {"sections": []}}
    curator._align_sections_to_keypoints(story_board, state, visual_context, classified_visuals)

    planner_state = create_state(
        str(tmp_path / "paper.pdf"),
        layout_template="cluster_96_landscape",
        width=54,
        height=27,
        visual_density="rich",
    )
    planner_state["output_dir"] = str(tmp_path / "output_planner")
    planner_state["template_fast_mode"] = True
    planner_state["resolved_layout_template"] = "cluster_96_landscape"
    planner_state["story_board"] = story_board
    planner_state["visual_assets"] = state["visual_assets"]
    planner_state["fast_block_contract"] = capacity_state["fast_block_contract"]
    planner_state["fast_visual_policy"] = capacity_state["fast_visual_policy"]
    planner_state["paper_poster_keypoints"] = state["paper_poster_keypoints"]
    planner_state["poster_reading_order"] = state["poster_reading_order"]

    result = TemplateBlockPlanner()(planner_state)

    visual_ids = [
        visual["visual_id"]
        for section in (result["story_board"].get("spatial_content_plan") or {}).get("sections", [])
        for visual in section.get("visual_assets", [])
    ]
    by_slot = {
        section.get("preferred_slot_id"): section
        for section in (result["story_board"].get("spatial_content_plan") or {}).get("sections", [])
    }
    slot_2_visual_ids = [visual["visual_id"] for visual in by_slot["slot_2"].get("visual_assets", [])]
    assert sum(visual_id.startswith("figure_") for visual_id in visual_ids) == 2
    assert sum(visual_id.startswith("table_") for visual_id in visual_ids) == capacity_state["fast_visual_policy"]["table_count"]
    assert "figure_2" in slot_2_visual_ids
    assert not any(visual_id.startswith("table_") for visual_id in slot_2_visual_ids)


def test_curator_block_template_key_visual_validation_uses_slot_mapping_not_middle_column():
    curator = StoryBoardCurator()
    story_board = {
        "spatial_content_plan": {
            "sections": [
                {
                    "section_id": f"section_{index}",
                    "section_title": f"Block {index}",
                    "column_assignment": "left" if index == 1 else "right",
                    "vertical_priority": "bottom",
                    "source_keypoint_ids": [index],
                    "text_content": ["Evidence-backed point.", "Capacity-aware detail."],
                    "visual_assets": [{"visual_id": "figure_1"}] if index == 1 else [],
                }
                for index in range(1, 5)
            ]
        }
    }
    visual_context = {
        "valid_visual_ids": ["figure_1"],
        "keypoint_target_count": 10,
        "keypoint_section_target_count": 4,
        "keypoint_grouping_mode": True,
        "requested_layout_template": "cluster_43_landscape",
        "visual_assets_heights": {},
    }

    assert curator._validate_story_board(
        story_board,
        classified_visuals={"key_visual": "figure_1"},
        visual_context=visual_context,
    )


def test_curator_portrait_standard_template_keeps_key_visual_and_result_table(tmp_path):
    capacity_state = create_state(
        str(tmp_path / "paper.pdf"),
        layout_template="cluster_8_portrait",
        width=36,
        height=50.88,
    )
    capacity_state["output_dir"] = str(tmp_path / "output")
    capacity_state = TemplateCapacityPlanner()(capacity_state)

    curator = StoryBoardCurator()
    state = create_state(str(tmp_path / "paper.pdf"))
    state["paper_poster_keypoints"] = [
        {"id": index, "key_point": f"HAGS portrait keypoint {index}.", "section": "Method" if index <= 6 else "Results"}
        for index in range(1, 11)
    ]
    state["poster_reading_order"] = list(range(1, 11))
    story_board = {"spatial_content_plan": {"sections": []}}
    visual_context = {
        "valid_visual_ids": ["figure_1", "figure_2", "table_1"],
        "keypoint_target_count": 10,
        "keypoint_section_target_count": 4,
        "keypoint_grouping_mode": True,
        "requested_layout_template": "cluster_8_portrait",
        "template_fast_mode": True,
        "fast_block_contract": capacity_state["fast_block_contract"],
        "fast_visual_policy": capacity_state["fast_visual_policy"],
        "template_layout": capacity_state["layout_template_metadata"],
        "visual_assets_heights": {
            "figure_1": {"aspect_ratio": 1.6},
            "figure_2": {"aspect_ratio": 2.0},
            "table_1": {"aspect_ratio": 2.4},
        },
    }
    classified_visuals = {
        "key_visual": "figure_1",
        "method_workflow": [],
        "main_results": ["figure_2", "table_1"],
        "comparative_results": [],
        "supporting": [],
    }

    guidance = curator._template_layout_guidance(visual_context)
    curator._align_sections_to_keypoints(story_board, state, visual_context, classified_visuals)

    sections = story_board["spatial_content_plan"]["sections"]
    visual_ids = [
        visual["visual_id"]
        for section in sections
        for visual in section.get("visual_assets", [])
    ]
    assert "Produce exactly 4 grouped poster sections" in guidance
    assert "Use up to 2 total visuals" in guidance
    assert visual_ids == ["figure_1", "table_1"]
    assert curator._validate_story_board(story_board, classified_visuals, visual_context)


def test_curator_fast_contract_adds_capacity_budget(tmp_path):
    capacity_state = create_state(
        str(tmp_path / "paper.pdf"),
        layout_template="cluster_43_landscape",
        width=54,
        height=27,
    )
    capacity_state["output_dir"] = str(tmp_path / "output")
    capacity_state = TemplateCapacityPlanner()(capacity_state)

    curator = StoryBoardCurator()
    state = create_state(str(tmp_path / "paper.pdf"))
    state["paper_poster_keypoints"] = [
        {"id": index, "key_point": f"Phishing detection keypoint {index} with method or result evidence.", "section": "Method" if index <= 6 else "Results"}
        for index in range(1, 11)
    ]
    state["poster_reading_order"] = list(range(1, 11))
    story_board = {"spatial_content_plan": {"sections": []}}
    visual_context = {
        "valid_visual_ids": ["figure_1", "figure_2", "table_1"],
        "keypoint_target_count": 10,
        "keypoint_section_target_count": 6,
        "keypoint_grouping_mode": True,
        "requested_layout_template": "cluster_43_landscape",
        "template_fast_mode": True,
        "fast_block_contract": capacity_state["fast_block_contract"],
        "fast_visual_policy": capacity_state["fast_visual_policy"],
        "visual_assets_heights": {
            "figure_1": {"aspect_ratio": 1.6},
            "figure_2": {"aspect_ratio": 2.0},
            "table_1": {"aspect_ratio": 2.4},
        },
    }
    classified_visuals = {
        "key_visual": "figure_1",
        "method_workflow": [],
        "main_results": ["figure_2", "table_1"],
        "comparative_results": [],
        "supporting": [],
    }

    curator._align_sections_to_keypoints(story_board, state, visual_context, classified_visuals)

    sections = story_board["spatial_content_plan"]["sections"]
    by_slot = {section["preferred_slot_id"]: section for section in sections}
    assert len(sections) == 6
    assert by_slot["slot_1"]["target_chars"] > 0
    assert by_slot["slot_5"]["min_chars"] > 0
    assert by_slot["slot_6"]["capacity_budget"]["visual_policy"] == "table_with_callouts"
    assert [keypoint_id for section in sections for keypoint_id in section["source_keypoint_ids"]] == list(range(1, 11))


def test_template_planner_uses_slot_mapping_and_preserves_method_visual():
    planner = TemplateBlockPlanner()
    cfg = load_config()
    template = LayoutTemplates(
        54,
        27,
        margin=cfg["layout"]["poster_margin"],
        col_gap=cfg["layout"]["column_spacing"],
    ).get_template("cluster_43_landscape")
    regions = template["regions"]
    sections = [
        {"section_id": "why", "section_title": "Why ELO Fails", "text_content": ["Motivation"], "visual_assets": [], "content_role": "overview", "preferred_slot_id": "slot_1"},
        {"section_id": "stable", "section_title": "Stable Estimation", "text_content": ["m-ELO"], "visual_assets": [{"visual_id": "figure_3"}], "content_role": "method", "preferred_slot_id": "slot_2"},
        {"section_id": "arena", "section_title": "Arena Setting", "text_content": ["Arena diagram"], "visual_assets": [{"visual_id": "figure_2"}], "content_role": "method", "preferred_slot_id": "slot_3"},
        {"section_id": "annotator", "section_title": "Annotator-Aware ELO", "text_content": ["Annotators"], "visual_assets": [], "content_role": "method", "preferred_slot_id": "slot_6"},
        {"section_id": "robust", "section_title": "Robustness Tests", "text_content": ["Robustness"], "visual_assets": [{"visual_id": "table_3"}], "content_role": "results", "preferred_slot_id": "slot_5"},
        {"section_id": "main", "section_title": "Main Results", "text_content": ["Results"], "visual_assets": [], "content_role": "results", "preferred_slot_id": "slot_4"},
    ]

    assigned = planner._assign_sections_to_regions(
        sections,
        regions,
        hero_section=sections[0],
        hero_region_id="slot_1",
        preserve_order=True,
        template_id="cluster_43_landscape",
    )

    by_id = {section["section_id"]: section for section in assigned}
    assert by_id["main"]["slot_id"] == "slot_4"
    assert by_id["robust"]["slot_id"] == "slot_5"
    assert by_id["annotator"]["slot_id"] == "slot_6"
    assert by_id["arena"]["slot_id"] == "slot_3"
    assert by_id["arena"]["visual_assets"][0]["visual_id"] == "figure_2"


def test_template_planner_preserves_fast_policy_table_in_low_density_slot():
    planner = TemplateBlockPlanner()
    cfg = load_config()
    template = LayoutTemplates(
        54,
        27,
        margin=cfg["layout"]["poster_margin"],
        col_gap=cfg["layout"]["column_spacing"],
    ).get_template("cluster_43_landscape")
    regions = template["regions"]
    sections = [
        {"section_id": "motivation", "section_title": "Why AGS", "text_content": ["Motivation"], "visual_assets": [], "content_role": "foundation", "preferred_slot_id": "slot_1"},
        {"section_id": "method", "section_title": "Search Tension", "text_content": ["Method"], "visual_assets": [{"visual_id": "figure_2"}], "content_role": "method", "preferred_slot_id": "slot_2"},
        {"section_id": "flow", "section_title": "HAGS Flow", "text_content": ["Flow"], "visual_assets": [{"visual_id": "figure_1"}], "content_role": "method", "preferred_slot_id": "slot_3"},
        {"section_id": "learns", "section_title": "How It Learns", "text_content": ["Learning"], "visual_assets": [], "content_role": "method", "preferred_slot_id": "slot_4"},
        {"section_id": "setup", "section_title": "Data Setup", "text_content": ["Setup"], "visual_assets": [], "content_role": "results", "preferred_slot_id": "slot_5"},
        {
            "section_id": "results",
            "section_title": "Key Results",
            "text_content": ["Results"],
            "visual_assets": [{"visual_id": "table_2"}],
            "content_role": "results",
            "preferred_slot_id": "slot_6",
            "capacity_budget": {"visual_policy": "table_with_callouts"},
        },
    ]

    assigned = planner._assign_sections_to_regions(
        sections,
        regions,
        hero_section=sections[1],
        hero_region_id="slot_2",
        preserve_order=True,
        template_id="cluster_43_landscape",
    )

    by_id = {section["section_id"]: section for section in assigned}
    assert by_id["results"]["slot_id"] == "slot_6"
    assert by_id["results"]["visual_assets"][0]["visual_id"] == "table_2"


def test_affiliation_logo_agent_creates_placeholder_when_download_fails(tmp_path, monkeypatch):
    state = create_state(str(tmp_path / "paper.pdf"), enable_affiliation_logos=True)
    state["output_dir"] = str(tmp_path / "output")
    state["affiliations"] = ["Example Research University"]

    agent = AffiliationLogoAgent()
    agent.config["include_placeholders"] = True
    monkeypatch.setattr(agent, "_download_clearbit_logo", lambda domain, output_path: None)
    monkeypatch.setattr(agent, "_download_wikidata_logo", lambda institution, output_path: None)
    monkeypatch.setattr(agent, "_download_known_commons_logo", lambda institution, output_path: None)

    result = agent(state)

    logos = result["affiliation_logos"]
    assert len(logos) == 1
    assert logos[0]["status"] == "placeholder"
    assert Path(logos[0]["logo_path"]).exists()
    assert (Path(state["output_dir"]) / "content" / "affiliation_logos.json").exists()


def test_layout_agent_places_affiliation_logos_in_title_right_region(tmp_path):
    logo_paths = []
    for index in range(3):
        path = tmp_path / f"logo_{index}.png"
        Image.new("RGBA", (300, 120), (255, 255, 255, 255)).save(path)
        logo_paths.append(str(path))

    state = create_state("/tmp/paper.pdf", enable_affiliation_logos=True)
    state["affiliation_logos"] = [
        {
            "institution": f"Institution {index}",
            "logo_path": path,
            "domain": None,
            "source": "test",
            "aspect": 2.5,
        }
        for index, path in enumerate(logo_paths)
    ]

    agent = LayoutAgent()
    template = agent._resolve_template_layout(state)
    elements = agent._create_logo_elements(state, state["poster_width"], template)

    assert len(elements) == 3
    assert {element["type"] for element in elements} == {"institution_logo"}
    title = agent._create_title_element(state, state["poster_width"], template["header"]["h"], template)
    assert all(element["x"] >= title["x"] + title["width"] for element in elements)


def test_layout_agent_avoids_title_conference_logo_overlap_for_portrait_templates(tmp_path):
    logo_path = tmp_path / "conference.png"
    Image.new("RGBA", (900, 420), (20, 80, 160, 255)).save(logo_path)

    agent = LayoutAgent()
    for template_id in list_block_template_ids():
        info = get_block_template_info(template_id)
        canvas = info["recommended_canvas_size"]
        state = create_state("/tmp/paper.pdf", layout_template=template_id)
        state["poster_width"] = canvas["width"]
        state["poster_height"] = canvas["height"]
        state["resolved_layout_template"] = template_id
        state["template_layout_mode"] = "template_prior"
        state["logo_path"] = str(logo_path)

        template = agent._resolve_template_layout(state)
        title = agent._create_title_element(state, state["poster_width"], template["header"]["h"], template)
        logo = next(
            element for element in agent._create_logo_elements(state, state["poster_width"], template)
            if element["type"] == "conf_logo"
        )

        title_right = title["x"] + title["width"]
        logo_left = logo["x"]
        assert title_right <= logo_left


def test_layout_agent_new_landscape_header_keeps_title_and_logo_zone(tmp_path):
    conf_path = tmp_path / "conference.png"
    aff_path = tmp_path / "affiliation.png"
    Image.new("RGBA", (900, 420), (20, 80, 160, 255)).save(conf_path)
    Image.new("RGBA", (700, 700), (160, 40, 60, 255)).save(aff_path)

    state = create_state(
        "/tmp/paper.pdf",
        layout_template="cluster_43_landscape",
        width=54,
        height=27,
        logo_path=str(conf_path),
        aff_logo_path=str(aff_path),
    )
    state["resolved_layout_template"] = "cluster_43_landscape"
    state["template_layout_mode"] = "template_prior"

    agent = LayoutAgent()
    template = agent._resolve_template_layout(state)
    title_box, logo_box = agent._header_title_logo_boxes(state, template)
    title = agent._create_title_element(state, state["poster_width"], template["header"]["h"], template)

    assert title_box["w"] >= template["header"]["w"] * 0.60
    assert logo_box["w"] >= template["header"]["w"] * 0.28
    assert title_box["x"] + title_box["w"] <= logo_box["x"]
    assert title["font_size"] >= 96
    assert title["author_font_size"] <= 72


def test_header_planner_generates_centered_subtitle_for_short_title(tmp_path):
    state = create_state(
        str(tmp_path / "paper.pdf"),
        layout_template="cluster_43_landscape",
        width=54,
        height=27,
        header_route="centered",
        header_subtitle_policy="always",
        header_seed=7,
    )
    state["narrative_content"] = {
        "meta": {
            "poster_title": "Fast Spatial Search",
            "authors": "A. Researcher and B. Scientist",
        }
    }
    state["story_board"] = {
        "spatial_content_plan": {
            "sections": [
                {
                    "section_id": "motivation",
                    "section_title": "Motivation",
                    "content_role": "overview",
                    "text_content": ["Prioritizes high-impact outreach cases while reducing unnecessary search effort."],
                }
            ]
        }
    }

    result = HeaderPlanner()(state)
    plan = result["header_plan"]

    assert plan["route"] == "centered"
    assert plan["title"]["alignment"] == "center"
    assert plan["subtitle"]["text"]
    assert plan["validation"]["passed"]


def test_header_planner_keeps_long_landscape_title_on_one_readable_line(tmp_path):
    conf_path = tmp_path / "conference.png"
    aff_path = tmp_path / "affiliation.png"
    Image.new("RGBA", (900, 420), (20, 80, 160, 255)).save(conf_path)
    Image.new("RGBA", (700, 700), (160, 40, 60, 255)).save(aff_path)
    long_title = "Audio-Assisted Face Video Restoration with Temporal and Identity Complementary Learning"
    state = create_state(
        str(tmp_path / "paper.pdf"),
        layout_template="cluster_43_landscape",
        width=54,
        height=27,
        logo_path=str(conf_path),
        aff_logo_path=str(aff_path),
        header_route="classic_left",
        header_subtitle_policy="off",
    )
    state["affiliation_logos"] = [
        {"institution": "Test University", "logo_path": str(aff_path), "domain": None, "source": "test", "aspect": 1.0}
    ]
    state["narrative_content"] = {"meta": {"poster_title": long_title, "authors": "A. One, B. Two, C. Three"}}

    plan = HeaderPlanner()(state)["header_plan"]
    title = plan["title"]

    single_line_size = HeaderPlanner()._fit_single_line_font_size(
        title["text"], plan["title_box"]["w"], 100, {"orientation": "landscape"}
    )
    assert title["wrap_policy"] == "single_line"
    assert title["single_line"] is True
    assert "\n" not in title["display_text"]
    assert title["font_size"] == single_line_size
    assert title["font_size"] >= 42
    logo_left = min(
        element["x"]
        for element in plan["logo_elements"]
        if element["type"] != "logo_divider"
    )
    title_right = plan["title_box"]["x"] + plan["title_box"]["w"]
    assert 0.20 <= logo_left - title_right <= 0.35
    assert plan["validation"]["passed"]


def test_header_planner_expands_realistic_long_title_toward_logo_zone(tmp_path):
    conf_path = tmp_path / "conference.png"
    aff_path = tmp_path / "affiliation.png"
    Image.new("RGBA", (1976, 645), (20, 80, 160, 255)).save(conf_path)
    Image.new("RGBA", (1200, 430), (10, 40, 110, 255)).save(aff_path)
    state = create_state(
        str(tmp_path / "paper.pdf"),
        layout_template="cluster_43_landscape",
        width=54,
        height=27,
        logo_path=str(conf_path),
        aff_logo_path=str(aff_path),
        header_route="classic_left",
        header_subtitle_policy="off",
    )
    state["narrative_content"] = {
        "meta": {
            "poster_title": "3DTopia-XL: Scaling High-quality 3D Asset Generation via Primitive Diffusion",
            "authors": "Z. Chen, J. Tang, Y. Dong, Z. Cao, F. Hong, Y. Lan, T. Wang",
        }
    }
    plan = HeaderPlanner()(state)["header_plan"]

    assert plan["title"]["single_line"] is True
    assert plan["title"]["font_size"] >= 64
    assert plan["title"]["font_size"] > plan["authors"]["font_size"]
    logo_left = min(
        element["x"]
        for element in plan["logo_elements"]
        if element["type"] != "logo_divider"
    )
    title_right = plan["title_box"]["x"] + plan["title_box"]["w"]
    assert 0.20 <= logo_left - title_right <= 0.35
    assert plan["validation"]["passed"]


def test_header_planner_wraps_only_when_single_line_would_be_too_small(tmp_path):
    state = create_state(
        str(tmp_path / "paper.pdf"),
        layout_template="cluster_43_landscape",
        width=54,
        height=27,
        header_route="classic_left",
        header_subtitle_policy="off",
    )
    state["narrative_content"] = {
        "meta": {
            "poster_title": (
                "A Unified Generalizable and Controllable Framework for High-Fidelity "
                "Multimodal Scientific Reasoning Generation Evaluation and Deployment "
                "Across Open-World Interactive Environments with Reliable Human Feedback"
            ),
            "authors": "A. One, B. Two, C. Three",
        }
    }

    title = HeaderPlanner()(state)["header_plan"]["title"]

    assert title["wrap_policy"] == "two_line"
    assert title["single_line"] is False
    assert "\n" in title["display_text"]


def test_header_planner_auto_defaults_to_stable_classic_left_without_seed(tmp_path):
    routes = []
    subtitles = []
    for _ in range(8):
        state = create_state(
            str(tmp_path / "paper.pdf"),
            layout_template="cluster_43_landscape",
            width=54,
            height=27,
            header_route="auto",
            header_subtitle_policy="auto",
        )
        state["narrative_content"] = {
            "meta": {
                "poster_title": "Fast Spatial Search",
                "authors": "A. Researcher and B. Scientist",
            }
        }
        state["story_board"] = {
            "spatial_content_plan": {
                "sections": [
                    {
                        "section_id": "motivation",
                        "section_title": "Motivation",
                        "content_role": "overview",
                        "text_content": ["Prioritizes high-impact outreach cases while reducing unnecessary search effort."],
                    }
                ]
            }
        }

        plan = HeaderPlanner()(state)["header_plan"]
        routes.append(plan["route"])
        subtitles.append(plan["subtitle"]["text"])

    assert routes == ["classic_left"] * 8
    assert len(set(subtitles)) == 1


def test_header_planner_auto_subtitle_stays_off_for_long_paper_title(tmp_path):
    class AlwaysSubtitleRng:
        def random(self):
            return 0.0

    state = create_state(
        str(tmp_path / "paper.pdf"),
        layout_template="cluster_43_landscape",
        width=54,
        height=27,
        header_subtitle_policy="auto",
    )
    title = "Active Geospatial Search for Efficient Tenant Eviction Outreach"
    state["story_board"] = {
        "spatial_content_plan": {
            "sections": [
                {
                    "section_id": "motivation",
                    "section_title": "Motivation",
                    "content_role": "overview",
                    "text_content": ["Eviction outreach requires sequential, budget-constrained search under uncertainty."],
                }
            ]
        }
    }

    subtitle = HeaderPlanner()._select_subtitle(state, title, AlwaysSubtitleRng())

    assert subtitle == ""


def test_header_planner_shortens_subtitle_to_complete_clause():
    agent = HeaderPlanner()

    subtitle = agent._shorten_subtitle(
        "Eviction-prevention outreach faces a sequential search problem: canvassers have limited budget, uncertain current risk labels, and travel costs.",
        86,
    )

    assert subtitle == "Eviction-prevention outreach faces a sequential search problem"


def test_header_planner_uses_readable_portrait_subtitle_size(tmp_path):
    conf_path = tmp_path / "conference.png"
    aff_path = tmp_path / "affiliation.png"
    Image.new("RGBA", (1200, 360), (20, 80, 160, 255)).save(conf_path)
    Image.new("RGBA", (900, 900), (150, 20, 35, 255)).save(aff_path)
    state = create_state(
        str(tmp_path / "paper.pdf"),
        layout_template="cluster_13_portrait",
        width=36,
        height=50.88,
        logo_path=str(conf_path),
        aff_logo_path=str(aff_path),
        header_route="split_logos",
        header_subtitle_policy="always",
    )
    state["narrative_content"] = {
        "meta": {
            "poster_title": "Active Geospatial Search for Efficient Tenant Eviction Outreach",
            "authors": "A. Sarkar, A. DiChristofano, S. Das, P.J. Fowler, N. Jacobs, Y. Vorobeychik",
        }
    }
    state["story_board"] = {
        "spatial_content_plan": {
            "sections": [
                {
                    "section_id": "motivation",
                    "section_title": "Search Problem",
                    "content_role": "foundation",
                    "text_content": [
                        "Eviction-prevention outreach faces a sequential search problem: canvassers have limited budget, uncertain current risk labels, and travel costs."
                    ],
                }
            ]
        }
    }

    plan = HeaderPlanner()(state)["header_plan"]

    assert plan["subtitle"]["text"] == "Eviction-prevention outreach faces a sequential search problem"
    assert plan["subtitle"]["font_size"] >= plan["title"]["font_size"] * 0.58
    expected_author_gap = load_config()["header_planner"]["portrait_title_author_gap_inches"]
    assert plan["authors"]["top_gap_inches"] == pytest.approx(expected_author_gap)


def test_header_planner_compacts_long_portrait_subtitle(tmp_path):
    conf_path = tmp_path / "conference.png"
    aff_path = tmp_path / "affiliation.png"
    Image.new("RGBA", (1200, 360), (20, 80, 160, 255)).save(conf_path)
    Image.new("RGBA", (900, 900), (150, 20, 35, 255)).save(aff_path)
    state = create_state(
        str(tmp_path / "paper.pdf"),
        layout_template="cluster_13_portrait",
        width=36,
        height=50.88,
        logo_path=str(conf_path),
        aff_logo_path=str(aff_path),
        header_route="split_logos",
        header_subtitle_policy="always",
    )
    state["narrative_content"] = {
        "meta": {
            "poster_title": "Active Geospatial Search for Efficient Tenant Eviction Outreach",
            "authors": "A. Sarkar, A. DiChristofano, S. Das, P.J. Fowler, N. Jacobs, Y. Vorobeychik",
        }
    }
    state["story_board"] = {
        "spatial_content_plan": {
            "sections": [
                {
                    "section_id": "motivation",
                    "section_title": "Search Problem",
                    "content_role": "foundation",
                    "text_content": [
                        "Tenant eviction outreach is budget-limited and must sequentially decide which properties to canvas without knowing current eviction risk in advance."
                    ],
                }
            ]
        }
    }

    plan = HeaderPlanner()(state)["header_plan"]

    assert plan["subtitle"]["text"] == "Tenant eviction outreach is budget-limited"
    assert plan["subtitle"]["font_size"] >= 42


def test_header_planner_balances_wrapped_author_lines(tmp_path):
    conf_path = tmp_path / "conference.png"
    aff_path = tmp_path / "affiliation.png"
    Image.new("RGBA", (1200, 360), (20, 80, 160, 255)).save(conf_path)
    Image.new("RGBA", (900, 900), (150, 20, 35, 255)).save(aff_path)
    state = create_state(
        str(tmp_path / "paper.pdf"),
        layout_template="cluster_13_portrait",
        width=36,
        height=50.88,
        logo_path=str(conf_path),
        aff_logo_path=str(aff_path),
        header_route="split_logos",
        header_subtitle_policy="always",
    )
    state["narrative_content"] = {
        "meta": {
            "poster_title": "Active Geospatial Search for Efficient Tenant Eviction Outreach",
            "authors": "A. Sarkar, A. DiChristofano, S. Das, P.J. Fowler, N. Jacobs, Y. Vorobeychik",
        }
    }
    state["story_board"] = {
        "spatial_content_plan": {
            "sections": [
                {
                    "section_id": "motivation",
                    "section_title": "Search Problem",
                    "content_role": "foundation",
                    "text_content": ["Tenant eviction outreach is budget-limited and must sequentially decide which properties to canvas."],
                }
            ]
        }
    }

    plan = HeaderPlanner()(state)["header_plan"]
    author_lines = plan["authors"]["text"].splitlines()

    assert len(author_lines) == 2
    assert abs(len(author_lines[0]) - len(author_lines[1])) < 16
    assert not author_lines[1].strip().startswith("Y. Vorobeychik")


def test_header_block_reviewer_repairs_small_subtitle_and_author_gap(tmp_path):
    state = create_state(
        str(tmp_path / "paper.pdf"),
        layout_template="cluster_13_portrait",
        width=36,
        height=50.88,
    )
    state["output_dir"] = str(tmp_path / "output")
    state["styled_layout"] = [
        {
            "type": "title",
            "x": 7.2,
            "y": 2.15,
            "width": 19.55,
            "height": 7.0,
            "title_text": "Active Geospatial Search for\nEfficient Tenant Eviction Outreach",
            "subtitle_text": "Eviction-prevention outreach faces a sequential search problem",
            "authors_text": "A. Sarkar, A. DiChristofano, S. Das",
            "content": "Active Geospatial Search for\nEfficient Tenant Eviction Outreach\nEviction-prevention outreach faces a sequential search problem\nA. Sarkar, A. DiChristofano, S. Das",
            "font_size": 67,
            "subtitle_font_size": 29,
            "subtitle_box_height": 0.45,
            "title_box_height": 1.96,
            "title_to_subtitle_gap_inches": 0.08,
            "author_font_size": 56,
            "author_box_height": 1.62,
            "author_top_gap_inches": 0.18,
        }
    ]

    result = HeaderBlockReviewer()(state)
    title = result["styled_layout"][0]

    assert result["header_block_patch_applied"] is True
    assert title["subtitle_font_size"] > 38
    assert title["subtitle_box_height"] > 0.60
    expected_author_gap = load_config()["header_planner"]["portrait_title_author_gap_inches"]
    assert title["author_top_gap_inches"] == pytest.approx(expected_author_gap)
    assert Path(state["output_dir"], "content", "header_block_review.json").exists()


def test_header_block_reviewer_applies_vlm_author_spacing_patch(tmp_path, monkeypatch):
    state = create_state(
        str(tmp_path / "paper.pdf"),
        layout_template="cluster_13_portrait",
        width=36,
        height=50.88,
    )
    state["output_dir"] = str(tmp_path / "output")
    initial_gap = load_config()["header_planner"]["portrait_title_author_gap_inches"]
    state["styled_layout"] = [
        {
            "type": "title",
            "x": 7.2,
            "y": 2.15,
            "width": 19.55,
            "height": 7.0,
            "title_text": "Active Geospatial Search for\nEfficient Tenant Eviction Outreach",
            "subtitle_text": "Eviction-prevention outreach faces a sequential search problem",
            "authors_text": "A. Sarkar, A. DiChristofano, S. Das",
            "content": "Active Geospatial Search for\nEfficient Tenant Eviction Outreach\nEviction-prevention outreach faces a sequential search problem\nA. Sarkar, A. DiChristofano, S. Das",
            "font_size": 67,
            "subtitle_font_size": 42,
            "subtitle_box_height": 0.66,
            "title_box_height": 1.96,
            "title_to_subtitle_gap_inches": 0.10,
            "author_font_size": 56,
            "author_box_height": 1.62,
            "author_top_gap_inches": initial_gap,
        }
    ]

    monkeypatch.setattr(
        HeaderBlockReviewer,
        "_review_with_vlm",
        lambda self, state, crop_path: {
            "status": "ok",
            "issues": [{"severity": "medium", "category": "author_spacing", "description": "authors too close"}],
            "recommendations": [{"target": "authors", "action": "move_down", "reason": "increase separation"}],
        },
    )

    result = HeaderBlockReviewer()(state)
    title = result["styled_layout"][0]

    assert result["header_block_patch_applied"] is True
    assert title["author_top_gap_inches"] > initial_gap
    assert title["author_top_gap_inches"] <= load_config()["header_block_review"]["max_author_gap_inches"]
    assert any(patch.get("source") == "vlm_header_review" for patch in result["header_block_review"]["patch"])


def test_layout_agent_uses_header_plan_for_title_and_logo_elements(tmp_path):
    conf_path = tmp_path / "conference.png"
    aff_path = tmp_path / "affiliation.png"
    Image.new("RGBA", (900, 420), (20, 80, 160, 255)).save(conf_path)
    Image.new("RGBA", (700, 240), (160, 40, 60, 255)).save(aff_path)

    state = create_state(
        str(tmp_path / "paper.pdf"),
        layout_template="cluster_43_landscape",
        width=54,
        height=27,
        logo_path=str(conf_path),
        aff_logo_path=str(aff_path),
        header_route="right_title",
        header_subtitle_policy="off",
    )
    state["narrative_content"] = {
        "meta": {
            "poster_title": "Header Planning for Posters",
            "authors": "A. Researcher",
        }
    }

    state = HeaderPlanner()(state)
    agent = LayoutAgent()
    template = agent._resolve_template_layout(state)
    title = agent._create_title_element(state, state["poster_width"], template["header"]["h"], template)
    logos = agent._create_logo_elements(state, state["poster_width"], template)

    assert title["alignment"] == "right"
    assert title["lock_header_typography"] is True
    assert {logo["type"] for logo in logos} >= {"conf_logo", "institution_logo"}
    assert all(logo["x"] + logo["width"] <= title["x"] for logo in logos if logo["type"] != "logo_divider")


def test_header_planner_centered_route_splits_affiliation_and_conference_logos(tmp_path):
    conf_path = tmp_path / "conference.png"
    aff_path = tmp_path / "affiliation.png"
    Image.new("RGBA", (900, 420), (20, 80, 160, 255)).save(conf_path)
    Image.new("RGBA", (700, 240), (160, 40, 60, 255)).save(aff_path)

    state = create_state(
        str(tmp_path / "paper.pdf"),
        layout_template="cluster_43_landscape",
        width=54,
        height=27,
        logo_path=str(conf_path),
        aff_logo_path=str(aff_path),
        header_route="centered",
        header_subtitle_policy="off",
    )
    state["narrative_content"] = {
        "meta": {
            "poster_title": "Centered Header Planning",
            "authors": "A. Researcher",
        }
    }

    plan = HeaderPlanner()(state)["header_plan"]
    title_box = plan["title_box"]
    aff_logo = next(element for element in plan["logo_elements"] if element["type"] == "institution_logo")
    conf_logo = next(element for element in plan["logo_elements"] if element["type"] == "conf_logo")

    assert plan["route"] == "centered"
    assert plan["title"]["alignment"] == "center"
    assert aff_logo["x"] + aff_logo["width"] < title_box["x"]
    assert conf_logo["x"] > title_box["x"] + title_box["w"]
    assert plan["validation"]["passed"]


def test_header_planner_maps_portrait_routes_to_logo_strip_full_title(tmp_path):
    conf_path = tmp_path / "conference.png"
    manual_aff_path = tmp_path / "manual_affiliation.png"
    auto_aff_paths = []
    Image.new("RGBA", (900, 420), (20, 80, 160, 255)).save(conf_path)
    Image.new("RGBA", (700, 700), (160, 40, 60, 255)).save(manual_aff_path)
    for index in range(3):
        path = tmp_path / f"auto_affiliation_{index}.png"
        Image.new("RGBA", (900, 240), (40 + index * 30, 90, 150, 255)).save(path)
        auto_aff_paths.append(path)

    state = create_state(
        str(tmp_path / "paper.pdf"),
        layout_template="cluster_8_portrait",
        width=27,
        height=54,
        logo_path=str(conf_path),
        aff_logo_path=str(manual_aff_path),
        header_route="right_title",
        header_subtitle_policy="off",
    )
    state["affiliation_logos"] = [
        {
            "institution": f"Auto Institution {index}",
            "logo_path": str(path),
            "domain": None,
            "source": "test",
            "aspect": 3.75,
        }
        for index, path in enumerate(auto_aff_paths)
    ]
    state["narrative_content"] = {
        "meta": {
            "poster_title": "Active Geospatial Search for Efficient Tenant Eviction Outreach",
            "authors": "A. Researcher, B. Scientist, and C. Collaborator",
        }
    }

    result = HeaderPlanner()(state)
    plan = result["header_plan"]
    header = result["layout_template_metadata"]["header"]
    title_box = plan["title_box"]
    institution_logos = [element for element in plan["logo_elements"] if element["type"] == "institution_logo"]
    non_divider_logos = [element for element in plan["logo_elements"] if element["type"] != "logo_divider"]

    assert plan["route"] == "right_title"
    assert plan["physical_route"] == "portrait_logo_strip_full_title"
    assert plan["title"]["alignment"] == "right"
    assert plan["title"]["single_line"] is True
    assert plan["title"]["font_size"] < 58
    estimated_title_width = len(plan["title"]["text"]) * (plan["title"]["font_size"] / 72) * 0.56
    assert estimated_title_width <= title_box["w"] * 0.95
    assert title_box["w"] == pytest.approx(header["w"])
    assert title_box["y"] > header["y"]
    assert len(institution_logos) == 1
    assert institution_logos[0]["source"] == "manual"
    assert plan["authors"]["font_size"] >= 38
    assert all(element["y"] + element["height"] <= title_box["y"] for element in non_divider_logos)
    assert plan["validation"]["passed"]


def test_header_planner_supports_portrait_split_logos_with_two_line_title(tmp_path):
    conf_path = tmp_path / "conference.png"
    manual_aff_path = tmp_path / "affiliation.png"
    auto_aff_path = tmp_path / "auto-affiliation.png"
    Image.new("RGBA", (1200, 360), (20, 80, 160, 255)).save(conf_path)
    Image.new("RGBA", (900, 900), (150, 20, 35, 255)).save(manual_aff_path)
    Image.new("RGBA", (900, 240), (30, 110, 80, 255)).save(auto_aff_path)

    state = create_state(
        str(tmp_path / "paper.pdf"),
        layout_template="cluster_22_portrait",
        width=27,
        height=54,
        logo_path=str(conf_path),
        aff_logo_path=str(manual_aff_path),
        header_route="split_logos",
        header_subtitle_policy="off",
        header_title_wrap_policy="two_line",
    )
    state["affiliation_logos"] = [
        {
            "institution": "Auto Institution",
            "logo_path": str(auto_aff_path),
            "domain": None,
            "source": "test",
            "aspect": 3.75,
        }
    ]
    state["narrative_content"] = {
        "meta": {
            "poster_title": "Active Geospatial Search for Efficient Tenant Eviction Outreach",
            "authors": "A. Researcher, B. Scientist, and C. Collaborator",
        }
    }

    result = HeaderPlanner()(state)
    plan = result["header_plan"]
    title_box = plan["title_box"]
    institution_logos = [element for element in plan["logo_elements"] if element["type"] == "institution_logo"]
    conference_logos = [element for element in plan["logo_elements"] if element["type"] == "conf_logo"]

    assert plan["route"] == "split_logos"
    assert plan["physical_route"] == "portrait_split_logos_title_center"
    assert plan["title"]["alignment"] == "center"
    assert plan["title"]["single_line"] is False
    assert plan["title"]["wrap_policy"] == "two_line"
    assert "\n" in plan["title"]["display_text"]
    assert plan["title"]["font_size"] >= 68
    assert plan["title"]["font_family"] == "Georgia"
    assert plan["title"]["box_height"] <= 2.0
    assert plan["authors"]["font_size"] >= 56
    expected_author_gap = load_config()["header_planner"]["portrait_title_author_gap_inches"]
    assert plan["authors"]["top_gap_inches"] == pytest.approx(expected_author_gap)
    assert plan["authors"]["word_wrap"] is True
    assert plan["authors"]["x"] == title_box["x"]
    assert plan["authors"]["w"] == title_box["w"]
    assert len(institution_logos) == 1
    assert institution_logos[0]["image_path"] == str(manual_aff_path)
    assert institution_logos[0]["source"] == "manual"
    assert len(conference_logos) == 1
    assert conference_logos[0]["height"] >= 1.2
    assert institution_logos[0]["x"] + institution_logos[0]["width"] < title_box["x"]
    assert title_box["x"] + title_box["w"] < conference_logos[0]["x"]
    assert plan["validation"]["passed"]


def test_header_planner_boosts_affiliation_logo_when_safe(tmp_path):
    conf_path = tmp_path / "conference.png"
    aff_path = tmp_path / "affiliation.png"
    Image.new("RGBA", (900, 420), (20, 80, 160, 255)).save(conf_path)
    Image.new("RGBA", (700, 700), (160, 40, 60, 255)).save(aff_path)

    state = create_state(
        str(tmp_path / "paper.pdf"),
        layout_template="cluster_43_landscape",
        width=54,
        height=27,
        logo_path=str(conf_path),
        aff_logo_path=str(aff_path),
        header_route="classic_left",
        header_subtitle_policy="off",
    )
    state["narrative_content"] = {
        "meta": {
            "poster_title": "Header Logo Sizing",
            "authors": "A. Researcher",
        }
    }

    plan = HeaderPlanner()(state)["header_plan"]
    institution_logo = next(element for element in plan["logo_elements"] if element["type"] == "institution_logo")

    assert plan["validation"]["passed"]
    assert plan["logo_resize_decision"] == "boosted_affiliation_logo"
    assert plan["affiliation_logo_scale"] > 1.0
    assert institution_logo["height"] > 2.0
    assert institution_logo["width"] == pytest.approx(institution_logo["height"], rel=0.05)
    assert [attempt["label"] for attempt in plan["logo_resize_attempts"]] == [
        "base",
        "boosted_affiliation_logo",
    ]


def test_header_planner_reverts_to_base_logo_when_boost_is_too_large(tmp_path):
    conf_path = tmp_path / "conference.png"
    aff_path = tmp_path / "affiliation.png"
    Image.new("RGBA", (900, 420), (20, 80, 160, 255)).save(conf_path)
    Image.new("RGBA", (700, 700), (160, 40, 60, 255)).save(aff_path)

    state = create_state(
        str(tmp_path / "paper.pdf"),
        layout_template="cluster_43_landscape",
        width=54,
        height=27,
        logo_path=str(conf_path),
        aff_logo_path=str(aff_path),
        header_route="classic_left",
        header_subtitle_policy="off",
    )
    state["narrative_content"] = {
        "meta": {
            "poster_title": "Header Logo Sizing",
            "authors": "A. Researcher",
        }
    }
    agent = HeaderPlanner()
    agent.header_config["preferred_affiliation_logo_scale"] = 9.0
    agent.header_config["hard_max_logo_header_fraction"] = 0.70

    plan = agent(state)["header_plan"]
    institution_logo = next(element for element in plan["logo_elements"] if element["type"] == "institution_logo")
    boosted_attempt = plan["logo_resize_attempts"][1]

    assert plan["validation"]["passed"]
    assert plan["logo_resize_decision"] == "base_after_boost_rejected"
    assert plan["affiliation_logo_scale"] == 1.0
    assert institution_logo["height"] == pytest.approx(1.95)
    assert boosted_attempt["passed"] is False


def test_font_agent_preserves_header_plan_typography():
    agent = FontAgent()
    state = create_state("/tmp/paper.pdf", poster_style_preset="teal_modern")
    element = {
        "type": "title",
        "content": "Short Title\nAuthors",
        "font_size": 88,
        "author_font_size": 42,
        "subtitle_font_size": 35,
        "alignment": "center",
        "lock_header_typography": True,
    }

    agent._apply_title_styling(element, {"text_on_theme": "#000000"}, state)

    assert element["font_size"] == 88
    assert element["author_font_size"] == 42
    assert element["subtitle_font_size"] == 35
    assert element["alignment"] == "center"


def test_font_agent_applies_portrait_header_wordart_override():
    agent = FontAgent()
    state = create_state("/tmp/paper.pdf", width=27, height=54, poster_style_preset="teal_modern")
    element = {
        "type": "title",
        "content": "Short Title\nAuthors",
        "font_size": 68,
        "author_font_size": 42,
        "header_route": "split_logos",
        "lock_header_typography": True,
    }

    agent._apply_title_styling(element, {"text_on_theme": "#000000"}, state)

    assert element["font_family"] == "Georgia"
    assert element["font_color"] == "#4A1020"
    assert element["main_title_style_override"]["font_family"] == "Georgia"
    assert element["main_title_style_override"]["author_font_color"] == "#211517"


def test_layout_agent_section_title_uses_navy_band_wordart():
    agent = LayoutAgent()
    state = create_state("/tmp/paper.pdf")
    section = {
        "section_id": "method",
        "section_title": "Method",
        "column_assignment": "slot_3",
        "slot_id": "slot_3",
    }

    elements = agent._create_section_title_design(section, column_x=1.0, start_y=2.0, column_width=6.0, state=state)

    bar = next(element for element in elements if element["type"] == "title_accent_block")
    title = next(element for element in elements if element["type"] == "section_title")
    assert bar["x"] == pytest.approx(1.0)
    assert bar["width"] == pytest.approx(6.0)
    assert bar["color"] == "#06134A"
    assert title["section_title"] == "Method"
    assert title["section_number"] is None
    assert title["section_numbering_mode"] == "off"
    assert title["font_family"] == "Georgia"
    assert title["font_color"] == "#FFFFFF"
    assert title["wordart_style"]["name"] == "navy_band_serif"


def test_layout_agent_section_title_supports_small_optional_numbering():
    agent = LayoutAgent()
    state = create_state("/tmp/paper.pdf", section_title_numbering="small")
    section = {
        "section_id": "method",
        "section_title": "Method",
        "column_assignment": "slot_3",
        "slot_id": "slot_3",
    }

    elements = agent._create_section_title_design(section, column_x=1.0, start_y=2.0, column_width=6.0, state=state)

    title = next(element for element in elements if element["type"] == "section_title")
    assert title["section_title"] == "Method"
    assert title["section_number"] == "3"
    assert title["section_numbering_mode"] == "small"
    assert title["section_number_font_scale"] < 1.0


def test_layout_agent_section_title_inline_numbering_keeps_legacy_label():
    agent = LayoutAgent()
    state = create_state("/tmp/paper.pdf", section_title_numbering="inline")
    section = {
        "section_id": "method",
        "section_title": "Method",
        "column_assignment": "slot_3",
        "slot_id": "slot_3",
    }

    elements = agent._create_section_title_design(section, column_x=1.0, start_y=2.0, column_width=6.0, state=state)

    title = next(element for element in elements if element["type"] == "section_title")
    assert title["section_title"] == "3. Method"
    assert title["section_numbering_mode"] == "inline"


def test_layout_agent_section_title_default_strips_existing_number_prefix():
    agent = LayoutAgent()
    state = create_state("/tmp/paper.pdf")
    section = {
        "section_id": "method",
        "section_title": "7. Method",
        "column_assignment": "slot_3",
        "slot_id": "slot_3",
    }

    elements = agent._create_section_title_design(section, column_x=1.0, start_y=2.0, column_width=6.0, state=state)

    title = next(element for element in elements if element["type"] == "section_title")
    assert title["section_title"] == "Method"
    assert title["section_numbering_mode"] == "off"


def test_layout_agent_portrait_section_title_uses_wordart_font_override():
    agent = LayoutAgent()
    state = create_state("/tmp/paper.pdf", width=27, height=54, poster_style_preset="teal_modern")
    agent._apply_state_style(state)
    state["section_title_design"] = {
        "section_title_design": {
            "section_applications": [
                {
                    "section_id": "method",
                    "title_styling": {"font_family": "Helvetica Neue", "color": "#FFFFFF"},
                    "accent_styling": {"color": "#0B4F5C"},
                }
            ]
        }
    }
    section = {
        "section_id": "method",
        "section_title": "Method",
        "column_assignment": "slot_3",
        "slot_id": "slot_3",
    }

    elements = agent._create_section_title_design(section, column_x=1.0, start_y=2.0, column_width=6.0, state=state)

    bar = next(element for element in elements if element["type"] == "title_accent_block")
    title = next(element for element in elements if element["type"] == "section_title")
    assert bar["color"] == "#0B4F5C"
    assert title["font_family"] == "Georgia"
    assert title["font_color"] == "#FFFFFF"
    assert title["wordart_style"]["shadow"]["color"] == "#9FB7BC"


def test_section_title_designer_emits_navy_band_template():
    state = create_state("/tmp/paper.pdf")
    state["story_board"] = {
        "spatial_content_plan": {
            "sections": [
                {"section_id": "method", "section_title": "Method"},
            ]
        }
    }
    state["color_scheme"] = {"theme": "#335f91", "mono_light": "#8AA0BA", "mono_dark": "#001E44"}

    result = SectionTitleDesigner()(state)

    design = result["section_title_design"]["section_title_design"]
    assert design["selected_template"] == "navy_band_wordart"
    application = design["section_applications"][0]
    assert application["accent_styling"]["type"] == "full_width_bar"
    assert application["title_styling"]["font_family"] == "Georgia"


def test_section_title_designer_uses_selected_style_preset():
    state = create_state("/tmp/paper.pdf", poster_style_preset="teal_modern")
    state["story_board"] = {
        "spatial_content_plan": {
            "sections": [
                {"section_id": "method", "section_title": "Method"},
            ]
        }
    }
    state["color_scheme"] = {"theme": "#335f91", "mono_light": "#8AA0BA", "mono_dark": "#001E44"}

    result = SectionTitleDesigner()(state)

    application = result["section_title_design"]["section_title_design"]["section_applications"][0]
    assert application["accent_styling"]["color"] == "#0B4F5C"
    assert application["title_styling"]["font_family"] == "Helvetica Neue"


def test_layout_agent_uses_manual_aff_logo_with_conference_logo(tmp_path):
    conf_path = tmp_path / "conference.png"
    aff_path = tmp_path / "affiliation.png"
    Image.new("RGBA", (900, 463), (20, 80, 160, 255)).save(conf_path)
    Image.new("RGBA", (1200, 1200), (160, 40, 60, 255)).save(aff_path)

    state = create_state("/tmp/paper.pdf", layout_template="cluster_3_portrait", width=36, height=50.88)
    state["logo_path"] = str(conf_path)
    state["aff_logo_path"] = str(aff_path)

    agent = LayoutAgent()
    template = agent._resolve_template_layout(state)
    elements = agent._create_logo_elements(state, state["poster_width"], template)

    assert {element["type"] for element in elements} == {
        "conf_logo",
        "institution_logo",
        "logo_divider",
    }
    institution_logo = next(element for element in elements if element["type"] == "institution_logo")
    conference_logo = next(element for element in elements if element["type"] == "conf_logo")
    assert institution_logo["image_path"] == str(aff_path)
    assert institution_logo["source"] == "manual"
    assert institution_logo["height"] >= 1.8
    assert conference_logo["height"] >= 2.4


def test_layout_agent_limits_affiliation_logos_to_three_without_overflow(tmp_path):
    logo_paths = []
    for index in range(5):
        path = tmp_path / f"logo_{index}.png"
        Image.new("RGBA", (420, 120), (20, 50, 120, 255)).save(path)
        logo_paths.append(str(path))

    state = create_state("/tmp/paper.pdf", enable_affiliation_logos=True)
    state["affiliation_logos"] = [
        {
            "institution": f"Institution {index}",
            "logo_path": path,
            "domain": None,
            "source": "test",
            "aspect": 3.5,
        }
        for index, path in enumerate(logo_paths)
    ]

    agent = LayoutAgent()
    template = agent._resolve_template_layout(state)
    elements = agent._create_logo_elements(state, state["poster_width"], template)
    logos = [element for element in elements if element.get("type") == "institution_logo"]
    region = agent._title_logo_region(template, False)

    assert len(logos) == 3
    assert all(region["x"] <= logo["x"] for logo in logos)
    assert all(logo["x"] + logo["width"] <= region["x"] + region["w"] + 1e-6 for logo in logos)
    assert all(region["y"] <= logo["y"] for logo in logos)
    assert all(logo["y"] + logo["height"] <= region["y"] + region["h"] + 1e-6 for logo in logos)


def test_micro_layout_refiner_keeps_logo_divider_as_global_element():
    state = create_state("/tmp/paper.pdf", layout_template="three_column_postergen")
    state["styled_layout"] = [
        {
            "type": "title",
            "x": 1.0,
            "y": 1.0,
            "width": 32.0,
            "height": 4.0,
            "priority": 1.0,
        },
        {
            "type": "logo_divider",
            "x": 42.0,
            "y": 1.5,
            "width": 0.04,
            "height": 3.5,
            "priority": 0.85,
        },
        {
            "type": "section_container",
            "section_id": "left::s1",
            "lane_id": "left",
            "x": 1.0,
            "y": 7.0,
            "width": 16.0,
            "height": 4.0,
            "priority": 0.1,
        },
        {
            "type": "text",
            "id": "left::s1_text",
            "section_id": "left::s1",
            "x": 1.3,
            "y": 7.6,
            "width": 15.4,
            "height": 2.5,
            "content": "Short text",
            "font_size": 40,
            "priority": 0.5,
        },
    ]

    result = MicroLayoutRefiner()(state)
    divider = next(element for element in result["styled_layout"] if element.get("type") == "logo_divider")

    assert divider["y"] == 1.5


def test_renderer_uses_element_path_for_institution_logo(tmp_path):
    logo_path = tmp_path / "logo.png"
    Image.new("RGBA", (100, 80), (255, 255, 255, 255)).save(logo_path)

    calls = []

    class FakeDirector:
        def add_image(self, *args, **kwargs):
            calls.append((args, kwargs))

    renderer = Renderer()
    renderer.director = FakeDirector()
    element = {
        "type": "institution_logo",
        "x": 1,
        "y": 1,
        "width": 2,
        "height": 1,
        "image_path": str(logo_path),
    }

    renderer._render_institution_logo(None, element, create_state("/tmp/paper.pdf"))

    assert calls
    assert calls[0][0][0] == str(logo_path)


def test_renderer_trims_logo_whitespace(tmp_path):
    logo_path = tmp_path / "logo_with_margin.png"
    image = Image.new("RGBA", (400, 400), (255, 255, 255, 255))
    for x in range(150, 250):
        for y in range(120, 280):
            image.putpixel((x, y), (20, 60, 120, 255))
    image.save(logo_path)

    renderer = Renderer()
    renderer._render_output_dir = tmp_path / "output"
    trimmed_path = Path(renderer._trim_logo_whitespace(str(logo_path)))

    assert trimmed_path.exists()
    with Image.open(trimmed_path) as trimmed:
        assert trimmed.size[0] < 160
        assert trimmed.size[1] < 220


def test_renderer_draws_explicit_section_container_fill_for_standard_templates():
    calls = []

    class FakeDirector:
        def add_shape(self, *args, **kwargs):
            calls.append((args, kwargs))

    renderer = Renderer()
    renderer.director = FakeDirector()
    element = {
        "type": "section_container",
        "x": 1,
        "y": 2,
        "width": 5,
        "height": 4,
        "fill_color": "#F0F6FF",
        "fill_opacity": 0.84,
        "border_color": "#C9DDF5",
        "border_opacity": 0.58,
        "border_width": 0.9,
        "border_style": "dashed",
        "shadow": {"enabled": True, "color": "#000000", "alpha": 0.16},
    }

    renderer._render_section_container(None, element, create_state("/tmp/paper.pdf"))

    assert calls
    assert calls[0][1]["fill_color"] == "#F0F6FF"
    assert calls[0][1]["fill_opacity"] == 0.84
    assert calls[0][1]["border_color"] == "#C9DDF5"
    assert calls[0][1]["border_opacity"] == 0.58
    assert calls[0][1]["border_style"] == "dashed"
    assert calls[0][1]["shadow"] == element["shadow"]


def test_pptx_director_add_shape_applies_outer_shadow_xml():
    director = PPTXDirector()

    shape = director.add_shape(
        MSO_SHAPE.RECTANGLE,
        0.1,
        0.1,
        1.0,
        1.0,
        fill_color="#F1F2F4",
        shadow={"enabled": True, "color": "#000000", "alpha": 0.16, "blur_pt": 5, "distance_pt": 2.4},
    )

    xml = shape._element.xml
    assert "outerShdw" in xml
    assert 'alpha val="16000"' in xml


def test_pptx_director_add_shape_applies_fill_and_border_opacity_xml():
    director = PPTXDirector()

    shape = director.add_shape(
        MSO_SHAPE.RECTANGLE,
        0.1,
        0.1,
        1.0,
        1.0,
        fill_color="#FFFFFF",
        fill_opacity=0.84,
        border_color="#D8DEE7",
        border_opacity=0.58,
        border_width=0.45,
    )

    xml = shape._element.xml
    assert 'alpha val="84000"' in xml
    assert 'alpha val="58000"' in xml


def test_renderer_draws_background_image_before_layout(tmp_path):
    background_path = tmp_path / "background.png"
    Image.new("RGB", (100, 140), (245, 247, 250)).save(background_path)
    calls = []

    class FakeDirector:
        def set_slide_dimensions(self, *args, **kwargs):
            pass

        @property
        def slide(self):
            return None

        def add_image(self, *args, **kwargs):
            calls.append((args, kwargs))

        def save(self, *args, **kwargs):
            pass

    state = create_state("/tmp/paper.pdf", width=36, height=50.88)
    state["background_image_path"] = str(background_path)
    state["styled_layout"] = [
        {
            "type": "title",
            "x": 1,
            "y": 1,
            "width": 10,
            "height": 2,
            "content": "Title",
        }
    ]

    renderer = Renderer()
    renderer.director = FakeDirector()
    renderer._render_background_image(None, state)

    assert calls[0][0][:5] == (str(background_path), 0, 0, 36, 50.88)
    assert calls[0][1]["keep_aspect_ratio"] is False


def test_renderer_separates_title_and_authors_with_physical_gap():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    renderer = Renderer()
    renderer.styling_interfaces = {"font_sizes": {"title": 100, "authors": 72}}
    before_shape_count = len(slide.shapes)
    element = {
        "type": "title",
        "x": 1,
        "y": 1,
        "width": 12,
        "height": 4,
        "content": "A Long Research Poster Title\nA. Author, B. Author",
        "font_size": 100,
        "author_font_size": 72,
    }

    renderer._render_title(slide, element, create_state("/tmp/paper.pdf"))

    assert len(slide.shapes) == before_shape_count + 3
    title_box = slide.shapes[-2]
    author_box = slide.shapes[-1]
    actual_gap = author_box.top.inches - (title_box.top.inches + title_box.height.inches)
    assert title_box.text_frame.word_wrap is False
    assert actual_gap == pytest.approx(16 / 72, abs=0.01)


def test_renderer_main_title_visual_style_overrides_element_color():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    renderer = Renderer()
    renderer.styling_interfaces = {"font_sizes": {"title": 100, "authors": 72}}
    element = {
        "type": "title",
        "x": 1,
        "y": 1,
        "width": 12,
        "height": 4,
        "content": "Styled Title\nA. Author",
        "font_color": "#FFFFFF",
        "font_family": "Helvetica Neue",
    }

    renderer._render_title(slide, element, create_state("/tmp/paper.pdf"))

    title_box = slide.shapes[-2]
    paragraph = title_box.text_frame.paragraphs[0]
    assert paragraph.font.name == "Georgia"
    assert str(paragraph.font.color.rgb) == "07164A"


def test_renderer_allows_portrait_header_wordart_override():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    renderer = Renderer()
    renderer.styling_interfaces = {"font_sizes": {"title": 100, "authors": 72}}
    renderer.visual_style_config = {
        "enabled": True,
        "main_title": {
            "font_family": "Helvetica Neue",
            "font_color": "#0B2F36",
            "author_font_family": "Arial",
            "author_font_color": "#27343A",
        },
    }
    element = {
        "type": "title",
        "x": 1,
        "y": 1,
        "width": 12,
        "height": 4,
        "content": "Styled Title\nA. Author",
        "font_color": "#FFFFFF",
        "font_family": "Helvetica Neue",
        "main_title_style_override": {
            "font_family": "Georgia",
            "font_color": "#4A1020",
            "author_font_color": "#211517",
        },
    }

    renderer._render_title(slide, element, create_state("/tmp/paper.pdf", poster_style_preset="teal_modern"))

    title_box = slide.shapes[-2]
    author_box = slide.shapes[-1]
    assert title_box.text_frame.paragraphs[0].font.name == "Georgia"
    assert str(title_box.text_frame.paragraphs[0].font.color.rgb) == "4A1020"
    assert str(author_box.text_frame.paragraphs[0].font.color.rgb) == "211517"


def test_renderer_tokenizer_consumes_malformed_markdown_marker():
    renderer = Renderer()
    text = "*Ablation:** Tabular records are slightly stronger than imagery alone."

    segments = renderer._tokenize_formatting(text)

    assert "".join(segment["text"] for segment in segments)
    assert any(segment["text"] == "*" for segment in segments)


def test_block_template_registry_exposes_cluster_templates():
    template_ids = set(list_block_template_ids())

    assert {"cluster_2_landscape", "cluster_43_landscape", "cluster_3_portrait", "cluster_8_portrait"}.issubset(template_ids)
    assert "cluster_27_landscape" in template_ids
    assert "cluster_27_portrait" in template_ids


def test_block_template_layout_identifies_header_and_content_slots():
    layout = load_block_template_layout("cluster_8_portrait", 36, 51, margin=1.0)

    assert layout["layout_mode"] == "template_prior"
    assert layout["orientation"] == "portrait"
    assert layout["template_aspect_ratio"] < 1
    assert layout["header_slot"]["slot_id"] == "slot_0"
    assert len(layout["content_slots"]) == 4
    assert len(layout["lanes"]) == 4
    assert all(slot["x"] + slot["w"] <= 36.05 for slot in layout["content_slots"])
    assert all(slot["y"] + slot["h"] <= 51.05 for slot in layout["content_slots"])
    assert all(slot["slot_id"] != layout["header_slot"]["slot_id"] for slot in layout["content_slots"])


def test_new_landscape_template_softly_expands_into_gutters():
    layout = load_block_template_layout("cluster_43_landscape", 54, 27, margin=1.0)
    slot_by_id = {slot["slot_id"]: slot for slot in layout["content_slots"]}
    report = layout["gap_absorption_report"]

    assert report["enabled"] is True
    assert report["mode"] == "equal_split"
    assert len(report["absorptions"]) >= 2
    assert len(report["edge_expansions"]) >= 2
    assert report["total_area_gain"] > 0
    assert any(absorption["orientation"] == "vertical" for absorption in report["absorptions"])
    assert any(absorption["orientation"] == "horizontal" for absorption in report["absorptions"])
    assert any(slot.get("gap_absorbed") is True for slot in slot_by_id.values())
    assert slot_by_id["slot_3"]["x"] + slot_by_id["slot_3"]["w"] <= 53.01


def test_standard_template_soft_geometry_is_enabled():
    layout = load_block_template_layout("cluster_43_landscape", 54, 27, margin=1.0)
    report = layout["gap_absorption_report"]

    assert report["enabled"] is True
    assert report["mode"] == "equal_split"
    assert report["area_after"] > report["area_before"]
    assert report["total_area_gain"] > 0
    assert report["edge_expansions"]
    assert {
        "affected_slot_ids",
        "orientation",
        "original_gap_inches",
        "split_boundary",
        "left_or_upper_gain",
        "right_or_lower_gain",
    }.issubset(report["absorptions"][0])


def test_equal_split_gap_absorption_splits_horizontal_neighbors_at_midline():
    raw = {
        "aspect_ratio": 2.0,
        "slots": [
            {"slot_id": 0, "bbox": [0, 0, 1000, 120]},
            {"slot_id": 1, "bbox": [0, 200, 400, 700]},
            {"slot_id": 2, "bbox": [600, 200, 1000, 700]},
        ],
    }

    layout = build_runtime_template(raw, "cluster_43_landscape", 20, 10, margin=1.0)
    slot_1, slot_2 = layout["content_slots"]
    report = layout["gap_absorption_report"]

    assert report["enabled"] is True
    assert report["mode"] == "equal_split"
    assert len(report["absorptions"]) == 1
    assert report["absorptions"][0]["orientation"] == "vertical"
    assert slot_1["x"] + slot_1["w"] == pytest.approx(slot_2["x"], abs=0.001)
    assert slot_1["w"] > 7.2
    assert slot_2["w"] > 7.2
    assert slot_1["slot_id"] == "slot_1"
    assert slot_2["slot_id"] == "slot_2"


def test_equal_split_gap_absorption_expands_outer_edges_to_safe_bounds():
    layout = load_block_template_layout("cluster_43_landscape", 54, 27, margin=1.0)
    slot_by_id = {slot["slot_id"]: slot for slot in layout["content_slots"]}
    report = layout["gap_absorption_report"]

    assert any(item["edge"] == "bottom" for item in report["edge_expansions"])
    assert any(item["edge"] == "left" for item in report["edge_expansions"])
    assert any(item["edge"] == "right" for item in report["edge_expansions"])
    assert slot_by_id["slot_4"]["x"] == pytest.approx(1.0)
    assert slot_by_id["slot_4"]["y"] + slot_by_id["slot_4"]["h"] == pytest.approx(26.0)
    assert slot_by_id["slot_6"]["x"] + slot_by_id["slot_6"]["w"] == pytest.approx(53.0)
    assert slot_by_id["slot_6"]["y"] + slot_by_id["slot_6"]["h"] == pytest.approx(26.0)


def test_equal_split_gap_absorption_splits_vertical_neighbors_at_midline():
    raw = {
        "aspect_ratio": 1.0,
        "slots": [
            {"slot_id": 0, "bbox": [0, 0, 1000, 120]},
            {"slot_id": 1, "bbox": [0, 200, 1000, 480]},
            {"slot_id": 2, "bbox": [0, 680, 1000, 1000]},
        ],
    }

    layout = build_runtime_template(raw, "cluster_43_landscape", 20, 10, margin=1.0)
    slot_1, slot_2 = layout["content_slots"]
    report = layout["gap_absorption_report"]

    assert report["enabled"] is True
    assert report["mode"] == "equal_split"
    assert len(report["absorptions"]) == 1
    assert report["absorptions"][0]["orientation"] == "horizontal"
    assert slot_1["y"] + slot_1["h"] == pytest.approx(slot_2["y"], abs=0.001)
    assert slot_1["h"] > 2.24
    assert slot_2["h"] > 2.56
    assert slot_1["slot_id"] == "slot_1"
    assert slot_2["slot_id"] == "slot_2"


def test_layout_agent_selects_single_primary_block_background():
    state = create_state("/tmp/paper.pdf", layout_template="cluster_3_portrait", width=36, height=50.88)
    state["color_scheme"] = {"theme": "#0057B8", "contrast": "#7F4B13"}
    agent = LayoutAgent()
    agent.config["selective_block_backgrounds"]["enabled"] = True
    agent.config["selective_block_backgrounds"]["max_highlight_blocks"] = 1
    template = agent._resolve_template_layout(state)
    column_assignments = [
        {
            "column_name": "slot_1",
            "sections": [
                {
                    "section_id": "stable_framework",
                    "section_title": "Stable Framework",
                    "content_role": "method",
                    "visual_assets": [{"visual_id": "figure_2"}],
                }
            ],
        },
        {
            "column_name": "slot_2",
            "sections": [
                {
                    "section_id": "main_results",
                    "section_title": "Main Results",
                    "content_role": "results",
                    "visual_assets": [{"visual_id": "figure_3"}],
                }
            ],
        },
        {
            "column_name": "slot_3",
            "sections": [
                {
                    "section_id": "background",
                    "section_title": "Background",
                    "content_role": "overview",
                    "visual_assets": [],
                }
            ],
        },
    ]

    selected = agent._select_highlight_section_ids(column_assignments, state, template)

    assert selected == {"main_results": 0}


def test_layout_agent_applies_primary_gray_highlight_panel():
    state = create_state("/tmp/paper.pdf", layout_template="cluster_3_portrait", width=36, height=50.88)
    state["color_scheme"] = {"theme": "#0057B8", "contrast": "#7F4B13"}
    agent = LayoutAgent()
    agent.config["selective_block_backgrounds"]["enabled"] = True
    container = {"section_id": "stable_framework", "priority": 0.1}
    section = {"section_id": "stable_framework", "content_role": "method"}

    agent._apply_selective_highlight_panel(container, section, state, {"stable_framework": 0})

    assert container["highlight_panel"] is True
    assert container["fill_color"] == "#EEF0F2"
    assert "border_color" not in container
    assert "border_style" not in container
    assert container["priority"] <= 0.08


def test_layout_agent_frames_normal_and_support_blocks():
    state = create_state("/tmp/paper.pdf", layout_template="cluster_3_portrait", width=36, height=50.88)
    agent = LayoutAgent()
    agent.config["selective_block_backgrounds"]["enabled"] = True
    agent.config["selective_block_backgrounds"]["frame_all_blocks"] = True
    normal = {
        "section_id": "method_details",
        "template_prior": True,
        "priority": 0.1,
    }
    support = {
        "section_id": "robustness_checks",
        "template_prior": True,
        "priority": 0.1,
    }

    agent._apply_selective_block_frame_style(
        normal,
        {"section_id": "method_details", "section_title": "Method Details", "content_role": "method"},
        state,
        {},
    )
    agent._apply_selective_block_frame_style(
        support,
        {"section_id": "robustness_checks", "section_title": "Robustness", "content_role": "overview"},
        state,
        {},
    )

    assert "fill_color" not in normal
    assert normal["border_color"] == "#D2D6DC"
    assert normal["border_style"] == "solid"
    assert support["border_color"] == "#D8DDE3"
    assert support["border_style"] == "dashed"


def test_layout_agent_uses_translucent_panels_for_generated_background():
    state = create_state("/tmp/paper.pdf", enable_generated_background=True)
    agent = LayoutAgent()
    container = {
        "section_id": "method_details",
        "template_prior": True,
        "priority": 0.1,
    }

    agent._apply_visual_block_panel_style(container, state)

    assert container["background_aware_panel"] is True
    assert container["fill_color"] == "#FFFFFF"
    assert container["fill_opacity"] == pytest.approx(0.64)
    assert container["border_color"] == "#D8DEE7"
    assert container["border_opacity"] == pytest.approx(0.44)
    assert container["shadow"]["alpha"] == pytest.approx(0.03)


def test_layout_agent_keeps_solid_panels_without_generated_background():
    state = create_state("/tmp/paper.pdf")
    agent = LayoutAgent()
    container = {
        "section_id": "method_details",
        "template_prior": True,
        "priority": 0.1,
    }

    agent._apply_visual_block_panel_style(container, state)

    assert container["fill_color"] == "#F1F2F4"
    assert "fill_opacity" not in container
    assert "background_aware_panel" not in container
    assert container["shadow"]["alpha"] == pytest.approx(0.16)


def test_create_state_uses_draft_stage_when_post_render_pass_is_enabled():
    assert create_state("/tmp/paper.pdf")["render_stage"] == "final"
    assert create_state("/tmp/paper.pdf", enable_generated_teaser=True)["render_stage"] == "final"
    assert create_state("/tmp/paper.pdf", enable_generated_background=True)["render_stage"] == "draft"
    assert create_state("/tmp/paper.pdf", enable_vlm_layout_review=True)["render_stage"] == "draft"
    assert create_state("/tmp/paper.pdf", enable_visual_legibility_review=True)["render_stage"] == "draft"
    assert create_state("/tmp/paper.pdf", enable_block_vlm_review=True)["render_stage"] == "draft"


def test_create_state_supports_isolated_output_root(tmp_path, monkeypatch):
    monkeypatch.setenv("PAPER2POSTER_OUTPUT_ROOT", str(tmp_path / "isolated"))

    state = create_state("/tmp/example_paper/paper.pdf")

    assert state["output_dir"] == str(tmp_path / "isolated" / "example_paper")


def test_generated_teaser_agent_is_noop_when_disabled(tmp_path, monkeypatch):
    def fail_generate_image(self, prompt, width, height, output_path):
        raise AssertionError("generated teaser API should not be called when disabled")

    monkeypatch.setattr("src.agents.generated_teaser_agent.ImageTools.generate_image", fail_generate_image)
    state = create_state(str(tmp_path / "paper.pdf"))
    state["output_dir"] = str(tmp_path / "output")
    state["story_board"] = {
        "spatial_content_plan": {
            "sections": [
                {
                    "section_id": "motivation",
                    "section_title": "Why Search",
                    "content_role": "foundation",
                    "text_content": ["Keep the original motivation text."],
                    "visual_assets": [],
                }
            ]
        },
    }
    state["visual_assets"] = {}

    result = GeneratedTeaserAgent()(state)

    section = result["story_board"]["spatial_content_plan"]["sections"][0]
    assert result.get("generated_teaser_report") is None
    assert section["visual_assets"] == []
    assert section["text_content"] == ["Keep the original motivation text."]


def test_generated_teaser_agent_injects_motivation_visual(tmp_path, monkeypatch):
    def fake_generate_image(self, prompt, width, height, output_path):
        Image.new("RGB", (width, height), color=(230, 240, 250)).save(output_path)
        assert "conceptual teaser/motivation figure" in prompt
        assert "no readable text" in prompt
        assert "no question marks" in prompt
        assert "Primary accent color: #0057B8" in prompt
        assert "Paper title: Active Search for Outreach" in prompt
        assert "poster slot aperture" in prompt
        assert width > height
        return output_path

    monkeypatch.setattr("src.agents.generated_teaser_agent.ImageTools.generate_image", fake_generate_image)
    state = create_state(str(tmp_path / "paper.pdf"), width=54, height=27, enable_generated_teaser=True)
    state["output_dir"] = str(tmp_path / "output")
    state["resolved_layout_template"] = "cluster_104_landscape"
    state["color_scheme"] = {"theme": "#0057B8"}
    state["narrative_content"] = {
        "meta": {
            "poster_title": "Active Search for Outreach",
        }
    }
    state["story_board"] = {
        "spatial_content_plan": {
            "sections": [
                {
                    "section_id": "motivation",
                    "section_title": "Why Search",
                    "content_role": "foundation",
                    "preferred_slot_id": "slot_1",
                    "column_assignment": "slot_1",
                    "text_content": ["Outreach teams must choose where to visit under limited budgets."],
                    "visual_assets": [],
                },
                {
                    "section_id": "method",
                    "section_title": "Method",
                    "content_role": "method",
                    "column_assignment": "slot_2",
                    "text_content": ["A model selects the next parcel."],
                    "visual_assets": [],
                },
            ]
        },
    }
    state["visual_assets"] = {}

    result = GeneratedTeaserAgent()(state)
    section = result["story_board"]["spatial_content_plan"]["sections"][0]
    visual_id = section["visual_assets"][0]["visual_id"]

    assert visual_id == "generated_teaser_1"
    assert result["visual_assets"][visual_id]["provenance"] == "generated_teaser"
    assert result["visual_assets"][visual_id]["aspect"] > 2.5
    assert Path(result["visual_assets"][visual_id]["source_path"]).exists()
    assert result["generated_teaser_report"]["target_section_id"] == "motivation"
    assert result["generated_teaser_report"]["geometry"]["source"] == "template_slot"
    assert result["generated_teaser_report"]["geometry"]["slot_id"] == "slot_1"
    geometry = result["generated_teaser_report"]["geometry"]
    assert geometry["target_height_inches"] == pytest.approx(geometry["slot_height_inches"] * 0.65, abs=0.01)
    assert section["generated_teaser_summary"] is True
    assert len(section["text_content"]) <= 2
    saved_visual_assets = json.loads((Path(state["output_dir"]) / "content" / "visual_assets.json").read_text(encoding="utf-8"))
    saved_story_board = json.loads((Path(state["output_dir"]) / "content" / "story_board.json").read_text(encoding="utf-8"))
    assert saved_visual_assets[visual_id]["provenance"] == "generated_teaser"
    assert saved_story_board["spatial_content_plan"]["sections"][0]["visual_assets"][0]["visual_id"] == visual_id


def test_generated_teaser_rejects_readable_text_with_ocr_guard(tmp_path, monkeypatch):
    calls = []

    def fake_generate_image(self, prompt, width, height, output_path):
        calls.append(prompt)
        image = Image.new("RGB", (width, height), color=(238, 246, 252))
        ImageDraw.Draw(image).line((0, 0, width, height), fill=(40, 90, 160), width=4)
        image.save(output_path)
        return output_path

    monkeypatch.setattr("src.agents.generated_teaser_agent.ImageTools.generate_image", fake_generate_image)
    monkeypatch.setattr(
        "src.agents.generated_teaser_agent.detect_readable_text",
        lambda *args, **kwargs: {
            "available": True,
            "rejected": True,
            "tokens": [{"text": "Architecture", "confidence": 88.0}],
            "reason": "readable_text_detected",
        },
    )
    state = create_state(str(tmp_path / "paper.pdf"), width=54, height=27, enable_generated_teaser=True)
    state["output_dir"] = str(tmp_path / "output")
    state["resolved_layout_template"] = "cluster_104_landscape"
    state["story_board"] = {
        "spatial_content_plan": {
            "sections": [
                {
                    "section_id": "motivation",
                    "section_title": "Motivation",
                    "content_role": "foundation",
                    "preferred_slot_id": "slot_1",
                    "column_assignment": "slot_1",
                    "text_content": ["The paper addresses a concrete scientific bottleneck."],
                    "visual_assets": [],
                }
            ]
        }
    }
    state["visual_assets"] = {}

    result = GeneratedTeaserAgent()(state)
    report = result["generated_teaser_report"]

    assert len(calls) == 3
    assert report["applied"] is False
    assert report["needs_regeneration"] is True
    assert report["used_procedural_fallback"] is False
    assert report["fallback_reason"] == "readable_text_artifacts"
    assert report["safety"]["readable_text_rejected"] is True
    assert result["story_board"]["spatial_content_plan"]["sections"][0]["visual_assets"] == []
    assert not Path(report["teaser_path"] or "missing").exists()


def test_generated_teaser_regenerates_after_content_rejection_without_fallback(tmp_path, monkeypatch):
    calls = []
    ocr_calls = []

    def fake_generate_image(self, prompt, width, height, output_path):
        calls.append(prompt)
        image = Image.new("RGB", (width, height), color=(238, 246, 252))
        ImageDraw.Draw(image).line((0, 0, width, height), fill=(40, 90, 160), width=4)
        image.save(output_path)
        return output_path

    def fake_detect_readable_text(*args, **kwargs):
        ocr_calls.append(args[0])
        rejected = len(ocr_calls) == 1
        return {
            "available": True,
            "rejected": rejected,
            "tokens": [{"text": "Architecture", "confidence": 88.0}] if rejected else [],
            "reason": "readable_text_detected" if rejected else "no_readable_text",
        }

    monkeypatch.setattr("src.agents.generated_teaser_agent.ImageTools.generate_image", fake_generate_image)
    monkeypatch.setattr("src.agents.generated_teaser_agent.detect_readable_text", fake_detect_readable_text)
    state = create_state(str(tmp_path / "paper.pdf"), width=54, height=27, enable_generated_teaser=True)
    state["output_dir"] = str(tmp_path / "output")
    state["resolved_layout_template"] = "cluster_104_landscape"
    state["story_board"] = {
        "spatial_content_plan": {
            "sections": [
                {
                    "section_id": "motivation",
                    "section_title": "Motivation",
                    "content_role": "foundation",
                    "preferred_slot_id": "slot_1",
                    "column_assignment": "slot_1",
                    "text_content": ["The paper addresses a concrete scientific bottleneck."],
                    "visual_assets": [],
                }
            ]
        }
    }
    state["visual_assets"] = {}

    result = GeneratedTeaserAgent()(state)
    report = result["generated_teaser_report"]

    assert len(calls) == 2
    assert "REGENERATION ATTEMPT 2" in calls[1]
    assert report["applied"] is True
    assert report["asset_source"] == "image_api"
    assert report["used_procedural_fallback"] is False
    assert report["needs_regeneration"] is False
    assert report["generation_attempt_count"] == 2
    assert Path(report["teaser_path"]).exists()


def test_final_artifact_gate_records_teaser_needing_regeneration(tmp_path):
    pptx_path = tmp_path / "poster.pptx"
    png_path = tmp_path / "poster.png"
    pptx_path.write_bytes(b"pptx")
    Image.new("RGB", (20, 10), "white").save(png_path)
    state = create_state(str(tmp_path / "paper.pdf"), enable_generated_teaser=True)
    state["pptx_output_path"] = str(pptx_path)
    state["poster_preview_path"] = str(png_path)
    state["generated_teaser_report"] = {
        "applied": False,
        "needs_regeneration": True,
        "fallback_reason": "readable_text_artifacts",
    }

    failures = _final_artifact_failures(state)

    assert failures == [
        {
            "category": "generated_asset",
            "asset": "teaser",
            "reason": "readable_text_artifacts",
            "needs_regeneration": True,
        }
    ]


def test_generated_teaser_uses_portrait_height_policy_for_tall_templates():
    state = create_state(
        "/tmp/paper.pdf",
        width=27,
        height=54,
        layout_template="cluster_22_portrait",
        enable_generated_teaser=True,
    )
    state["resolved_layout_template"] = "cluster_22_portrait"
    section = {
        "section_id": "motivation",
        "section_title": "Why Search",
        "content_role": "foundation",
        "preferred_slot_id": "slot_1",
        "column_assignment": "slot_1",
    }

    geometry = GeneratedTeaserAgent()._resolve_teaser_geometry(state, section)
    section["text_content"] = [
        "Outreach teams must decide where to visit with limited canvassing budget.",
        "The search policy balances exploration with exploitation as labels arrive online.",
    ]
    summary = GeneratedTeaserAgent()._compress_target_section_text(section, geometry)
    max_fraction = load_config()["generated_teaser"]["portrait_max_block_height_fraction"]
    layout_fraction = LayoutAgent()._max_visual_height_fraction("generated_teaser_1", state)

    assert geometry["orientation"] == "portrait"
    assert geometry["target_height_inches"] <= geometry["slot_height_inches"] * max_fraction + 1e-6
    assert geometry["target_height_inches"] < 8.6
    assert layout_fraction == pytest.approx(load_config()["generated_teaser"]["portrait_layout_max_height_fraction"])
    assert len(summary) == 1


def test_generated_teaser_summary_uses_complete_sentences_within_total_budget():
    first = (
        "Face video restoration can exploit strong audio-visual correlation, especially lip-motion synchronization, "
        "but prior audio-aided methods mainly target compression artifacts."
    )
    second = (
        "The paper introduces GAVN, a general audio-assisted face video restoration network for compression artifact "
        "removal, deblurring, and super-resolution."
    )
    section = {"text_content": [first, second]}

    summary = GeneratedTeaserAgent()._compress_target_section_text(
        section,
        {"orientation": "landscape"},
    )

    assert summary == [first]
    assert "audio-aided." not in summary[0]
    assert "compression artifact." not in summary[0]


def test_generated_teaser_agent_skips_sections_with_existing_visuals_by_default(tmp_path, monkeypatch):
    def fake_generate_image(self, prompt, width, height, output_path):
        Image.new("RGB", (width, height), color=(230, 240, 250)).save(output_path)
        return output_path

    monkeypatch.setattr("src.agents.generated_teaser_agent.ImageTools.generate_image", fake_generate_image)
    state = create_state(str(tmp_path / "paper.pdf"), enable_generated_teaser=True)
    state["output_dir"] = str(tmp_path / "output")
    state["story_board"] = {
        "spatial_content_plan": {
            "sections": [
                {
                    "section_id": "motivation",
                    "section_title": "Why Search",
                    "content_role": "foundation",
                    "column_assignment": "slot_1",
                    "text_content": ["Motivation"],
                    "visual_assets": [{"visual_id": "figure_1"}],
                },
                {
                    "section_id": "method",
                    "section_title": "Method",
                    "content_role": "method",
                    "column_assignment": "slot_2",
                    "text_content": ["Method"],
                    "visual_assets": [],
                },
            ]
        },
    }
    state["visual_assets"] = {"figure_1": {"asset_type": "figure"}}

    result = GeneratedTeaserAgent()(state)

    assert result["generated_teaser_report"]["applied"] is False
    assert result["story_board"]["spatial_content_plan"]["sections"][0]["visual_assets"][0]["visual_id"] == "figure_1"


def test_generated_teaser_summary_blocks_are_protected_from_expansion(tmp_path):
    state = create_state(str(tmp_path / "paper.pdf"), enable_block_vlm_review=True)
    state["output_dir"] = str(tmp_path / "output")
    state["story_board"] = {
        "spatial_content_plan": {
            "sections": [
                {
                    "section_id": "slot_1_problem",
                    "section_title": "Search Problem",
                    "generated_teaser_summary": True,
                    "visual_assets": [{"visual_id": "generated_teaser_1"}],
                    "text_content": ["Short motivation summary."],
                }
            ]
        }
    }
    state["block_occupancy_report"] = {
        "settings": {"acceptable_min": 0.97, "hard_max": 0.995},
        "blocks": [
            {
                "slot_id": "slot_1",
                "section_id": "slot_1_problem",
                "utilization": 0.72,
                "target_extra_chars": 240,
                "action": "expand",
                "visual_count": 1,
            }
        ],
    }
    state["block_vlm_review"] = {"blocks": []}

    report = BlockContentRefiner().refine(state)

    assert report["applied"] is False
    assert report["actions_considered"] == []
    assert state["story_board"]["spatial_content_plan"]["sections"][0]["text_content"] == ["Short motivation summary."]


def test_generated_teaser_summary_blocks_allow_concise_fill_repair(tmp_path, monkeypatch):
    state = create_state(str(tmp_path / "paper.pdf"), enable_block_vlm_review=True)
    state["output_dir"] = str(tmp_path / "output")
    state["story_board"] = {
        "spatial_content_plan": {
            "sections": [
                {
                    "section_id": "slot_1_problem",
                    "section_title": "Search Problem",
                    "generated_teaser_summary": True,
                    "visual_assets": [{"visual_id": "generated_teaser_1"}],
                    "text_content": ["Short motivation summary."],
                }
            ]
        }
    }
    state["block_occupancy_report"] = {
        "settings": {"acceptable_min": 0.96, "hard_max": 0.995},
        "blocks": [
            {
                "slot_id": "slot_1",
                "section_id": "slot_1_problem",
                "available_height": 12.91,
                "visible_content_height": 11.9153,
                "used_height": 11.9153,
                "bottom_whitespace": 0.9947,
                "line_height": 0.7667,
                "chars_per_line": 58,
                "utilization": 0.923,
                "target_extra_chars": 47,
                "action": "expand",
                "visual_count": 1,
            }
        ],
    }
    state["block_vlm_review"] = {
        "blocks": [
            {
                "slot_id": "slot_1",
                "section_id": "slot_1_problem",
                "status": "ok",
                "severity": "low",
                "description": "caption fits comfortably",
            }
        ]
    }
    state["block_refinement_count"] = 1

    def fake_expansion(self, state, actions, section_by_id):
        assert actions[0]["target_extra_chars"] == 47
        return {"slot_1_problem": {"new_bullets": ["Adds budget and cost context for the teaser."]}}

    monkeypatch.setattr(BlockContentRefiner, "_generate_expansion_patches", fake_expansion)

    report = BlockContentRefiner().refine(state)

    assert report["applied"] is True
    assert state["story_board"]["spatial_content_plan"]["sections"][0]["text_content"][-1] == "Adds budget and cost context for the teaser."


def test_background_image_agent_prompt_is_background_only():
    state = create_state("/tmp/paper.pdf")
    state["color_scheme"] = {"theme": "#0057B8", "mono_light": "#E6EAEF"}
    state["background_palette"] = "light_blue"
    state["story_board"] = {
        "spatial_content_plan": {
            "sections": [
                {
                    "section_id": "main_results",
                    "section_title": "Main Results",
                    "content_role": "results",
                }
            ]
        }
    }
    prompt = BackgroundImageAgent()._build_prompt(state)

    assert "BACKGROUND ONLY" in prompt
    assert "no text" in prompt
    assert "not plain white" in prompt
    assert "Selected background palette: light_blue" in prompt
    assert "#0057B8" in prompt
    assert "Main Results" not in prompt
    assert "Paper title context" not in prompt


def test_background_image_agent_switches_palette_prompt():
    state = create_state("/tmp/paper.pdf", background_palette="light_gray")

    prompt = BackgroundImageAgent()._build_prompt(state)

    assert "pale neutral gray" in prompt
    assert "Selected background palette: light_gray" in prompt


def test_background_image_agent_explicit_style_prompt_preserves_safety():
    state = create_state("/tmp/paper.pdf", background_style="tech_grid", background_palette="light_blue")
    state["color_scheme"] = {"theme": "#0057B8", "mono_light": "#E6EAEF"}

    agent = BackgroundImageAgent()
    decision = agent._background_style_decision(state)
    prompt = agent._build_prompt(state, decision, "light_blue")

    assert decision["requested_style"] == "tech_grid"
    assert decision["resolved_style"] == "tech_grid"
    assert "Resolved background style: tech_grid" in prompt
    assert "thin technical grid lines" in prompt
    assert "BACKGROUND ONLY" in prompt
    assert "no text" in prompt
    assert "Do not create block frames" in prompt


def test_background_image_agent_auto_selects_cartographic_for_geospatial_paper():
    state = create_state("/tmp/paper.pdf", background_style="auto", background_palette="auto")
    state["narrative_content"] = {
        "title": "Active Geospatial Search for Efficient Tenant Eviction Outreach",
    }
    state["visual_density"] = "rich"
    state["story_board"] = {
        "spatial_content_plan": {
            "sections": [
                {
                    "section_id": "motivation",
                    "section_title": "Urban Parcel Search",
                    "content_role": "foundation",
                    "text_content": ["Outreach teams map rental properties and select city regions under limited search budgets."],
                }
            ]
        }
    }

    agent = BackgroundImageAgent()
    decision = agent._background_style_decision(state)
    palette = agent._palette_name(state, decision)
    prompt = agent._build_prompt(state, decision, palette)

    assert decision["requested_style"] == "auto"
    assert decision["resolved_style"] == "cartographic"
    assert palette == "mint"
    assert "cartographic contour lines" in prompt
    assert "Selected background palette: mint" in prompt


def test_background_image_agent_auto_palette_uses_style_default():
    state = create_state("/tmp/paper.pdf", background_style="flat_cartoon", background_palette="auto")
    agent = BackgroundImageAgent()
    decision = agent._background_style_decision(state)

    assert agent._palette_name(state, decision) == "warm_ivory"


def test_background_image_agent_enforces_visibility_floor_for_pale_background():
    agent = BackgroundImageAgent()
    img = Image.new("RGB", (120, 80), (252, 253, 254))
    for x in range(0, 120, 8):
        for y in range(80):
            img.putpixel((x, y), (244, 249, 253))

    before = agent._background_visibility_metrics(img)
    boosted = agent._enforce_background_visibility(
        img,
        {
            "min_average_distance_from_white": 12.0,
            "min_channel_stddev": 4.0,
            "max_boost_factor": 5.0,
        },
    )
    after = agent._background_visibility_metrics(boosted)

    assert after["average_distance_from_white"] > before["average_distance_from_white"] * 2
    assert after["channel_stddev"] > before["channel_stddev"] * 2
    assert after["average_distance_from_white"] >= 10.0


def test_background_image_agent_placeholder_is_marked_for_regeneration_without_fallback(tmp_path, monkeypatch):
    calls = []

    def fake_generate_image(self, prompt, width, height, output_path):
        calls.append(prompt)
        Image.new("RGB", (width, height), color=(200, 200, 200)).save(output_path)
        return output_path

    monkeypatch.setattr("src.agents.background_image_agent.ImageTools.generate_image", fake_generate_image)
    state = create_state(str(tmp_path / "paper.pdf"), enable_generated_background=True, background_palette="light_blue")
    state["output_dir"] = str(tmp_path / "output")
    state["color_scheme"] = {"theme": "#0057B8", "mono_light": "#E6EAEF"}
    agent = BackgroundImageAgent()
    agent.background_config["width_px"] = 160
    agent.background_config["height_px"] = 120

    result = agent(state)

    assert result["errors"] == []
    assert len(calls) == 3
    assert result["background_image_path"] is None
    assert result["background_image_report"]["palette"] == "light_blue"
    assert result["background_image_report"]["resolved_palette"] == "light_blue"
    assert result["background_image_report"]["requested_style"] == "auto"
    assert result["background_image_report"]["resolved_style"]
    assert result["background_image_report"]["used_procedural_fallback"] is False
    assert result["background_image_report"]["applied"] is False
    assert result["background_image_report"]["needs_regeneration"] is True
    assert result["background_image_report"]["asset_source"] == "none"


def test_background_image_agent_uses_poster_preview_as_reference(tmp_path, monkeypatch):
    def fake_edit_image(self, image_path, prompt, output_path):
        Image.new("RGB", (160, 120), color=(242, 248, 255)).save(output_path)
        assert image_path.endswith("draft.png")
        assert "provided poster image" in prompt
        return output_path

    monkeypatch.setattr("src.agents.background_image_agent.ImageTools.edit_image", fake_edit_image)
    state = create_state(str(tmp_path / "paper.pdf"), enable_generated_background=True, background_palette="light_blue")
    state["output_dir"] = str(tmp_path / "output")
    state["color_scheme"] = {"theme": "#0057B8", "mono_light": "#E6EAEF"}
    preview_path = tmp_path / "draft.png"
    Image.new("RGB", (160, 120), color=(255, 255, 255)).save(preview_path)
    state["poster_preview_path"] = str(preview_path)
    agent = BackgroundImageAgent()
    agent.background_config["width_px"] = 160
    agent.background_config["height_px"] = 120
    agent.background_config["procedural_only"] = False
    agent.background_config["condition_on_poster"] = True  # opt in to poster conditioning

    result = agent(state)

    assert result["errors"] == []
    assert Path(result["background_image_path"]).exists()
    assert result["background_image_report"]["generation_mode"] == "poster_conditioned_image_api"
    assert result["background_image_report"]["reference_poster_path"] == str(preview_path)


def test_background_image_agent_defaults_to_text_to_image_without_poster_reference(tmp_path, monkeypatch):
    """By default the background must be generated text-to-image, never conditioned on
    the rendered poster, so the image model cannot copy poster text into the background
    (the cause of ghosted/duplicated text)."""

    def fail_edit_image(self, image_path, prompt, output_path):
        raise AssertionError("edit_image must not be called when conditioning is off")

    def fake_generate_image(self, prompt, width, height, output_path):
        assert "No reference image is provided" in prompt
        Image.new("RGB", (width or 160, height or 120), color=(244, 249, 255)).save(output_path)
        return output_path

    monkeypatch.setattr("src.agents.background_image_agent.ImageTools.edit_image", fail_edit_image)
    monkeypatch.setattr("src.agents.background_image_agent.ImageTools.generate_image", fake_generate_image)
    state = create_state(str(tmp_path / "paper.pdf"), enable_generated_background=True, background_palette="light_blue")
    state["output_dir"] = str(tmp_path / "output")
    state["color_scheme"] = {"theme": "#0057B8", "mono_light": "#E6EAEF"}
    preview_path = tmp_path / "draft.png"
    Image.new("RGB", (160, 120), color=(255, 255, 255)).save(preview_path)
    state["poster_preview_path"] = str(preview_path)
    agent = BackgroundImageAgent()
    agent.background_config["width_px"] = 160
    agent.background_config["height_px"] = 120
    agent.background_config["procedural_only"] = False

    result = agent(state)

    assert result["errors"] == []
    assert Path(result["background_image_path"]).exists()
    assert result["background_image_report"]["generation_mode"] == "image_api"
    assert result["background_image_report"]["used_procedural_fallback"] is False


def test_background_image_agent_rejects_background_that_copies_layout_text(tmp_path, monkeypatch):
    from PIL import ImageDraw

    def fake_edit_image(self, image_path, prompt, output_path):
        img = Image.new("RGB", (320, 160), color=(246, 251, 255))
        draw = ImageDraw.Draw(img)
        draw.text((78, 8), "Copied Poster Title", fill=(55, 55, 55))
        for x, y, w, h, label in [
            (20, 42, 110, 12, "Motivation"),
            (150, 42, 130, 12, "Core Method"),
            (20, 102, 110, 12, "Results"),
        ]:
            draw.rectangle([x, y, x + w, y + h], fill=(218, 225, 232))
            draw.text((x + 28, y + 1), label, fill=(64, 64, 64))
        img.save(output_path)
        return output_path

    monkeypatch.setattr("src.agents.background_image_agent.ImageTools.edit_image", fake_edit_image)
    state = create_state(
        str(tmp_path / "paper.pdf"),
        width=16,
        height=8,
        enable_generated_background=True,
        background_palette="light_blue",
    )
    state["output_dir"] = str(tmp_path / "output")
    state["color_scheme"] = {"theme": "#0057B8", "mono_light": "#E6EAEF"}
    state["styled_layout"] = [
        {"type": "title", "x": 1.0, "y": 0.3, "width": 14.0, "height": 1.0},
        {"type": "section_title", "section_id": "motivation", "x": 1.0, "y": 2.1, "width": 5.5, "height": 0.6},
        {"type": "section_title", "section_id": "method", "x": 7.5, "y": 2.1, "width": 6.5, "height": 0.6},
        {"type": "section_title", "section_id": "results", "x": 1.0, "y": 5.1, "width": 5.5, "height": 0.6},
    ]
    preview_path = tmp_path / "draft.png"
    Image.new("RGB", (320, 160), color=(255, 255, 255)).save(preview_path)
    state["poster_preview_path"] = str(preview_path)
    agent = BackgroundImageAgent()
    agent.background_config["width_px"] = 320
    agent.background_config["height_px"] = 160
    agent.background_config["procedural_only"] = False
    agent.background_config["condition_on_poster"] = True  # contamination check only runs on the conditioned path

    result = agent(state)
    report = result["background_image_report"]

    assert result["errors"] == []
    assert report["used_procedural_fallback"] is False
    assert report["needs_regeneration"] is True
    assert report["applied"] is False
    assert report["generation_mode"] == "poster_conditioned_image_api_rejected_no_fallback"
    assert report["postprocess"]["fallback_reason"] == "layout_copy_artifacts"
    assert report["postprocess"]["copy_artifact_report"]["rejected"] is True
    assert report["safety"]["layout_copy_artifacts_rejected"] is True
    assert result["degraded_quality_states"][-1]["category"] == "generated_background"
    assert Path(report["raw_path"]).exists()
    assert report["background_image_path"] == ""


def test_background_image_agent_rejects_ocr_text_artifacts(tmp_path, monkeypatch):
    def fake_generate_image(self, prompt, width, height, output_path):
        image = Image.new("RGB", (width, height), color=(238, 246, 252))
        ImageDraw.Draw(image).line((0, 0, width, height), fill=(80, 120, 180), width=3)
        image.save(output_path)
        return output_path

    monkeypatch.setattr("src.agents.background_image_agent.ImageTools.generate_image", fake_generate_image)
    monkeypatch.setattr(
        "src.agents.background_image_agent.detect_readable_text",
        lambda *args, **kwargs: {
            "available": True,
            "rejected": True,
            "tokens": [{"text": "Motivation", "confidence": 91.0}],
            "reason": "readable_text_detected",
        },
    )
    state = create_state(str(tmp_path / "paper.pdf"), enable_generated_background=True)
    state["output_dir"] = str(tmp_path / "output")
    agent = BackgroundImageAgent()
    agent.background_config["width_px"] = 320
    agent.background_config["height_px"] = 160

    report = agent(state)["background_image_report"]

    assert report["used_procedural_fallback"] is False
    assert report["needs_regeneration"] is True
    assert report["applied"] is False
    assert report["postprocess"]["fallback_reason"] == "readable_text_artifacts"
    assert report["safety"]["readable_text_rejected"] is True


def test_background_image_agent_regenerates_after_content_rejection(tmp_path, monkeypatch):
    calls = []
    ocr_calls = []

    def fake_generate_image(self, prompt, width, height, output_path):
        calls.append(prompt)
        image = Image.new("RGB", (width, height), color=(238, 246, 252))
        ImageDraw.Draw(image).line((0, 0, width, height), fill=(80, 120, 180), width=3)
        image.save(output_path)
        return output_path

    def fake_detect_readable_text(*args, **kwargs):
        ocr_calls.append(args[0])
        rejected = len(ocr_calls) == 1
        return {
            "available": True,
            "rejected": rejected,
            "tokens": [{"text": "Motivation", "confidence": 91.0}] if rejected else [],
            "reason": "readable_text_detected" if rejected else "no_readable_text",
        }

    monkeypatch.setattr("src.agents.background_image_agent.ImageTools.generate_image", fake_generate_image)
    monkeypatch.setattr("src.agents.background_image_agent.detect_readable_text", fake_detect_readable_text)
    state = create_state(str(tmp_path / "paper.pdf"), enable_generated_background=True)
    state["output_dir"] = str(tmp_path / "output")
    agent = BackgroundImageAgent()
    agent.background_config["width_px"] = 320
    agent.background_config["height_px"] = 160

    result = agent(state)
    report = result["background_image_report"]

    assert len(calls) == 2
    assert "REGENERATION ATTEMPT 2" in calls[1]
    assert report["asset_source"] == "image_api"
    assert report["applied"] is True
    assert report["needs_regeneration"] is False
    assert report["generation_attempt_count"] == 2
    assert Path(report["background_image_path"]).exists()


def test_color_agent_records_degraded_state_when_visual_color_extraction_fails(tmp_path, monkeypatch):
    visual_path = tmp_path / "figure.png"
    Image.new("RGB", (120, 80), color=(240, 240, 255)).save(visual_path)
    state = create_state(str(tmp_path / "paper.pdf"))
    state["output_dir"] = str(tmp_path / "output")
    state["classified_visuals"] = {"key_visual": "figure_1"}
    state["visual_assets"] = {"figure_1": {"source_path": str(visual_path)}}

    monkeypatch.setattr(
        ColorAgent,
        "_analyze_figure_for_color",
        lambda self, image_path, state: (_ for _ in ()).throw(RuntimeError("vision unavailable")),
    )

    result = ColorAgent()(state)

    assert result["errors"] == []
    assert result["color_scheme"]["theme"] == load_config()["colors"]["fallback_theme"]
    assert result["degraded_quality_states"][-1]["component"] == "color_agent"
    assert result["degraded_quality_states"][-1]["category"] == "color_extraction"
    assert result["degraded_quality_states"][-1]["fallback"] == "default_theme"


def test_background_image_agent_api_failure_is_recorded_without_fallback(tmp_path, monkeypatch):
    def fail_edit_image(self, image_path, prompt, output_path):
        raise TimeoutError("stuck image API")

    monkeypatch.setattr("src.agents.background_image_agent.ImageTools.edit_image", fail_edit_image)
    state = create_state(str(tmp_path / "paper.pdf"), enable_generated_background=True, background_palette="light_blue")
    state["output_dir"] = str(tmp_path / "output")
    state["color_scheme"] = {"theme": "#0057B8", "mono_light": "#E6EAEF"}
    preview_path = tmp_path / "draft.png"
    Image.new("RGB", (160, 120), color=(255, 255, 255)).save(preview_path)
    state["poster_preview_path"] = str(preview_path)
    agent = BackgroundImageAgent()
    agent.background_config["width_px"] = 160
    agent.background_config["height_px"] = 120
    agent.background_config["procedural_only"] = False
    agent.background_config["api_timeout_seconds"] = 1
    agent.background_config["condition_on_poster"] = True  # exercise the poster-conditioned failure path

    result = agent(state)

    assert result["errors"] == []
    assert result["background_image_path"] is None
    assert result["background_image_report"]["used_procedural_fallback"] is False
    assert result["background_image_report"]["needs_regeneration"] is True
    assert result["background_image_report"]["generation_mode"] == "poster_conditioned_image_api_rejected_no_fallback"
    assert result["background_image_report"]["raw_path"] == ""


def test_background_image_agent_timeout_can_use_environment_override(monkeypatch):
    monkeypatch.setenv("BACKGROUND_IMAGE_API_TIMEOUT_SECONDS", "2.5")
    agent = BackgroundImageAgent()
    agent.background_config["api_timeout_seconds"] = 75

    assert agent._api_timeout_seconds() == 2.5


def test_background_image_agent_procedural_only_skips_image_api(tmp_path, monkeypatch):
    def fail_generate_image(self, prompt, width, height, output_path):
        raise AssertionError("image generation API should not be called")

    monkeypatch.setattr("src.agents.background_image_agent.ImageTools.generate_image", fail_generate_image)
    state = create_state(
        str(tmp_path / "paper.pdf"),
        enable_generated_background=True,
        background_palette="auto",
        background_style="minimal_solid",
    )
    state["output_dir"] = str(tmp_path / "output")
    agent = BackgroundImageAgent()
    agent.background_config["procedural_only"] = True
    agent.background_config["width_px"] = 160
    agent.background_config["height_px"] = 120

    result = agent(state)

    assert result["errors"] == []
    assert Path(result["background_image_path"]).exists()
    assert result["background_image_report"]["generation_mode"] == "procedural_only"
    assert result["background_image_report"]["resolved_style"] == "minimal_solid"
    assert result["background_image_report"]["palette"] == "light_gray"
    assert result["background_image_report"]["raw_path"] == ""


def test_background_image_agent_matches_landscape_poster_aspect():
    state = create_state("/tmp/paper.pdf", width=54, height=27, enable_generated_background=True)
    agent = BackgroundImageAgent()
    agent.background_config["width_px"] = 160
    agent.background_config["height_px"] = 120

    assert agent._background_dimensions(state) == (160, 80)


def test_image_tools_failover_retries_each_base_url(tmp_path, monkeypatch):
    import base64
    import io

    img = Image.new("RGB", (1, 1), color=(240, 248, 255))
    buffer = io.BytesIO()
    img.save(buffer, format="PNG")
    b64 = base64.b64encode(buffer.getvalue()).decode("ascii")
    calls = []

    class FakeResponse:
        def raise_for_status(self):
            return None

        def json(self):
            return {"data": [{"b64_json": b64}]}

    def fake_post(url, **kwargs):
        calls.append(url)
        if url.startswith("https://first.example"):
            raise RuntimeError("temporary upstream failure")
        return FakeResponse()

    monkeypatch.setattr("src.tools.image_api.requests.post", fake_post)
    tool = ImageTools(
        api_key="test-key",
        base_url="https://first.example/v1, https://second.example/v1",
        model="gpt-image-2",
        retry_attempts=2,
        retry_delay=0,
    )
    output_path = tmp_path / "generated.png"

    result = tool.generate_image("plain academic background", width=1600, height=900, output_path=str(output_path))

    assert result == str(output_path)
    assert output_path.exists()
    assert calls == [
        "https://first.example/v1/images/generations",
        "https://first.example/v1/images/generations",
        "https://second.example/v1/images/generations",
    ]


def test_image_tools_retries_transient_errors_five_times_with_six_second_intervals(monkeypatch):
    tool = ImageTools(
        api_key="test-key",
        base_url="https://example.test/v1",
        model="test-image-model",
        fallback_models=[],
        retry_attempts=5,
        retry_delay=6,
    )
    calls = []
    sleeps = []

    monkeypatch.setattr("src.tools.image_api.time.sleep", sleeps.append)

    def operation(base_url):
        calls.append(base_url)
        if len(calls) < 5:
            raise TimeoutError("temporary image generation timeout")
        return "ok"

    assert tool._request_with_failover("image retry test", operation) == "ok"
    assert calls == ["https://example.test/v1"] * 5
    assert sleeps == [6, 6, 6, 6]


def test_image_tools_does_not_retry_hard_balance_errors(monkeypatch):
    tool = ImageTools(
        api_key="test-key",
        base_url="https://first.example/v1, https://second.example/v1",
        model="test-image-model",
        fallback_models=[],
        retry_attempts=5,
        retry_delay=6,
    )
    calls = []
    sleeps = []

    monkeypatch.setattr("src.tools.image_api.time.sleep", sleeps.append)

    def operation(base_url):
        calls.append(base_url)
        raise RuntimeError("Your account balance is insufficient; please recharge before continuing")

    with pytest.raises(RuntimeError, match="account balance is insufficient"):
        tool._request_with_failover("image quota test", operation)

    assert calls == ["https://first.example/v1"]
    assert sleeps == []


def test_image_tools_default_timeout_allows_slow_generation(monkeypatch):
    monkeypatch.delenv("IMAGE_REQUEST_TIMEOUT_SECONDS", raising=False)

    tool = ImageTools(
        api_key="test-key",
        base_url="https://example.test/v1",
        model="test-image-model",
        fallback_models=[],
    )

    assert tool.request_timeout == 120


def test_image_tools_gpt_image_uses_supported_aspect_size():
    tool = ImageTools(
        api_key="test-key",
        base_url="https://example.test/v1",
        model="gpt-image-2",
        retry_attempts=1,
        retry_delay=0,
    )

    assert tool._request_size(2035, 1018) == "1536x1024"
    assert tool._request_size(1018, 2035) == "1024x1536"
    assert tool._request_size(1024, 1024) == "1024x1024"


def test_image_tools_base_url_list_takes_priority_over_legacy_env(monkeypatch):
    monkeypatch.setenv("IMAGE_BASE_URLS", "https://first.example/v1, https://second.example/v1")
    monkeypatch.setenv("IMAGE_BASE_URL", "https://legacy-image.example/v1")
    monkeypatch.setenv("VLM_BASE_URL", "https://legacy-vlm.example/v1")

    tool = ImageTools(api_key="test-key", model="gpt-image-2", retry_attempts=1, retry_delay=0)

    assert tool.base_urls == ["https://first.example/v1", "https://second.example/v1"]


def test_image_tools_falls_back_from_gpt_image_to_gemini_without_retrying_model_errors():
    tool = ImageTools(
        api_key="test-key",
        base_url="https://example.test/v1",
        model="gpt-image-2",
        retry_attempts=5,
        retry_delay=0,
    )
    calls = []

    def operation(base_url):
        calls.append((tool.model, base_url))
        if tool.model == "gpt-image-2":
            raise RuntimeError("No available channel for model gpt-image-2 under group auto")
        return "ok"

    result = tool._request_with_failover("image model fallback test", operation)

    assert result == "ok"
    assert tool.models[:2] == ["gpt-image-2", "gemini-3.1-flash-image-preview"]
    assert calls == [
        ("gpt-image-2", "https://example.test/v1"),
        ("gemini-3.1-flash-image-preview", "https://example.test/v1"),
    ]


def test_image_tools_does_not_retry_unsupported_operations():
    tool = ImageTools(
        api_key="test-key",
        base_url="https://first.example/v1, https://second.example/v1",
        model="unsupported-image-model",
        retry_attempts=5,
        retry_delay=0,
    )
    calls = []

    def operation(base_url):
        calls.append((tool.model, base_url))
        raise RuntimeError("The requested operation is unsupported.")

    with pytest.raises(RuntimeError):
        tool._request_with_failover("unsupported image operation", operation)

    assert calls == [
        ("unsupported-image-model", "https://first.example/v1"),
        ("unsupported-image-model", "https://second.example/v1"),
    ]


def test_template_block_planner_matches_block_count_to_content_slots(monkeypatch):
    json_response = """{
      "blocks": [
        {"target_title": "Overview", "target_bullets": ["Problem framing.", "Core idea."]},
        {"target_title": "Method", "target_bullets": ["Framework details.", "Optimization flow."]},
        {"target_title": "Results", "target_bullets": ["Main result.", "Comparison summary."]},
        {"target_title": "Takeaways", "target_bullets": ["Deployment note.", "Conclusion."]}
      ]
    }"""

    class FakeResponse:
        input_tokens = 1
        output_tokens = 1
        content = json_response

    class FakeAgent:
        def __init__(self, *args, **kwargs):
            pass

        def step(self, message):
            return FakeResponse()

    monkeypatch.setattr("src.agents.template_block_planner.LangGraphAgent", FakeAgent)

    state = create_state("/tmp/paper.pdf", layout_template="cluster_2_landscape", width=54, height=27)
    state["resolved_layout_template"] = "cluster_2_landscape"
    state["story_board"] = {
        "spatial_content_plan": {
            "sections": [
                {
                    "section_id": "overview",
                    "section_title": "Overview",
                    "column_assignment": "left",
                    "vertical_priority": "top",
                    "text_content": ["Problem framing.", "Core idea."],
                    "visual_assets": [],
                },
                {
                    "section_id": "method",
                    "section_title": "Method",
                    "column_assignment": "middle",
                    "vertical_priority": "top",
                    "text_content": ["Framework details.", "Optimization flow."],
                    "visual_assets": [{"visual_id": "figure_1"}],
                },
                {
                    "section_id": "results",
                    "section_title": "Results",
                    "column_assignment": "right",
                    "vertical_priority": "middle",
                    "text_content": ["Main result.", "Comparison summary."],
                    "visual_assets": [],
                },
            ]
        }
    }

    result = TemplateBlockPlanner()(state)

    blocks = result["template_block_plan"]["blocks"]
    assert len(blocks) == 3
    assert len({block["slot_id"] for block in blocks}) == 3
    rewritten_sections = result["story_board"]["spatial_content_plan"]["sections"]
    assert len(rewritten_sections) == 3
    assert all(section["column_assignment"].startswith("slot_") for section in rewritten_sections)
    assert all(section["slot_id"] == section["column_assignment"] for section in rewritten_sections)
    assert all(section.get("capacity_budget") for section in rewritten_sections)


def test_block_capacity_contract_scales_with_slot_area():
    planner = TemplateBlockPlanner()
    state = create_state("/tmp/paper.pdf", width=36, height=51, layout_template="cluster_3_portrait")
    layout = load_block_template_layout("cluster_3_portrait", 36, 51, margin=1.0)
    regions = sorted(layout["regions"], key=lambda item: float(item["area_ratio"]))
    small_region = regions[0]
    large_region = regions[-1]
    sections = [
        {
            "section_id": "large_method",
            "section_title": "Large Method",
            "content_role": "method",
            "region_id": large_region["region_id"],
            "slot_id": large_region["region_id"],
            "visual_assets": [],
        },
        {
            "section_id": "small_note",
            "section_title": "Small Note",
            "content_role": "takeaway",
            "region_id": small_region["region_id"],
            "slot_id": small_region["region_id"],
            "visual_assets": [],
        },
    ]

    contract = planner._build_block_capacity_contract(sections, layout, state)
    by_section = contract["by_section"]

    assert by_section["large_method"]["target_chars"] > by_section["small_note"]["target_chars"]
    assert by_section["large_method"]["target_bullets"] >= by_section["small_note"]["target_bullets"]


def test_block_capacity_contract_reserves_visual_space():
    planner = TemplateBlockPlanner()
    state = create_state("/tmp/paper.pdf", width=36, height=51, layout_template="cluster_3_portrait")
    state["visual_assets"] = {
        "figure_1": {"asset_id": "figure_1", "asset_type": "figure", "aspect": 1.4}
    }
    layout = load_block_template_layout("cluster_3_portrait", 36, 51, margin=1.0)
    region = max(layout["regions"], key=lambda item: float(item["area_ratio"]))
    settings = planner._capacity_settings()
    text_only = {
        "section_id": "text_only",
        "section_title": "Text Only",
        "content_role": "overview",
        "visual_assets": [],
    }
    with_visual = {
        "section_id": "with_visual",
        "section_title": "With Visual",
        "content_role": "method",
        "visual_assets": [{"visual_id": "figure_1"}],
    }

    text_budget = planner._capacity_budget_for_section(text_only, region, state, settings)
    visual_budget = planner._capacity_budget_for_section(with_visual, region, state, settings)

    assert visual_budget["reserved_visual_height"] > text_budget["reserved_visual_height"]
    assert visual_budget["available_text_height"] < text_budget["available_text_height"]
    assert visual_budget["raw_target_chars"] < text_budget["raw_target_chars"]
    assert visual_budget["visual_policy"] in {"reserve_visual_space", "prioritize_visual_scale"}


def test_template_block_planner_capacity_rewrite_preserves_refs_and_expands(tmp_path, monkeypatch):
    class FailingAgent:
        def __init__(self, *args, **kwargs):
            pass

        def step(self, message):
            raise RuntimeError("offline")

    monkeypatch.setattr("src.agents.template_block_planner.LangGraphAgent", FailingAgent)
    state = create_state(str(tmp_path / "paper.pdf"), width=36, height=51, layout_template="cluster_3_portrait")
    state["output_dir"] = str(tmp_path / "output")
    state["resolved_layout_template"] = "cluster_3_portrait"
    state["raw_text"] = (
        "The method uses active geospatial search to choose a sequence of rental units. "
        "The search policy updates after each query and uses property-level information. "
        "The approach accounts for budget and travel-cost constraints. "
        "Experiments compare HAGS with greedy and conventional active search baselines. "
        "Results show improved targeting under uniform and travel-aware budgets."
    )
    state["story_board"] = {
        "spatial_content_plan": {
            "sections": [
                {
                    "section_id": "method",
                    "section_title": "Active Search Method",
                    "column_assignment": "middle",
                    "vertical_priority": "top",
                    "text_content": ["Active search selects parcels."],
                    "visual_assets": [],
                    "source_sections": ["method"],
                },
                {
                    "section_id": "results",
                    "section_title": "Results",
                    "column_assignment": "right",
                    "vertical_priority": "middle",
                    "text_content": ["HAGS improves targeting."],
                    "visual_assets": [],
                    "source_sections": ["results"],
                },
                {
                    "section_id": "problem",
                    "section_title": "Problem",
                    "column_assignment": "left",
                    "vertical_priority": "middle",
                    "text_content": ["Outreach teams have limited budgets."],
                    "visual_assets": [],
                    "source_sections": ["problem"],
                },
            ]
        }
    }

    result = TemplateBlockPlanner()(state)
    sections = result["story_board"]["spatial_content_plan"]["sections"]
    method = next(section for section in sections if section["section_id"] == "method")

    assert result["block_capacity_contract"]["blocks"]
    assert result["capacity_planning_report"]["blocks"]
    assert method["slot_id"] == method["column_assignment"]
    assert method["source_sections"] == ["method"]
    assert "capacity_budget" in method
    assert sum(len(item) for item in method["text_content"]) >= len("Active search selects parcels.")


def test_template_block_planner_keypoint_mode_uses_available_unique_slots(tmp_path, monkeypatch):
    class FailingAgent:
        def __init__(self, *args, **kwargs):
            pass

        def step(self, message):
            raise RuntimeError("offline")

    monkeypatch.setattr("src.agents.template_block_planner.LangGraphAgent", FailingAgent)
    template_id = "cluster_62_landscape"
    info = get_block_template_info(template_id)
    canvas = info["recommended_canvas_size"]
    state = create_state(
        str(tmp_path / "paper.pdf"),
        width=canvas["width"],
        height=canvas["height"],
        layout_template=template_id,
    )
    state["output_dir"] = str(tmp_path / "output")
    state["resolved_layout_template"] = template_id
    state["paper_poster_keypoints"] = [
        {"id": index, "key_point": f"Keypoint factual claim {index}", "section": "Method" if index < 7 else "Experiments"}
        for index in range(1, 11)
    ]
    state["poster_reading_order"] = list(range(1, 11))
    state["raw_text"] = " ".join(
        f"Keypoint factual claim {index} is supported by the paper with additional implementation and evaluation context."
        for index in range(1, 11)
    )
    state["story_board"] = {
        "spatial_content_plan": {
            "sections": [
                {
                    "section_id": f"keypoint_{index}",
                    "section_title": f"Point {index}",
                    "column_assignment": "middle",
                    "vertical_priority": "middle",
                    "text_content": [f"Keypoint factual claim {index}."],
                    "visual_assets": [],
                    "keypoint_id": index,
                    "source_section": "Method" if index < 7 else "Experiments",
                    "source_sections": ["Method" if index < 7 else "Experiments"],
                }
                for index in range(1, 11)
            ]
        }
    }

    result = TemplateBlockPlanner()(state)

    blocks = result["template_block_plan"]["blocks"]
    available_slots = len(result["layout_template_metadata"]["regions"])
    assert len(blocks) == available_slots
    assert len({block["slot_id"] for block in blocks}) == available_slots
    assert [block["keypoint_id"] for block in blocks] == list(range(1, available_slots + 1))
    rewritten_sections = result["story_board"]["spatial_content_plan"]["sections"]
    assert [section["keypoint_id"] for section in rewritten_sections] == list(range(1, available_slots + 1))
    assert all(section.get("capacity_budget") for section in rewritten_sections)


def test_template_block_planner_grouped_keypoints_preserve_visuals(tmp_path, monkeypatch):
    class FailingAgent:
        def __init__(self, *args, **kwargs):
            pass

        def step(self, message):
            raise RuntimeError("offline")

    monkeypatch.setattr("src.agents.template_block_planner.LangGraphAgent", FailingAgent)
    template_id = "cluster_104_landscape"
    info = get_block_template_info(template_id)
    canvas = info["recommended_canvas_size"]
    state = create_state(
        str(tmp_path / "paper.pdf"),
        width=canvas["width"],
        height=canvas["height"],
        layout_template=template_id,
    )
    state["output_dir"] = str(tmp_path / "output")
    state["resolved_layout_template"] = template_id
    state["paper_poster_keypoints"] = [
        {"id": index, "key_point": f"Poster keypoint {index}", "section": "Method" if index < 6 else "Results"}
        for index in range(1, 11)
    ]
    state["poster_reading_order"] = list(range(1, 11))
    state["raw_text"] = (
        "am-ELO reformulates Elo scoring with maximum likelihood estimation. "
        "The method models annotator ability and estimates model scores jointly. "
        "Results show lower loss and better prediction performance than baselines. "
        "Robustness tests evaluate perturbations to arena outcomes."
    )
    state["visual_assets"] = {
        "figure_2": {"asset_type": "figure", "aspect": 1.7, "source_path": "/tmp/figure-2.png"},
        "figure_3": {"asset_type": "figure", "aspect": 2.3, "source_path": "/tmp/figure-3.png"},
        "table_3": {"asset_type": "table", "aspect": 0.94, "source_path": "/tmp/table-3.png"},
    }
    state["story_board"] = {
        "spatial_content_plan": {
            "sections": [
                {
                    "section_id": "motivation",
                    "section_title": "Motivation",
                    "column_assignment": "left",
                    "vertical_priority": "top",
                    "text_content": ["Arena comparisons can create unstable Elo rankings."],
                    "visual_assets": [],
                    "keypoint_id": 1,
                    "source_keypoint_ids": [1, 2],
                    "source_sections": ["Introduction"],
                    "content_type": "foundation",
                },
                {
                    "section_id": "method",
                    "section_title": "am-ELO Method",
                    "column_assignment": "middle",
                    "vertical_priority": "top",
                    "text_content": ["am-ELO models annotator ability in pairwise comparisons."],
                    "visual_assets": [{"visual_id": "figure_2"}],
                    "keypoint_id": 3,
                    "source_keypoint_ids": [3, 4],
                    "source_sections": ["Method"],
                    "content_type": "method",
                },
                {
                    "section_id": "estimation",
                    "section_title": "MLE Estimate",
                    "column_assignment": "middle",
                    "vertical_priority": "middle",
                    "text_content": ["Likelihood-based estimation avoids sequential update sensitivity."],
                    "visual_assets": [],
                    "keypoint_id": 5,
                    "source_keypoint_ids": [5],
                    "source_sections": ["Method"],
                    "content_type": "method",
                },
                {
                    "section_id": "setup",
                    "section_title": "Experiment Setup",
                    "column_assignment": "left",
                    "vertical_priority": "middle",
                    "text_content": ["Experiments compare Elo, m-ELO, and am-ELO on arena records."],
                    "visual_assets": [],
                    "keypoint_id": 6,
                    "source_keypoint_ids": [6],
                    "source_sections": ["Experiments"],
                    "content_type": "foundation",
                },
                {
                    "section_id": "results",
                    "section_title": "Key Results",
                    "column_assignment": "right",
                    "vertical_priority": "top",
                    "text_content": ["am-ELO obtains lower loss than baseline estimators."],
                    "visual_assets": [{"visual_id": "figure_3"}],
                    "keypoint_id": 7,
                    "source_keypoint_ids": [7],
                    "source_sections": ["Results"],
                    "content_type": "results",
                },
                {
                    "section_id": "table",
                    "section_title": "Prediction",
                    "column_assignment": "right",
                    "vertical_priority": "middle",
                    "text_content": ["Prediction experiments indicate better generalization."],
                    "visual_assets": [{"visual_id": "table_3"}],
                    "keypoint_id": 8,
                    "source_keypoint_ids": [8],
                    "source_sections": ["Results"],
                    "content_type": "results",
                },
                {
                    "section_id": "takeaway",
                    "section_title": "Takeaway",
                    "column_assignment": "right",
                    "vertical_priority": "bottom",
                    "text_content": ["Annotator-aware MLE improves stability for arena-based LLM evaluation."],
                    "visual_assets": [],
                    "keypoint_id": 9,
                    "source_keypoint_ids": [9, 10],
                    "source_sections": ["Robustness", "Conclusion"],
                    "content_type": "takeaway",
                },
            ]
        }
    }

    result = TemplateBlockPlanner()(state)

    sections = result["story_board"]["spatial_content_plan"]["sections"]
    visual_ids = [
        visual["visual_id"]
        for section in sections
        for visual in section.get("visual_assets", [])
    ]
    region_by_id = {
        region["region_id"]: region
        for region in result["layout_template_metadata"]["regions"]
    }
    visual_sections = [section for section in sections if section.get("visual_assets")]
    assert len(sections) == 7
    assert len({section["slot_id"] for section in sections}) == 7
    assert {"figure_2", "figure_3", "table_3"}.issubset(set(visual_ids))
    assert all(section.get("source_keypoint_ids") for section in sections)
    assert all(region_by_id[section["slot_id"]]["can_host_visual"] for section in visual_sections)
    assert all(region_by_id[section["slot_id"]]["text_density_limit"] != "low" for section in visual_sections)


def test_layout_templates_support_block_template_ids():
    template_names = LayoutTemplates.available_template_names()

    assert "cluster_2_landscape" in template_names
    assert "cluster_104_landscape" in template_names
    assert "cluster_3_portrait" in template_names
    info = get_block_template_info("cluster_3_portrait")
    assert info["orientation"] == "portrait"
    layout = LayoutTemplates(36, 51, margin=1.0, col_gap=1.0).get_template("cluster_3_portrait")
    assert layout["layout_mode"] == "template_prior"
    assert layout["orientation"] == "portrait"
    assert [lane["id"] for lane in layout["lanes"]][:2] == ["slot_1", "slot_2"]
    assert "base_layout_template" not in layout


def test_dense_landscape_template_exposes_seven_visual_capable_content_blocks():
    layout = load_block_template_layout("cluster_104_landscape", 54, 27, margin=1.0)
    regions = layout["regions"]

    assert layout["slot_count"] == 7
    assert len(regions) == 7
    assert sum(1 for region in regions if region.get("can_host_visual")) >= 5
    assert any(region["region_id"] == "slot_7" and region.get("text_density_limit") == "low" for region in regions)


def test_template_block_planner_keeps_main_results_out_of_small_low_density_slot():
    layout = load_block_template_layout("cluster_104_landscape", 54, 27, margin=1.0)
    planner = TemplateBlockPlanner()
    region = planner._region_for_ordered_section(
        {
            "section_id": "sec_main_results",
            "section_title": "Main Results",
            "content_role": "results",
            "visual_assets": [],
        },
        layout["regions"],
        {"slot_1", "slot_2", "slot_3", "slot_4", "slot_5"},
        later_visual_count=0,
    )

    assert region["region_id"] == "slot_6"
    slot_7 = next(item for item in layout["regions"] if item["region_id"] == "slot_7")
    assert region["area_ratio"] > slot_7["area_ratio"]
    assert region["can_host_visual"] is True


def test_resolve_poster_dimensions_allows_imported_wide_block_templates():
    width, height = resolve_poster_dimensions("cluster_104_landscape", None, None)

    assert width / height <= 2.1


def test_font_agent_keyword_prompt_uses_narrative_content(monkeypatch):
    captured = {}

    class FakeResponse:
        content = '{"section_keywords": {}, "formatting_summary": {}}'
        input_tokens = 1
        output_tokens = 1

    class FakeAgent:
        def __init__(self, *args, **kwargs):
            pass

        def step(self, message):
            captured["message"] = message
            return FakeResponse()

    state = create_state("/tmp/paper.pdf")
    state["narrative_content"] = {"and": "poster narrative signal"}
    agent = FontAgent()
    monkeypatch.setattr("src.agents.font_agent.LangGraphAgent", FakeAgent)

    agent._identify_keywords({"spatial_content_plan": {"sections": []}}, state)

    assert "poster narrative signal" in captured["message"]


def test_font_agent_normalizes_keyword_section_aliases_to_layout_ids():
    agent = FontAgent()
    story_board = {
        "spatial_content_plan": {
            "sections": [
                {
                    "section_id": "sec_method_ags",
                    "section_title": "AGS Formulation",
                    "content_role": "method",
                    "text_content": ["Define AGS as a budget-constrained MDP."],
                },
                {
                    "section_id": "sec_motivation_foundation",
                    "section_title": "Eviction Outreach Challenge",
                    "content_role": "foundation",
                    "text_content": ["Motivate tenant outreach."],
                }
            ]
        }
    }
    keywords = {
        "section_keywords": {
            "method_ags": {
                "bold_contrast": ["Active Geospatial Search"],
                "bold": ["MDP"],
                "italic": ["policy"],
            },
            "motivation": {
                "bold_contrast": ["Active Geospatial Search"],
                "bold": ["at-risk"],
                "italic": [],
            }
        }
    }

    normalized = agent._normalize_keyword_section_ids(keywords, story_board)

    assert "sec_method_ags" in normalized["section_keywords"]
    assert normalized["section_keywords"]["sec_method_ags"]["bold"] == ["MDP"]
    assert "sec_motivation_foundation" in normalized["section_keywords"]
    assert normalized["section_keywords"]["sec_motivation_foundation"]["bold"] == ["at-risk"]


def test_micro_layout_refiner_styles_added_fill_lines_with_section_keywords():
    refiner = MicroLayoutRefiner()
    state = {
        "keywords": {
            "section_keywords": {
                "sec_results": {
                    "bold_contrast": ["HAGS"],
                    "bold": ["ANT"],
                    "italic": ["sparse"],
                }
            }
        },
        "color_scheme": {"contrast": "#7A1F2B"},
    }

    styled = refiner._apply_fill_keyword_highlighting(
        "HAGS improves ANT under sparse targets.",
        "sec_results",
        state,
    )

    assert "<color:#7A1F2B>HAGS</color>" in styled
    assert "**ANT**" in styled
    assert "*sparse*" in styled


def test_text_cleanup_repairs_mojibake_bullets_and_common_ocr_typos():
    text = "â¢ **Realistic Setting: When costs matter.\\nâ¦ Effcient search improves 42%â70% with î»L_BCE."

    cleaned = normalize_text_for_poster(text)

    assert "â" not in cleaned
    assert "• **Realistic Setting:** When costs matter." in cleaned
    assert "◦ Efficient search improves 42%-70% with lambda L_BCE." in cleaned


def test_text_cleanup_removes_ocr_paths_tables_and_metadata_noise():
    noisy = (
        "A detailed formal presentation follows. ![](_page_4_Figure_0.jpeg) "
        "Figure 2: HAGS policy network architecture. | Search Budget | 15 | 20 |"
    )

    cleaned = normalize_text_for_poster(noisy)

    assert cleaned == "A detailed formal presentation follows."
    assert "_page_" not in cleaned
    assert ".jpeg" not in cleaned
    assert "|" not in cleaned
    assert "Figure 2" not in cleaned


def test_text_cleanup_removes_table_references_and_section_metadata():
    metadata = (
        "HAGS uses a hierarchical policy importance high contains_figures contains_tables "
        "section_name Related Work section_type foundation content."
    )

    assert normalize_text_for_poster(metadata) == "HAGS uses a hierarchical policy"
    assert normalize_text_for_poster("The results are presented in Tables 3 and 4.") == ""
    assert normalize_text_for_poster("Table 5: Main comparison across baselines.") == ""
    assert normalize_text_for_poster("Challenge:** Outreach teams cannot visit every parcel.") == "**Challenge:** Outreach teams cannot visit every parcel."
    assert normalize_text_for_poster("A detailed presentation of the complete method is provided in Algorithm 2 in Supplement.") == ""
    assert normalize_text_for_poster("HAGS is best across budgets, outperforming.") == "HAGS is best across budgets."


def test_renderer_resolves_figure_and_table_paths():
    renderer = Renderer()
    state = create_state("/tmp/paper.pdf")
    state["resolved_visual_assets"] = {
        "methods_figure_1": {"resolved_path": "/tmp/figure-slot.png"},
        "results_table_1": {"resolved_path": "/tmp/table-slot.png"},
    }

    assert renderer._get_resolved_visual_entry("methods_figure_1", "figure_1", state)["resolved_path"] == "/tmp/figure-slot.png"
    assert renderer._get_resolved_visual_entry("results_table_1", "table_1", state)["resolved_path"] == "/tmp/table-slot.png"


def test_visual_asset_agent_disabled_is_slot_preserving_crop_only(tmp_path):
    source_path = tmp_path / "source.png"
    from PIL import Image

    Image.new("RGB", (100, 80), color=(255, 0, 0)).save(source_path)

    state = create_state(str(tmp_path / "paper.pdf"))
    state["output_dir"] = str(tmp_path / "output")
    state["enable_visual_refinement"] = False
    state["visual_assets"] = {
        "figure_1": {
            "asset_id": "figure_1",
            "asset_type": "figure",
            "source_path": str(source_path),
            "resolved_path": None,
            "caption": "Figure 1",
            "aspect": 1.25,
            "provenance": "paper_extracted",
        }
    }
    state["styled_layout"] = [
        {
            "type": "visual",
            "slot_id": "method_figure_1",
            "id": "method_figure_1",
            "visual_id": "figure_1",
            "width": 2.0,
            "height": 1.0,
        }
    ]

    result = VisualAssetAgent()(state)

    assert result["visual_plan"][0]["action"] == "crop_only"
    assert "method_figure_1" in result["resolved_visual_assets"]
    assert Path(result["resolved_visual_assets"]["method_figure_1"]["resolved_path"]).exists()


def test_visual_asset_agent_table_crop_only_preserves_full_table_edges(tmp_path):
    source_path = tmp_path / "wide_table.png"
    image = Image.new("RGB", (400, 100), color=(255, 255, 255))
    for x in range(0, 40):
        for y in range(0, 100):
            image.putpixel((x, y), (220, 0, 0))
    for x in range(360, 400):
        for y in range(0, 100):
            image.putpixel((x, y), (0, 0, 220))
    image.save(source_path)

    state = create_state(str(tmp_path / "paper.pdf"))
    state["output_dir"] = str(tmp_path / "output_table")
    state["enable_visual_refinement"] = False
    state["visual_assets"] = {
        "table_1": {
            "asset_id": "table_1",
            "asset_type": "table",
            "source_path": str(source_path),
            "resolved_path": None,
            "caption": "Table 1",
            "aspect": 4.0,
            "provenance": "paper_extracted",
        }
    }
    state["styled_layout"] = [
        {
            "type": "visual",
            "slot_id": "results_table_1",
            "id": "results_table_1",
            "visual_id": "table_1",
            "width": 2.0,
            "height": 1.0,
        }
    ]

    result = VisualAssetAgent()(state)

    resolved = result["resolved_visual_assets"]["results_table_1"]
    assert resolved["provenance"] == "fit_resized"
    with Image.open(resolved["resolved_path"]).convert("RGB") as rendered:
        mid_y = rendered.height // 2
        left_pixel = rendered.getpixel((5, mid_y))
        right_pixel = rendered.getpixel((rendered.width - 6, mid_y))
    assert left_pixel[0] > 150 and left_pixel[1] < 80 and left_pixel[2] < 80
    assert right_pixel[2] > 150 and right_pixel[0] < 80 and right_pixel[1] < 80


def test_visual_asset_agent_enabled_generates_for_missing_source_slot(tmp_path, monkeypatch):
    def fake_generate_image(self, prompt, width, height, output_path):
        Image.new("RGB", (width, height), "white").save(output_path)
        return output_path

    monkeypatch.setattr(ImageTools, "generate_image", fake_generate_image)
    state = create_state(str(tmp_path / "paper.pdf"), enable_visual_refinement=True)
    state["output_dir"] = str(tmp_path / "output")
    state["visual_assets"] = {}
    state["styled_layout"] = [
        {
            "type": "visual",
            "slot_id": "method_generated_visual",
            "id": "method_generated_visual",
            "width": 2.0,
            "height": 1.0,
        }
    ]

    result = VisualAssetAgent()(state)

    assert result["visual_plan"][0]["action"] == "generate_new"
    resolved = result["resolved_visual_assets"]["method_generated_visual"]
    assert resolved["asset_type"] == "generated"
    assert Path(resolved["resolved_path"]).exists()


def test_vlm_layout_reviewer_disabled_is_noop():
    state = create_state("/tmp/paper.pdf")
    state["styled_layout"] = [{"type": "title", "id": "title", "x": 1, "y": 1, "width": 10, "height": 2}]

    result = VLMLayoutReviewer()(state)

    assert result.get("vlm_layout_review") is None
    assert result["styled_layout"] == state["styled_layout"]


def test_vlm_layout_reviewer_uses_responses_endpoint(monkeypatch):
    captured = {}

    class FakeResponse:
        headers = {}

        def raise_for_status(self):
            pass

        def json(self):
            return {
                "output_text": '{"overall_score": 90, "accept": true, "issues": [], "patch": [], "visual_asset_recommendations": []}'
            }

    def fake_post(url, headers, json, timeout, stream=False):
        captured["url"] = url
        captured["payload"] = json
        captured["stream"] = stream
        return FakeResponse()

    monkeypatch.setattr("src.agents.vlm_layout_reviewer.requests.post", fake_post)
    reviewer = VLMLayoutReviewer()
    response = reviewer._post_vlm_request(
        "https://example.com/api/v1/responses",
        {"Authorization": "Bearer test"},
        "gpt-5.4",
        "review",
        "data:image/png;base64,abc",
    )
    text = reviewer._extract_response_text(response)

    assert captured["url"] == "https://example.com/api/v1/responses"
    assert captured["payload"]["store"] is False
    assert captured["payload"]["stream"] is True
    assert captured["stream"] is True
    assert captured["payload"]["input"][0]["content"][0]["type"] == "input_text"
    assert captured["payload"]["input"][0]["content"][1]["type"] == "input_image"
    assert "overall_score" in text


def test_vlm_layout_reviewer_records_response_usage():
    class FakeResponse:
        headers = {"content-type": "application/json"}

        def json(self):
            return {
                "output_text": '{"accept": true}',
                "usage": {"input_tokens": 321, "output_tokens": 45},
            }

    reviewer = VLMLayoutReviewer()
    assert reviewer._extract_response_text(FakeResponse()) == '{"accept": true}'

    state = create_state("/tmp/paper.pdf")
    reviewer._record_usage(state, "vlm_layout_reviewer")

    calls = state["timing_metrics"].api_calls
    assert len(calls) == 1
    assert calls[0].agent == "vlm_layout_reviewer"
    assert calls[0].input_tokens == 321
    assert calls[0].output_tokens == 45
    assert state["tokens"].input_vision == 321
    assert state["tokens"].output_vision == 45


def test_vlm_layout_reviewer_falls_back_from_responses_to_chat(monkeypatch):
    transports = []

    class FakeResponse:
        headers = {"content-type": "application/json"}

        def raise_for_status(self):
            pass

        def json(self):
            return {
                "choices": [{"message": {"content": '{"accept": true}'}}],
                "usage": {"prompt_tokens": 12, "completion_tokens": 3},
            }

    def fake_post(self, base_url, headers, model, prompt, image_data, *, transport=None):
        transports.append(transport)
        if transport == "responses":
            raise ValueError("response.failed")
        return FakeResponse()

    reviewer = VLMLayoutReviewer()
    monkeypatch.setattr(VLMLayoutReviewer, "_post_vlm_request", fake_post)

    content = reviewer._request_vlm_text(
        "https://example.com/v1",
        {"Authorization": "Bearer test"},
        "gpt-5.4",
        "review",
        "data:image/png;base64,abc",
    )

    assert content == '{"accept": true}'
    assert transports == ["responses", "chat"]


def test_vlm_layout_reviewer_falls_back_on_request_failure(tmp_path, monkeypatch):
    state = create_state(str(tmp_path / "paper.pdf"), enable_vlm_layout_review=True)
    state["output_dir"] = str(tmp_path / "output")
    state["poster_preview_path"] = str(tmp_path / "preview.png")
    state["styled_layout"] = [{"type": "title", "id": "title", "x": 1, "y": 1, "width": 10, "height": 2}]
    Path(state["output_dir"], "content").mkdir(parents=True)
    Image.new("RGB", (200, 120), color=(255, 255, 255)).save(state["poster_preview_path"])

    monkeypatch.setenv("VLM_BASE_URL", "https://example.com/v1")
    monkeypatch.setenv("VLM_API_KEY", "test")
    monkeypatch.setenv("VLM_MODEL", "gpt-5.4")

    def fail_post(self, base_url, headers, model, prompt, image_data, **kwargs):
        raise ConnectionResetError("connection reset")

    monkeypatch.setattr(VLMLayoutReviewer, "_post_vlm_request", fail_post)

    result = VLMLayoutReviewer()(state)

    assert result["errors"] == []
    assert result["vlm_layout_review"]["source"] == "fallback"
    assert result["vlm_layout_review"]["accept"] is False
    assert "VLM layout request failed" in result["vlm_layout_review"]["warnings"][0]


def test_visual_legibility_reviewer_records_degraded_state_on_vlm_fallback(tmp_path, monkeypatch):
    state = create_state(str(tmp_path / "paper.pdf"), enable_visual_legibility_review=True)
    state["output_dir"] = str(tmp_path / "output")
    state["poster_preview_path"] = str(tmp_path / "preview.png")
    state["styled_layout"] = []
    state["layout_template_metadata"] = {"lanes": []}
    Path(state["output_dir"], "content").mkdir(parents=True)
    Image.new("RGB", (200, 120), color=(255, 255, 255)).save(state["poster_preview_path"])

    monkeypatch.setenv("VLM_BASE_URL", "https://example.com/v1")
    monkeypatch.setenv("VLM_API_KEY", "test")
    monkeypatch.setenv("VLM_MODEL", "gpt-5.4")

    def fail_post(self, base_url, headers, model, prompt, image_data, **kwargs):
        raise ConnectionResetError("connection reset")

    monkeypatch.setattr(VLMLayoutReviewer, "_post_vlm_request", fail_post)

    result = VisualLegibilityReviewer()(state)

    assert result["errors"] == []
    review = result["visual_legibility_review"]
    assert review["source"] == "fallback"
    assert review["degraded"] is True
    assert result["degraded_quality_states"][-1]["component"] == "visual_legibility_reviewer"
    assert result["degraded_quality_states"][-1]["category"] == "visual_legibility_review"
    assert result["degraded_quality_states"][-1]["fallback"] == "deterministic_visual_heuristic"


def test_gpt_54_uses_openai_chat_provider():
    config = _get_model_config("gpt-5.4")

    assert config.provider == "openai"
    assert config.model_name == "gpt-5.4"


def test_langgraph_agent_does_not_retry_authentication_errors():
    class FakeAuthenticationError(Exception):
        status_code = 401

    class FailingModel:
        def __init__(self):
            self.calls = 0

        def invoke(self, _history):
            self.calls += 1
            raise FakeAuthenticationError("401 Unauthorized: 无效的令牌")

    model = FailingModel()
    agent = LangGraphAgent.__new__(LangGraphAgent)
    agent.system_msg = "system"
    agent.config = _get_model_config("gpt-5.4")
    agent.model = model
    agent.history = []
    agent.state = None
    agent.agent_name = "test"

    with pytest.raises(FakeAuthenticationError):
        agent.step("hello")

    assert model.calls == 1


def test_vlm_layout_reviewer_applies_single_safe_patch(tmp_path, monkeypatch):
    state = create_state(str(tmp_path / "paper.pdf"), enable_vlm_layout_review=True)
    state["output_dir"] = str(tmp_path / "output")
    Path(state["output_dir"], "content").mkdir(parents=True)
    preview = tmp_path / "preview.png"
    from PIL import Image

    Image.new("RGB", (200, 120), color=(255, 255, 255)).save(preview)
    state["poster_preview_path"] = str(preview)
    state["layout_template_metadata"] = LayoutTemplates(54, 36, margin=1.0, col_gap=1.0).get_template(
        "three_column_postergen",
        header_height=(36 - 2) * 0.18,
    )
    left_lane = state["layout_template_metadata"]["lanes"][0]
    state["styled_layout"] = [
        {
            "type": "section_container",
            "section_id": "intro",
            "lane_id": "left",
            "x": left_lane["x"],
            "y": left_lane["y"],
            "width": left_lane["w"],
            "height": 4.0,
            "priority": 0.1,
        },
        {
            "type": "visual",
            "id": "intro_visual",
            "slot_id": "intro_visual",
            "visual_id": "figure_1",
            "x": left_lane["x"] + 1.0,
            "y": left_lane["y"] + 1.0,
            "width": 4.0,
            "height": 2.0,
            "priority": 0.4,
        },
    ]

    def fake_review(self, state):
        return {
            "overall_score": 82,
            "accept": False,
            "issues": [{"severity": "medium", "category": "whitespace", "target": "intro_visual"}],
            "patch": [{"target": "intro_visual", "op": "increase_visual_scale", "value": 1.1}],
            "visual_asset_recommendations": [],
        }

    monkeypatch.setattr(VLMLayoutReviewer, "_review_or_fallback", fake_review)
    result = VLMLayoutReviewer()(state)

    visual = next(element for element in result["styled_layout"] if element.get("id") == "intro_visual")
    assert result["vlm_patch_applied"] is True
    assert result["vlm_reflow_required"] is True
    assert result["vlm_review_count"] == 1
    assert visual["width"] > 4.0


def test_vlm_layout_reviewer_fast_mode_records_patch_without_reflow(tmp_path, monkeypatch):
    state = create_state(str(tmp_path / "paper.pdf"), enable_vlm_layout_review=True)
    state["output_dir"] = str(tmp_path / "output")
    state["template_fast_mode"] = True
    state["template_layout_mode"] = "template_prior"
    state["styled_layout"] = [{"type": "title", "id": "title", "x": 1, "y": 1, "width": 10, "height": 2}]
    Path(state["output_dir"], "content").mkdir(parents=True)

    reviewer = VLMLayoutReviewer()
    monkeypatch.setattr(
        reviewer,
        "_review_or_fallback",
        lambda _state: {
            "source": "test",
            "overall_score": 62,
            "accept": False,
            "issues": [{"severity": "medium", "category": "whitespace", "description": "underfilled"}],
            "patch": [{"operation": "move", "target": "title", "dx": 0.1, "dy": 0.0}],
            "warnings": [],
        },
    )

    result = reviewer(state)

    assert result["vlm_reflow_required"] is False
    assert result["vlm_patch_applied"] is False
    assert result["template_repair_required"] is False
    assert result["vlm_layout_patch"]
    assert "did not apply" in result["vlm_layout_review"]["warnings"][0]


def test_visual_legibility_fast_mode_is_report_only(tmp_path, monkeypatch):
    state = create_state(str(tmp_path / "paper.pdf"), enable_visual_legibility_review=True)
    state["output_dir"] = str(tmp_path / "output")
    state["template_fast_mode"] = True
    state["template_layout_mode"] = "template_prior"
    Path(state["output_dir"], "content").mkdir(parents=True)

    reviewer = VisualLegibilityReviewer()
    monkeypatch.setattr(
        reviewer,
        "_review_or_fallback",
        lambda _state: {
            "needs_relayout": True,
            "issues": [{"severity": "medium", "target": "slot_3", "description": "small figure text"}],
            "layout_recommendation": {"target_region": "slot_3", "action": "promote_region", "reason": "small"},
            "warnings": [],
        },
    )
    monkeypatch.setattr(reviewer, "_merge_heuristic_review", lambda _state, review: review)

    result = reviewer(state)

    assert result["template_repair_required"] is False
    assert result["adaptive_relayout_required"] is False
    assert "does not trigger automatic relayout" in result["visual_legibility_review"]["warnings"][0]


def test_layout_templates_expose_multiple_geometry_families():
    templates = LayoutTemplates(54, 36, margin=1.0, col_gap=1.0)

    three_col = templates.get_template("three_column_postergen", header_height=6.0)
    two_plus_one = templates.get_template("two_plus_one_mixed", header_height=6.0)
    one_plus_two = templates.get_template("one_plus_two_mixed", header_height=6.0)
    single_col = templates.get_template("single_column_vertical", header_height=6.0)
    adaptive = templates.get_template(
        "adaptive_three_column",
        header_height=6.0,
        width_ratios={"left": 0.85, "middle": 1.30, "right": 0.85},
    )

    assert len(three_col["lanes"]) == 3
    assert round(three_col["lanes"][0]["w"], 4) == round(three_col["lanes"][1]["w"], 4)

    assert two_plus_one["lanes"][2]["w"] > two_plus_one["lanes"][0]["w"]
    assert one_plus_two["lanes"][0]["w"] > one_plus_two["lanes"][1]["w"]

    assert len({round(lane["x"], 4) for lane in single_col["lanes"]}) == 1
    assert single_col["lanes"][0]["y"] < single_col["lanes"][1]["y"] < single_col["lanes"][2]["y"]

    assert adaptive["template_name"] == "adaptive_three_column"
    assert adaptive["lanes"][1]["w"] > adaptive["lanes"][0]["w"]
    assert adaptive["lanes"][1]["w"] > adaptive["lanes"][2]["w"]


def test_template_extractor_builds_expected_template_schema():
    image_path = Path("template/poster(1).png")
    if not image_path.exists():
        return

    template, raw = build_template(image_path)

    assert template["geometry_policy"] == "soft"
    assert template["source_lanes"]
    assert template["panel_style_tokens"]
    assert len(template["lanes"]) == 3
    assert raw["ocr_blocks"]
    for box in [template["header"], *template["lanes"], *template["logo_regions"]]:
        assert 0 <= box["x"] <= 1
        assert 0 <= box["y"] <= 1
        assert 0 < box["w"] <= 1
        assert 0 < box["h"] <= 1


def test_extracted_template_registry_and_scaling():
    template_ids = set(list_extracted_template_ids())
    if not template_ids:
        return

    assert "extracted_poster1_landscape_three_panel" in template_ids
    assert "extracted_poster2_landscape_multi_panel" in template_ids
    assert "extracted_poster3_portrait_section_band" in template_ids

    loaded = load_extracted_template("extracted_poster3_portrait_section_band")
    assert loaded is not None
    assert loaded["orientation"] == "portrait"

    layout = LayoutTemplates(36, 54, margin=1.0, col_gap=1.0).get_template(
        "extracted_poster3_portrait_section_band"
    )
    assert layout["template_name"] == "extracted_poster3_portrait_section_band"
    assert layout["orientation"] == "portrait"
    assert layout["geometry_policy"] == "soft"
    assert len(layout["lanes"]) == 3
    assert layout["lanes"][0]["w"] > layout["lanes"][0]["h"]
    assert layout["lanes"][0]["y"] < layout["lanes"][1]["y"] < layout["lanes"][2]["y"]
    assert layout["source_lanes"][0]["h"] != layout["lanes"][0]["h"]

    horizontal = LayoutTemplates(54, 36, margin=1.0, col_gap=1.0).get_template(
        "extracted_poster1_landscape_three_panel",
        header_height=6.0,
    )
    body_height = 36 - 1.0 - horizontal["lanes"][0]["y"]
    assert len(horizontal["lanes"]) == 3
    assert horizontal["lanes"][0]["h"] >= body_height * 0.95
    assert horizontal["lanes"][0]["y"] < horizontal["source_lanes"][0]["y"]


def test_layout_agent_resolves_extracted_template_metadata():
    if "extracted_poster3_portrait_section_band" not in set(list_extracted_template_ids()):
        return

    state = create_state(
        "/tmp/paper.pdf",
        width=36,
        height=54,
        layout_template="extracted_poster3_portrait_section_band",
    )

    template = LayoutAgent()._resolve_template_layout(state)

    assert template["template_name"] == "extracted_poster3_portrait_section_band"
    assert state["resolved_layout_template"] == "extracted_poster3_portrait_section_band"
    assert template["header"]["w"] == 36 * 0.95
    assert template["lanes"][0]["w"] == 34


def test_layout_agent_respects_requested_template_geometry():
    state = create_state("/tmp/paper.pdf", layout_template="one_plus_two_mixed")
    state["narrative_content"] = {"meta": {"poster_title": "Paper", "authors": "Authors"}}
    state["story_board"] = {
        "spatial_content_plan": {
            "sections": [
                {
                    "section_id": "intro",
                    "section_title": "Intro",
                    "column_assignment": "left",
                    "vertical_priority": "top",
                    "text_content": ["Point A", "Point B"],
                    "visual_assets": [],
                    "importance_level": 2,
                },
                {
                    "section_id": "method",
                    "section_title": "Method",
                    "column_assignment": "middle",
                    "vertical_priority": "top",
                    "text_content": ["Point A", "Point B"],
                    "visual_assets": [],
                    "importance_level": 1,
                },
                {
                    "section_id": "results",
                    "section_title": "Results",
                    "column_assignment": "right",
                    "vertical_priority": "top",
                    "text_content": ["Point A", "Point B"],
                    "visual_assets": [],
                    "importance_level": 2,
                },
            ]
        }
    }

    result = LayoutAgent()(state, mode="initial")
    sections = {
        element["section_id"]: element
        for element in result["initial_layout_data"]
        if element.get("type") == "section_container"
    }

    assert sections["intro"]["width"] > sections["method"]["width"]
    assert sections["intro"]["width"] > sections["results"]["width"]


def _six_slot_story_board(*, long_slot_4: bool = False):
    sections = []
    for idx, slot_id in enumerate(["slot_1", "slot_2", "slot_3", "slot_4", "slot_5", "slot_6"], start=1):
        text = [f"Concise section {idx} point."]
        if long_slot_4 and slot_id == "slot_4":
            text = [
                "This deliberately long paragraph repeats enough content to exceed the fixed "
                "template slot before micro-layout has a chance to compress it. " * 8
            ]
        sections.append(
            {
                "section_id": f"section_{idx}",
                "section_title": f"Section {idx}",
                "column_assignment": slot_id,
                "preferred_slot_id": slot_id,
                "vertical_priority": "top",
                "text_content": text,
                "visual_assets": [],
                "importance_level": 2,
            }
        )
    return {"spatial_content_plan": {"sections": sections}}


def test_template_prior_layout_validation_does_not_poison_global_errors():
    state = create_state("/tmp/paper.pdf", layout_template="cluster_43_landscape", width=54, height=27)
    state["narrative_content"] = {"meta": {"poster_title": "Paper", "authors": "Authors"}}
    state["optimized_story_board"] = _six_slot_story_board(long_slot_4=True)

    result = LayoutAgent()(state, mode="final")

    assert result["layout_validation"]["valid"] is False
    assert result["design_layout"]
    assert not result["errors"]


def test_layout_with_balancer_ignores_stale_errors_when_retrying_template_prior_layout():
    state = create_state("/tmp/paper.pdf", layout_template="cluster_43_landscape", width=54, height=27)
    state["narrative_content"] = {"meta": {"poster_title": "Paper", "authors": "Authors"}}
    state["template_layout_mode"] = "template_prior"
    state["story_board"] = _six_slot_story_board()
    state["errors"] = ["layout_agent: final layout validation failed: stale pre-micro diagnostic"]

    result = LayoutWithBalancerAgent()(state)

    assert result["design_layout"]
    assert result["optimized_story_board"]
    assert result["errors"] == ["layout_agent: final layout validation failed: stale pre-micro diagnostic"]


def test_layout_agent_respects_adaptive_lane_widths():
    state = create_state("/tmp/paper.pdf", layout_template="three_column_postergen")
    state["adaptive_lane_widths"] = {"left": 0.85, "middle": 1.30, "right": 0.85}
    state["narrative_content"] = {"meta": {"poster_title": "Paper", "authors": "Authors"}}
    state["story_board"] = {
        "spatial_content_plan": {
            "sections": [
                {
                    "section_id": "intro",
                    "section_title": "Intro",
                    "column_assignment": "left",
                    "vertical_priority": "top",
                    "text_content": ["Point A"],
                    "visual_assets": [],
                    "importance_level": 2,
                },
                {
                    "section_id": "method",
                    "section_title": "Method",
                    "column_assignment": "middle",
                    "vertical_priority": "top",
                    "text_content": ["Point A"],
                    "visual_assets": [],
                    "importance_level": 1,
                },
                {
                    "section_id": "results",
                    "section_title": "Results",
                    "column_assignment": "right",
                    "vertical_priority": "top",
                    "text_content": ["Point A"],
                    "visual_assets": [],
                    "importance_level": 2,
                },
            ]
        }
    }

    result = LayoutAgent()(state, mode="initial")
    sections = {
        element["section_id"]: element
        for element in result["initial_layout_data"]
        if element.get("type") == "section_container"
    }

    assert result["resolved_layout_template"] == "adaptive_three_column"
    assert sections["method"]["width"] > sections["intro"]["width"]
    assert sections["method"]["width"] > sections["results"]["width"]


def test_visual_legibility_heuristic_requests_middle_lane_for_wide_visual():
    state = create_state("/tmp/paper.pdf", enable_visual_legibility_review=True, enable_adaptive_column_width=True)
    state["layout_template_metadata"] = LayoutTemplates(54, 36, margin=1.0, col_gap=1.0).get_template(
        "three_column_postergen",
        header_height=6.0,
    )
    middle_lane = state["layout_template_metadata"]["lanes"][1]
    state["visual_assets"] = {
        "figure_2": {
            "asset_id": "figure_2",
            "caption": "Hierarchical method pipeline overview",
            "asset_type": "figure",
        }
    }
    state["styled_layout"] = [
        {
            "type": "visual",
            "id": "method_figure",
            "slot_id": "method_figure",
            "visual_id": "figure_2",
            "x": middle_lane["x"] + 0.3,
            "y": middle_lane["y"] + 1.0,
            "width": 15.0,
            "height": 5.0,
        }
    ]

    review = VisualLegibilityReviewer()._heuristic_review(state)

    assert review["needs_relayout"] is True
    assert review["layout_recommendation"]["target_lane"] == "middle"


def test_visual_legibility_heuristic_flags_low_resolution_table(tmp_path):
    source = tmp_path / "tiny_table.png"
    Image.new("RGB", (360, 120), "white").save(source)
    state = create_state(str(tmp_path / "paper.pdf"), enable_visual_legibility_review=True)
    state["template_layout_mode"] = "template_prior"
    state["layout_template_metadata"] = {"lanes": []}
    state["visual_assets"] = {
        "table_1": {
            "asset_id": "table_1",
            "asset_type": "table",
            "source_path": str(source),
            "caption": "Primary quantitative comparison",
        }
    }
    state["styled_layout"] = [
        {
            "type": "visual",
            "id": "results_table",
            "slot_id": "results_table",
            "visual_id": "table_1",
            "width": 12.0,
            "height": 4.0,
        }
    ]

    review = VisualLegibilityReviewer()._heuristic_review(state)

    assert review["needs_relayout"] is True
    assert any(issue["severity"] == "high" for issue in review["issues"])


def test_adaptive_column_relayout_sets_template_and_saves_decision(tmp_path):
    state = create_state(
        str(tmp_path / "paper.pdf"),
        enable_visual_legibility_review=True,
        enable_adaptive_column_width=True,
    )
    state["output_dir"] = str(tmp_path / "output")
    state["adaptive_relayout_required"] = True
    state["visual_legibility_review"] = {
        "needs_relayout": True,
        "layout_recommendation": {
            "target_lane": "middle",
            "action": "widen_lane",
            "preferred_width_ratio": 1.3,
            "reason": "Middle visual text is too small.",
        },
        "issues": [],
    }

    result = AdaptiveColumnRelayoutAgent()(state)

    assert result["layout_template"] == "adaptive_three_column"
    assert result["adaptive_relayout_count"] == 1
    assert result["adaptive_lane_widths"]["middle"] > result["adaptive_lane_widths"]["left"]
    assert Path(state["output_dir"], "content", "adaptive_layout_decision.json").exists()


def test_layout_agent_single_column_template_stacks_semantic_lanes():
    state = create_state("/tmp/paper.pdf", layout_template="single_column_vertical")
    state["narrative_content"] = {"meta": {"poster_title": "Paper", "authors": "Authors"}}
    state["story_board"] = {
        "spatial_content_plan": {
            "sections": [
                {
                    "section_id": "intro",
                    "section_title": "Intro",
                    "column_assignment": "left",
                    "vertical_priority": "top",
                    "text_content": ["Point A", "Point B"],
                    "visual_assets": [],
                    "importance_level": 2,
                },
                {
                    "section_id": "method",
                    "section_title": "Method",
                    "column_assignment": "middle",
                    "vertical_priority": "top",
                    "text_content": ["Point A", "Point B"],
                    "visual_assets": [],
                    "importance_level": 1,
                },
                {
                    "section_id": "results",
                    "section_title": "Results",
                    "column_assignment": "right",
                    "vertical_priority": "top",
                    "text_content": ["Point A", "Point B"],
                    "visual_assets": [],
                    "importance_level": 2,
                },
            ]
        }
    }

    result = LayoutAgent()(state, mode="initial")
    sections = {
        element["section_id"]: element
        for element in result["initial_layout_data"]
        if element.get("type") == "section_container"
    }

    assert sections["intro"]["x"] == sections["method"]["x"] == sections["results"]["x"]
    assert sections["intro"]["y"] < sections["method"]["y"] < sections["results"]["y"]


def test_micro_layout_refiner_packs_overflowing_lane_without_lane_overflow():
    state = create_state("/tmp/paper.pdf", layout_template="three_column_postergen")
    state["layout_template_metadata"] = LayoutTemplates(54, 36, margin=1.0, col_gap=1.0).get_template(
        "three_column_postergen",
        header_height=(36 - 2) * 0.18,
    )
    state["styled_layout"] = [
        {
            "type": "section_container",
            "section_id": "s1",
            "lane_id": "left",
            "x": 1.0,
            "y": 7.12,
            "width": 16.66,
            "height": 12.0,
            "importance_level": 2,
            "priority": 0.1,
        },
        {
            "type": "text",
            "id": "s1_text",
            "x": 1.3,
            "y": 8.0,
            "width": 16.06,
            "height": 11.0,
            "content": "Point A\nPoint B\nPoint C\nPoint D\nPoint E\nPoint F",
            "font_family": "Arial",
            "font_size": 44,
            "font_color": "#000000",
            "priority": 0.5,
        },
        {
            "type": "section_container",
            "section_id": "s2",
            "lane_id": "left",
            "x": 1.0,
            "y": 20.5,
            "width": 16.66,
            "height": 12.0,
            "importance_level": 2,
            "priority": 0.1,
        },
        {
            "type": "text",
            "id": "s2_text",
            "x": 1.3,
            "y": 21.3,
            "width": 16.06,
            "height": 10.8,
            "content": "Point A\nPoint B\nPoint C\nPoint D\nPoint E\nPoint F",
            "font_family": "Arial",
            "font_size": 44,
            "font_color": "#000000",
            "priority": 0.5,
        },
    ]

    refiner = MicroLayoutRefiner()
    refined_layout, report = refiner._refine_layout(
        state["styled_layout"],
        state["layout_template_metadata"],
        state,
    )
    result = {"styled_layout": refined_layout}
    left_lane = state["layout_template_metadata"]["lanes"][0]
    left_sections = [
        element for element in result["styled_layout"]
        if element.get("type") == "section_container" and element.get("lane_id") == "left"
    ]

    assert left_sections
    assert max(section["y"] + section["height"] for section in left_sections) <= left_lane["y"] + left_lane["h"] + 0.05


def test_micro_layout_refiner_reports_post_force_fit_overflow():
    refiner = MicroLayoutRefiner()
    state = create_state("/tmp/paper.pdf", layout_template="cluster_43_landscape", width=54, height=27)
    state["template_fast_mode"] = True
    lane = {"id": "slot_1", "x": 1.0, "y": 5.1, "w": 14.04, "h": 6.0}
    group = {
        "section_id": "method",
        "lane_id": "slot_1",
        "container": {
            "type": "section_container",
            "section_id": "method",
            "lane_id": "slot_1",
            "x": lane["x"],
            "y": lane["y"],
            "width": lane["w"],
            "height": 12.0,
        },
        "children": [
            {
                "type": "title_accent_block",
                "section_id": "method",
                "lane_id": "slot_1",
                "x": lane["x"],
                "y": lane["y"],
                "width": lane["w"],
                "height": 0.78,
            },
            {
                "type": "section_title",
                "section_id": "method",
                "lane_id": "slot_1",
                "x": lane["x"] + 0.28,
                "y": lane["y"] + 0.04,
                "width": lane["w"] - 0.56,
                "height": 0.7,
                "font_size": 48,
            },
            {
                "type": "text",
                "id": "method_text",
                "section_id": "method",
                "lane_id": "slot_1",
                "x": lane["x"] + 0.3,
                "y": lane["y"] + 1.3,
                "width": lane["w"] - 0.6,
                "height": 12.0,
                "content": "\n".join(
                    [
                        "Line one with enough words to wrap across the narrow panel and consume height.",
                        "Line two with enough words to wrap across the narrow panel and consume height.",
                        "Line three with enough words to wrap across the narrow panel and consume height.",
                        "Line four with enough words to wrap across the narrow panel and consume height.",
                        "Line five with enough words to wrap across the narrow panel and consume height.",
                    ]
                ),
                "font_family": "Arial",
                "font_size": 44,
                "font_color": "#000000",
            },
        ],
    }

    result = refiner._refine_lane(
        [group],
        lane,
        state,
        {"template_name": "cluster_43_landscape", "orientation": "landscape", "layout_mode": "template_prior", "lanes": [lane]},
    )
    report = result["report"]

    assert report["force_fit_used"] is True
    assert report["pre_force_fit_overflow"] > report["final_overflow"]
    assert report["final_overflow"] == pytest.approx(refiner._lane_overflow(result["elements"], lane))


def test_micro_layout_refiner_expands_underfilled_lane():
    state = create_state("/tmp/paper.pdf", layout_template="three_column_postergen")
    state["layout_template_metadata"] = LayoutTemplates(54, 36, margin=1.0, col_gap=1.0).get_template(
        "three_column_postergen",
        header_height=(36 - 2) * 0.18,
    )
    left_lane = state["layout_template_metadata"]["lanes"][0]
    state["styled_layout"] = [
        {
            "type": "section_container",
            "section_id": "s1",
            "lane_id": "left",
            "x": left_lane["x"],
            "y": left_lane["y"],
            "width": left_lane["w"],
            "height": 4.0,
            "importance_level": 2,
            "priority": 0.1,
        },
        {
            "type": "section_title",
            "id": "s1_title",
            "x": left_lane["x"] + 0.2,
            "y": left_lane["y"] + 0.2,
            "width": left_lane["w"] - 0.4,
            "height": 1.0,
            "font_size": 64,
            "priority": 0.2,
        },
        {
            "type": "text",
            "id": "s1_text",
            "x": left_lane["x"] + 0.3,
            "y": left_lane["y"] + 1.4,
            "width": left_lane["w"] - 0.6,
            "height": 2.0,
            "content": "Point A\nPoint B\nPoint C",
            "font_family": "Arial",
            "font_size": 44,
            "font_color": "#000000",
            "priority": 0.5,
        },
        {
            "type": "section_container",
            "section_id": "s2",
            "lane_id": "left",
            "x": left_lane["x"],
            "y": left_lane["y"] + 5.0,
            "width": left_lane["w"],
            "height": 4.0,
            "importance_level": 2,
            "priority": 0.1,
        },
        {
            "type": "section_title",
            "id": "s2_title",
            "x": left_lane["x"] + 0.2,
            "y": left_lane["y"] + 5.2,
            "width": left_lane["w"] - 0.4,
            "height": 1.0,
            "font_size": 64,
            "priority": 0.2,
        },
        {
            "type": "text",
            "id": "s2_text",
            "x": left_lane["x"] + 0.3,
            "y": left_lane["y"] + 6.4,
            "width": left_lane["w"] - 0.6,
            "height": 2.0,
            "content": "Point D\nPoint E\nPoint F",
            "font_family": "Arial",
            "font_size": 44,
            "font_color": "#000000",
            "priority": 0.5,
        },
    ]

    original_bottom = max(
        element["y"] + element["height"]
        for element in state["styled_layout"]
        if element.get("type") == "section_container"
    )
    result = MicroLayoutRefiner()(state)
    refined_sections = [
        element for element in result["styled_layout"]
        if element.get("type") == "section_container" and element.get("lane_id") == "left"
    ]
    refined_bottom = max(section["y"] + section["height"] for section in refined_sections)

    assert refined_bottom > original_bottom
    assert refined_bottom <= left_lane["y"] + left_lane["h"] + 0.05


def test_micro_layout_refiner_stretches_title_bar_to_full_block_width():
    refiner = MicroLayoutRefiner()
    state = create_state("/tmp/paper.pdf", layout_template="cluster_104_landscape")
    lane = {"id": "slot_5", "x": 1.0, "y": 5.0, "w": 10.0, "h": 10.0}
    group = {
        "section_id": "slot_5_eval_setup",
        "container": {
            "type": "section_container",
            "section_id": "slot_5_eval_setup",
            "lane_id": "slot_5",
            "x": lane["x"],
            "y": lane["y"],
            "width": lane["w"],
            "height": 3.0,
        },
        "children": [
            {
                "type": "title_accent_block",
                "section_id": "slot_5_eval_setup",
                "lane_id": "slot_5",
                "x": lane["x"] + 0.4,
                "y": lane["y"],
                "width": lane["w"] - 0.8,
                "height": 0.78,
            },
            {
                "type": "section_title",
                "section_id": "slot_5_eval_setup",
                "lane_id": "slot_5",
                "x": lane["x"] + 0.28,
                "y": lane["y"] + 0.04,
                "width": lane["w"] - 0.56,
                "height": 0.7,
                "font_size": 48,
            },
        ],
    }
    params = {
        "section_gap": 0.5,
        "title_to_content_gap": 0.25,
        "visual_gap": 0.18,
        "text_padding": 0.24,
        "body_font_reduction": 0,
        "title_font_reduction": 0,
        "body_font_boost": 0,
        "title_font_boost": 4,
        "visual_scale": 1.0,
    }

    elements, _ = refiner._layout_section(group, lane, lane["y"], state, params, {"template_name": "cluster_104_landscape"})

    bar = next(element for element in elements if element["type"] == "title_accent_block")
    assert bar["x"] == pytest.approx(lane["x"])
    assert bar["width"] == pytest.approx(lane["w"])


def test_micro_layout_refiner_force_fit_keeps_title_bar_full_width():
    refiner = MicroLayoutRefiner()
    state = create_state("/tmp/paper.pdf", layout_template="cluster_104_landscape")
    lane = {"id": "slot_5", "x": 1.0, "y": 5.0, "w": 10.0, "h": 5.0}
    lane_layout = [
        {
            "type": "section_container",
            "section_id": "slot_5_eval_setup",
            "lane_id": "slot_5",
            "x": lane["x"],
            "y": lane["y"],
            "width": lane["w"],
            "height": 6.0,
        },
        {
            "type": "title_accent_block",
            "section_id": "slot_5_eval_setup",
            "lane_id": "slot_5",
            "x": lane["x"],
            "y": lane["y"],
            "width": lane["w"],
            "height": 0.78,
        },
        {
            "type": "text",
            "section_id": "slot_5_eval_setup",
            "lane_id": "slot_5",
            "x": lane["x"] + 0.24,
            "y": lane["y"] + 1.0,
            "width": lane["w"] - 0.48,
            "height": 5.2,
            "font_size": 44,
        },
    ]

    compressed = refiner._force_fit_lane(
        lane_layout,
        lane,
        state,
        {"template_name": "cluster_104_landscape", "orientation": "landscape"},
    )

    bar = next(element for element in compressed if element["type"] == "title_accent_block")
    assert bar["x"] == pytest.approx(lane["x"])
    assert bar["width"] == pytest.approx(lane["w"])


def test_micro_layout_refiner_preserves_visual_scale_floor_for_key_visual():
    refiner = MicroLayoutRefiner()
    state = create_state("/tmp/paper.pdf", layout_template="cluster_104_landscape")
    state["template_fast_mode"] = True
    groups = [
        {
            "container": {"importance_level": 1},
            "children": [{"type": "visual", "visual_id": "figure_1"}],
        }
    ]
    params = {
        "section_gap": 0.5,
        "title_to_content_gap": 0.25,
        "visual_gap": 0.18,
        "text_padding": 0.24,
        "body_font_reduction": 0,
        "title_font_reduction": 0,
        "body_font_boost": 0,
        "title_font_boost": 0,
        "visual_scale": 1.0,
    }

    tightened = refiner._tighten_params(params, groups, state, {"template_name": "cluster_104_landscape"})

    assert tightened["visual_scale"] == pytest.approx(1.0)


def test_micro_layout_refiner_enlarges_visual_to_footprint_contract():
    refiner = MicroLayoutRefiner()
    state = create_state("/tmp/paper.pdf", layout_template="cluster_104_landscape")
    state["template_fast_mode"] = True
    state["visual_assets"] = {"figure_2": {"asset_type": "figure", "aspect": 2.84}}
    lane = {"id": "slot_2", "x": 1.0, "y": 5.0, "w": 17.16, "h": 7.15}
    group = {
        "section_id": "slot_2_method",
        "container": {
            "type": "section_container",
            "section_id": "slot_2_method",
            "lane_id": "slot_2",
            "x": lane["x"],
            "y": lane["y"],
            "width": lane["w"],
            "height": 4.0,
            "importance_level": 1,
        },
        "children": [
            {
                "type": "visual",
                "id": "slot_2_method_visual_figure_2",
                "visual_id": "figure_2",
                "section_id": "slot_2_method",
                "lane_id": "slot_2",
                "x": lane["x"] + 4.0,
                "y": lane["y"],
                "width": 7.84,
                "height": 2.76,
            }
        ],
    }
    params = {
        "section_gap": 0.5,
        "title_to_content_gap": 0.25,
        "visual_gap": 0.18,
        "text_padding": 0.24,
        "body_font_reduction": 0,
        "title_font_reduction": 0,
        "body_font_boost": 0,
        "title_font_boost": 0,
        "visual_scale": 0.85,
    }

    elements, _ = refiner._layout_section(group, lane, lane["y"], state, params, {"template_name": "cluster_104_landscape"})

    visual = next(element for element in elements if element.get("type") == "visual")
    assert visual["width"] >= 10.5
    assert visual["visual_footprint"]["ok"] is True


def test_micro_layout_refiner_uses_portrait_split_for_wide_shallow_figure_block():
    refiner = MicroLayoutRefiner()
    state = create_state("/tmp/paper.pdf", layout_template="cluster_8_portrait", width=36, height=50.88)
    state["template_fast_mode"] = True
    state["visual_assets"] = {"figure_1": {"asset_type": "figure", "aspect": 2.52}}
    lane = {"id": "slot_2", "x": 1.0, "y": 15.32, "w": 34.0, "h": 5.68}
    group = {
        "section_id": "sec_core_hags",
        "container": {
            "type": "section_container",
            "section_id": "sec_core_hags",
            "lane_id": "slot_2",
            "x": lane["x"],
            "y": lane["y"],
            "width": lane["w"],
            "height": lane["h"],
            "importance_level": 1,
        },
        "children": [
            {
                "type": "title_accent_block",
                "section_id": "sec_core_hags",
                "lane_id": "slot_2",
                "x": lane["x"],
                "y": lane["y"],
                "width": lane["w"],
                "height": 0.78,
            },
            {
                "type": "section_title",
                "section_id": "sec_core_hags",
                "lane_id": "slot_2",
                "x": lane["x"] + 0.28,
                "y": lane["y"] + 0.04,
                "width": lane["w"] - 0.56,
                "height": 0.7,
                "font_size": 48,
            },
            {
                "type": "visual",
                "id": "sec_core_hags_figure_1",
                "visual_id": "figure_1",
                "section_id": "sec_core_hags",
                "lane_id": "slot_2",
                "x": lane["x"] + 12.75,
                "y": lane["y"] + 1.1,
                "width": 8.5,
                "height": 3.37,
            },
            {
                "type": "text",
                "id": "sec_core_hags_text",
                "section_id": "sec_core_hags",
                "lane_id": "slot_2",
                "x": lane["x"] + 0.24,
                "y": lane["y"] + 4.7,
                "width": lane["w"] - 0.48,
                "height": 0.8,
                "font_size": 44,
                "content": "AGS updates beliefs online as new eviction labels are discovered.",
            },
        ],
    }
    params = {
        "section_gap": 0.5,
        "title_to_content_gap": 0.25,
        "visual_gap": 0.18,
        "text_padding": 0.24,
        "body_font_reduction": 0,
        "title_font_reduction": 0,
        "body_font_boost": 0,
        "title_font_boost": 0,
        "visual_scale": 0.95,
    }

    elements, _ = refiner._layout_section(
        group,
        lane,
        lane["y"],
        state,
        params,
        {"template_name": "cluster_8_portrait", "orientation": "portrait"},
    )

    visual = next(element for element in elements if element.get("type") == "visual")
    text = next(element for element in elements if element.get("type") == "text")
    assert visual["portrait_split_layout"] == "image_left_text_right"
    assert visual["width"] > 10.5
    assert visual["height"] > 4.0
    assert text["x"] > visual["x"] + visual["width"]
    assert visual["visual_footprint"]["ok"] is True


def test_micro_layout_refiner_uses_portrait_split_for_moderately_wide_portrait_block():
    refiner = MicroLayoutRefiner()
    state = create_state("/tmp/paper.pdf", layout_template="cluster_8_portrait", width=36, height=50.88)
    state["template_fast_mode"] = True
    state["visual_assets"] = {"figure_2": {"asset_type": "figure", "aspect": 2.84}}
    lane = {"id": "slot_3", "x": 1.0, "y": 21.088, "w": 34.0, "h": 15.2493}
    group = {
        "section_id": "sec_hags",
        "container": {
            "type": "section_container",
            "section_id": "sec_hags",
            "lane_id": "slot_3",
            "x": lane["x"],
            "y": lane["y"],
            "width": lane["w"],
            "height": lane["h"],
            "importance_level": 1,
        },
        "children": [
            {
                "type": "title_accent_block",
                "section_id": "sec_hags",
                "lane_id": "slot_3",
                "x": lane["x"],
                "y": lane["y"],
                "width": lane["w"],
                "height": 0.78,
            },
            {
                "type": "section_title",
                "section_id": "sec_hags",
                "lane_id": "slot_3",
                "x": lane["x"] + 0.28,
                "y": lane["y"] + 0.04,
                "width": lane["w"] - 0.56,
                "height": 0.7,
                "font_size": 48,
            },
            {
                "type": "visual",
                "id": "sec_hags_figure_2",
                "visual_id": "figure_2",
                "section_id": "sec_hags",
                "lane_id": "slot_3",
                "x": lane["x"] + 5.1,
                "y": lane["y"] + 1.1,
                "width": 23.8,
                "height": 8.38,
            },
            {
                "type": "text",
                "id": "sec_hags_text",
                "section_id": "sec_hags",
                "lane_id": "slot_3",
                "x": lane["x"] + 0.24,
                "y": lane["y"] + 9.8,
                "width": lane["w"] - 0.48,
                "height": 2.0,
                "font_size": 44,
                "content": (
                    "HAGS first chooses a region and then selects a parcel within that region.\n"
                    "The hierarchy shares predictors across regions and reduces city-scale search complexity."
                ),
            },
        ],
    }
    params = {
        "section_gap": 0.5,
        "title_to_content_gap": 0.25,
        "visual_gap": 0.18,
        "text_padding": 0.24,
        "body_font_reduction": 0,
        "title_font_reduction": 0,
        "body_font_boost": 0,
        "title_font_boost": 0,
        "visual_scale": 0.95,
    }

    elements, _ = refiner._layout_section(
        group,
        lane,
        lane["y"],
        state,
        params,
        {"template_name": "cluster_8_portrait", "orientation": "portrait", "layout_mode": "template_prior"},
    )

    visual = next(element for element in elements if element.get("type") == "visual")
    text = next(element for element in elements if element.get("type") == "text")
    assert visual["portrait_split_layout"] == "image_left_text_right"
    assert visual["x"] < text["x"]
    assert visual["width"] >= 23.0
    assert text["x"] > visual["x"] + visual["width"]
    top_gap = visual["y"] - text["y"]
    bottom_gap = lane["y"] + lane["h"] - 0.03 - (visual["y"] + visual["height"])
    assert top_gap > 1.0
    assert abs(top_gap - bottom_gap) <= 0.05
    assert visual["visual_footprint"]["ok"] is True
    content_bottom = max(
        element.get("y", 0) + element.get("height", 0)
        for element in elements
        if element.get("type") in {"section_title", "title_accent_block", "text", "visual"}
    )
    assert content_bottom <= lane["y"] + lane["h"]


def test_micro_layout_refiner_keeps_generated_teaser_as_banner_in_portrait_block():
    refiner = MicroLayoutRefiner()
    state = create_state("/tmp/paper.pdf", layout_template="cluster_8_portrait", width=36, height=50.88)
    state["template_fast_mode"] = True
    state["visual_assets"] = {"generated_teaser_1": {"asset_type": "figure", "aspect": 6.66}}
    lane = {"id": "slot_1", "x": 1.0, "y": 7.46, "w": 34.0, "h": 7.8612}
    group = {
        "section_id": "sec_problem",
        "container": {
            "type": "section_container",
            "section_id": "sec_problem",
            "lane_id": "slot_1",
            "x": lane["x"],
            "y": lane["y"],
            "width": lane["w"],
            "height": lane["h"],
            "importance_level": 1,
        },
        "children": [
            {
                "type": "title_accent_block",
                "section_id": "sec_problem",
                "lane_id": "slot_1",
                "x": lane["x"],
                "y": lane["y"],
                "width": lane["w"],
                "height": 0.78,
            },
            {
                "type": "section_title",
                "section_id": "sec_problem",
                "lane_id": "slot_1",
                "x": lane["x"] + 0.28,
                "y": lane["y"] + 0.04,
                "width": lane["w"] - 0.56,
                "height": 0.7,
                "font_size": 48,
            },
            {
                "type": "visual",
                "id": "sec_problem_generated_teaser_1",
                "visual_id": "generated_teaser_1",
                "section_id": "sec_problem",
                "lane_id": "slot_1",
                "x": lane["x"] + 1.0,
                "y": lane["y"] + 1.1,
                "width": 31.73,
                "height": 4.76,
            },
            {
                "type": "text",
                "id": "sec_problem_text",
                "section_id": "sec_problem",
                "lane_id": "slot_1",
                "x": lane["x"] + 0.24,
                "y": lane["y"] + 6.0,
                "width": lane["w"] - 0.48,
                "height": 0.8,
                "font_size": 44,
                "content": "Eviction outreach must identify at-risk properties quickly under a limited budget.",
            },
        ],
    }
    params = {
        "section_gap": 0.5,
        "title_to_content_gap": 0.25,
        "visual_gap": 0.18,
        "text_padding": 0.24,
        "body_font_reduction": 0,
        "title_font_reduction": 0,
        "body_font_boost": 0,
        "title_font_boost": 0,
        "visual_scale": 0.95,
    }

    elements, _ = refiner._layout_section(
        group,
        lane,
        lane["y"],
        state,
        params,
        {"template_name": "cluster_8_portrait", "orientation": "portrait", "layout_mode": "template_prior"},
    )

    visual = next(element for element in elements if element.get("type") == "visual")
    text = next(element for element in elements if element.get("type") == "text")
    assert "portrait_split_layout" not in visual
    assert visual["width"] >= 30.0
    assert text["y"] > visual["y"] + visual["height"]


def test_micro_layout_refiner_does_not_fill_portrait_split_text_column_from_source():
    refiner = MicroLayoutRefiner()
    state = create_state("/tmp/paper.pdf", layout_template="cluster_8_portrait", width=36, height=50.88)
    state["template_fast_mode"] = True
    state["visual_assets"] = {"figure_1": {"asset_type": "figure", "aspect": 2.52}}
    state["story_board"] = {
        "spatial_content_plan": {
            "sections": [
                {
                    "section_id": "sec_core_hags",
                    "section_title": "Adaptive Policy",
                    "content_role": "method",
                    "source_sections": ["Method"],
                    "text_content": [
                        "AGS updates beliefs online as new eviction labels are discovered.",
                        "The predictor estimates parcel risk from tabular records and overhead imagery.",
                        "The search policy chooses the next parcel under remaining budget and observed labels.",
                    ],
                }
            ]
        }
    }
    state["structured_sections"] = {
        "paper_sections": [
            {
                "section_name": "Method",
                "section_type": "method",
                "content": (
                    "The adaptive geospatial search policy alternates between querying promising parcels "
                    "and updating the risk model with newly observed eviction labels. The prediction module "
                    "scores parcels from administrative features, neighborhood context, and imagery. The policy "
                    "uses the remaining outreach budget to balance exploration of uncertain areas with exploitation "
                    "of parcels that already appear high risk."
                ),
                "key_points": [
                    "Each query changes both the classifier and the search policy.",
                    "Adaptive search balances exploration with immediate outreach yield.",
                    "The budget-aware policy avoids spending visits in low-yield areas.",
                ],
            }
        ]
    }
    lane = {"id": "slot_2", "x": 1.0, "y": 15.32, "w": 34.0, "h": 5.68}
    group = {
        "section_id": "sec_core_hags",
        "container": {
            "type": "section_container",
            "section_id": "sec_core_hags",
            "lane_id": "slot_2",
            "x": lane["x"],
            "y": lane["y"],
            "width": lane["w"],
            "height": lane["h"],
            "importance_level": 1,
        },
        "children": [
            {
                "type": "title_accent_block",
                "section_id": "sec_core_hags",
                "lane_id": "slot_2",
                "x": lane["x"],
                "y": lane["y"],
                "width": lane["w"],
                "height": 0.78,
            },
            {
                "type": "section_title",
                "section_id": "sec_core_hags",
                "lane_id": "slot_2",
                "x": lane["x"] + 0.28,
                "y": lane["y"] + 0.04,
                "width": lane["w"] - 0.56,
                "height": 0.7,
                "font_size": 48,
            },
            {
                "type": "visual",
                "id": "sec_core_hags_figure_1",
                "visual_id": "figure_1",
                "section_id": "sec_core_hags",
                "lane_id": "slot_2",
                "x": lane["x"] + 12.75,
                "y": lane["y"] + 1.1,
                "width": 8.5,
                "height": 3.37,
            },
            {
                "type": "text",
                "id": "sec_core_hags_text",
                "section_id": "sec_core_hags",
                "lane_id": "slot_2",
                "x": lane["x"] + 0.24,
                "y": lane["y"] + 4.7,
                "width": lane["w"] - 0.48,
                "height": 0.8,
                "font_size": 44,
                "content": "AGS updates beliefs online as new eviction labels are discovered.",
            },
        ],
    }
    params = {
        "section_gap": 0.5,
        "title_to_content_gap": 0.25,
        "visual_gap": 0.18,
        "text_padding": 0.24,
        "body_font_reduction": 0,
        "title_font_reduction": 0,
        "body_font_boost": 0,
        "title_font_boost": 0,
        "visual_scale": 0.95,
    }

    elements, _ = refiner._layout_section(
        group,
        lane,
        lane["y"],
        state,
        params,
        {"template_name": "cluster_8_portrait", "orientation": "portrait", "layout_mode": "template_prior"},
    )

    text = next(element for element in elements if element.get("type") == "text")
    assert text["portrait_split_layout"] == "image_left_text_right"
    assert text["portrait_split_text_fill_ratio"] < 0.90
    assert text["content"].splitlines() == ["AGS updates beliefs online as new eviction labels are discovered."]
    assert "predictor estimates" not in text["content"].lower()


def test_micro_layout_refiner_portrait_wide_columns_fill_bottom_with_dense_content():
    refiner = MicroLayoutRefiner()
    state = create_state("/tmp/paper.pdf", layout_template="cluster_13_portrait", width=36, height=51)
    section_lines = [
        "Evaluation uses 16,000 St. Louis residential rental properties with tabular records and overhead imagery.",
        "<color:#7A1F2B>HAGS</color> gives the strongest yield across all shown search budgets.",
        "The gains are largest when outreach budgets are tight and every visit must count.",
        "Compared with greedy and conventional active-search baselines, HAGS keeps improving as the budget grows.",
        "Operational takeaway: prioritize regions first, then parcels, to turn scarce canvassing visits into more tenant-support opportunities.",
        "In uniform-cost search, HAGS consistently outperforms all baselines, including conventional active search, greedy adaptive search, and the non-hierarchical AGS model.",
    ]
    state["story_board"] = {
        "spatial_content_plan": {
            "sections": [
                {
                    "section_id": "sec_results",
                    "section_title": "City Results",
                    "content_role": "results",
                    "text_content": section_lines,
                    "source_sections": ["Results"],
                }
            ]
        }
    }
    state["structured_sections"] = {
        "paper_sections": [
            {
                "section_name": "Results",
                "section_type": "results",
                "key_points": [
                    "The results are for positive rates of 5% and 10%, respectively.",
                    "HAGS stays strongest across the displayed query budgets.",
                    "The hierarchy matters most when outreach visits are scarce.",
                    "Region-first search avoids spending budget in low-yield areas.",
                    "Parcel-level updates exploit each new observation immediately.",
                    "Travel-aware settings amplify the gains over greedy baselines.",
                ],
            }
        ]
    }
    lane = {"id": "slot_4", "x": 1.0, "y": 30.0, "w": 34.0, "h": 20.0}
    params = {
        "section_gap": 1.0,
        "title_to_content_gap": 0.4,
        "visual_gap": 0.3,
        "text_padding": 0.3,
        "body_font_reduction": 0,
        "title_font_reduction": 0,
        "body_font_boost": 0,
        "title_font_boost": 0,
        "visual_scale": 1.0,
    }
    text_elements = [
        {
            "type": "text",
            "id": "sec_results_text_col_1",
            "section_id": "sec_results",
            "lane_id": "slot_4",
            "content": "\n".join(section_lines[:3]),
            "font_family": "Arial",
            "font_size": 40,
            "line_spacing": 0.96,
        },
        {
            "type": "text",
            "id": "sec_results_text_col_2",
            "section_id": "sec_results",
            "lane_id": "slot_4",
            "content": "\n".join(section_lines[3:]),
            "font_family": "Arial",
            "font_size": 40,
            "line_spacing": 0.96,
        }
    ]

    result = refiner._layout_wide_text_columns_for_fill(
        text_elements,
        lane,
        current_y=39.5,
        state=state,
        params=params,
        template_layout={"template_name": "cluster_13_portrait", "orientation": "portrait", "layout_mode": "template_prior"},
        section_id="sec_results",
        is_last_group=True,
    )

    assert result is not None
    column_elements, _, content_bottom = result
    assert lane["y"] + lane["h"] - content_bottom <= 0.18
    combined = "\n".join(element["content"] for element in column_elements)
    assert sum(len(element["content"].splitlines()) for element in column_elements) == len(section_lines)
    assert "Travel-aware settings amplify" not in combined


def test_micro_layout_refiner_portrait_expands_visual_to_absorb_bottom_gap():
    refiner = MicroLayoutRefiner()
    lane = {"id": "slot_3", "x": 1.0, "y": 20.0, "w": 34.0, "h": 12.0}
    elements = [
        {
            "type": "visual",
            "id": "method_figure",
            "section_id": "method",
            "lane_id": "slot_3",
            "x": 9.0,
            "y": 22.0,
            "width": 16.0,
            "height": 5.0,
        },
        {
            "type": "text",
            "id": "method_text",
            "section_id": "method",
            "lane_id": "slot_3",
            "x": 1.3,
            "y": 27.3,
            "width": 33.4,
            "height": 4.0,
        },
    ]
    params = {"text_padding": 0.3}

    updated, content_bottom = refiner._portrait_expand_visual_to_absorb_bottom_gap(
        elements,
        lane,
        params,
        {"template_name": "cluster_8_portrait", "orientation": "portrait"},
        content_bottom=31.3,
    )

    visual = next(element for element in updated if element["type"] == "visual")
    text = next(element for element in updated if element["type"] == "text")
    assert visual["height"] > 5.0
    assert text["y"] > 27.3
    assert lane["y"] + lane["h"] - content_bottom <= 0.18


def test_micro_layout_refiner_portrait_text_fill_uses_fine_line_spacing_step():
    refiner = MicroLayoutRefiner()
    text_element = {
        "type": "text",
        "id": "sec_core_method_text",
        "section_id": "sec_core_method",
        "lane_id": "slot_2",
        "width": 16.4,
        "content": "\n".join(
            [
                "<color:#7A1F2B>AGS</color> couples exploration and exploitation: each query discovers a current eviction-risk case and improves future predictions.",
                "The predictor estimates parcel risk from tabular records and overhead imagery.",
                "The search policy chooses the next parcel under remaining budget and observed labels.",
                "Each query updates the model, so the route can shift as new evidence arrives.",
                "Region-level choices avoid spending visits in low-yield parts of the city.",
                "AGS explicitly couples exploration and exploitation: each query both discovers a current eviction-risk case and improves future predictions from newly observed labels.",
                "Predict risk; query a parcel, then update the route.",
            ]
        ),
        "font_family": "Arial",
        "font_size": 42,
        "line_spacing": 0.98,
    }
    state = {"story_board": {}, "structured_sections": {}, "narrative_content": {}}

    updated, font_size, line_spacing = refiner._expand_text_content_to_fill(
        text_element,
        "sec_core_method",
        state,
        text_width=16.4,
        font_size=42,
        line_spacing=0.98,
        max_height=12.7025,
        template_layout={"template_name": "cluster_13_portrait", "orientation": "portrait", "layout_mode": "template_prior"},
    )

    assert updated["content"] == text_element["content"]
    assert font_size == 42
    assert line_spacing == pytest.approx(1.01)


def test_micro_layout_refiner_caps_body_font_in_underfilled_template_slots():
    refiner = MicroLayoutRefiner()
    lane = {"id": "slot_3", "x": 1.0, "y": 2.0, "w": 16.0, "h": 12.0}
    group = {
        "section_id": "sec_application",
        "lane_id": "slot_3",
        "container": {
            "type": "section_container",
            "section_id": "sec_application",
            "lane_id": "slot_3",
            "slot_id": "slot_3",
            "x": 1.0,
            "y": 2.0,
            "width": 16.0,
            "height": 6.0,
        },
        "children": [
            {
                "type": "title_accent_block",
                "section_id": "sec_application",
                "lane_id": "slot_3",
                "slot_id": "slot_3",
                "x": 1.0,
                "y": 2.0,
                "width": 16.0,
                "height": 0.7,
            },
            {
                "type": "section_title",
                "id": "sec_application_title",
                "section_id": "sec_application",
                "lane_id": "slot_3",
                "slot_id": "slot_3",
                "x": 1.2,
                "y": 2.05,
                "width": 15.6,
                "height": 0.7,
                "font_size": 48,
                "content": "Application",
            },
            {
                "type": "visual",
                "id": "sec_application_figure",
                "section_id": "sec_application",
                "lane_id": "slot_3",
                "slot_id": "slot_3",
                "x": 1.3,
                "y": 3.2,
                "width": 14.0,
                "height": 4.0,
            },
            {
                "type": "text",
                "id": "sec_application_text",
                "section_id": "sec_application",
                "lane_id": "slot_3",
                "slot_id": "slot_3",
                "x": 1.3,
                "y": 7.5,
                "width": 15.4,
                "height": 1.0,
                "font_size": 44,
                "font_family": "Arial",
                "line_spacing": 1.0,
                "content": "Short application takeaway.",
            },
        ],
    }
    state = {
        "poster_width": 54,
        "poster_height": 27,
        "story_board": {},
        "structured_sections": {},
        "narrative_content": {},
        "template_fast_mode": True,
    }

    result = refiner._refine_lane(
        [group],
        lane,
        state,
        {
            "template_name": "cluster_43_landscape",
            "orientation": "landscape",
            "layout_mode": "template_prior",
            "lanes": [lane],
        },
    )

    body_text = next(element for element in result["elements"] if element.get("id") == "sec_application_text")
    max_body_font = load_config()["micro_layout_refinement"]["max_body_font_size"]
    assert max_body_font <= 48
    assert body_text["font_size"] <= max_body_font
    assert result["report"]["underflow_expanded"] is True


def test_micro_layout_refiner_does_not_create_semantic_bottom_fill_text():
    refiner = MicroLayoutRefiner()
    group = {"section_id": "sec_results", "lane_id": "slot_1"}
    lane = {"id": "slot_1", "x": 1.0, "y": 2.0, "w": 14.0, "h": 12.0}
    state = {
        "story_board": {
            "spatial_content_plan": {
                "sections": [
                    {
                        "section_id": "sec_results",
                        "section_title": "Results",
                        "content_role": "results",
                        "text_content": ["Existing result statement from the rewritten block."],
                    }
                ]
            }
        }
    }

    fill = refiner._bottom_fill_elements(
        group,
        lane,
        state,
        {"text_padding": 0.3},
        {"template_name": "cluster_86_landscape", "orientation": "landscape", "layout_mode": "template_prior"},
        content_bottom=6.0,
    )

    assert fill == []


def test_micro_layout_refiner_real_content_fill_reuses_only_existing_lines_by_default():
    refiner = MicroLayoutRefiner()
    content = "The rewritten block explains how HAGS improves parcel search decisions under limited outreach budgets."
    state = {
        "story_board": {
            "spatial_content_plan": {
                "sections": [
                    {
                        "section_id": "sec_results",
                        "section_title": "Results",
                        "content_role": "results",
                        "text_content": [content],
                    }
                ]
            }
        },
        "structured_sections": {
            "paper_sections": [
                {
                    "section_name": "Results",
                    "section_type": "results",
                    "content": "A separate source sentence should not be injected during micro layout.",
                }
            ]
        },
    }

    lines = refiner._content_lines_for_fill("sec_results", state, content, max_items=5)

    assert lines == [content]
    assert all("Operational takeaway" not in line for line in lines)


def test_micro_layout_refiner_teaser_fill_rejects_unrelated_source_sentences():
    refiner = MicroLayoutRefiner()
    content = (
        "Eviction outreach requires sequential, budget-constrained search under uncertainty.\n"
        "AGS formalizes the exploration-exploitation tradeoff with geospatial and travel-cost considerations."
    )
    state = {
        "story_board": {
            "spatial_content_plan": {
                "sections": [
                    {
                        "section_id": "motivation",
                        "section_title": "Motivation",
                        "content_role": "foundation",
                        "generated_teaser_summary": True,
                        "text_content": content.splitlines(),
                    }
                ]
            }
        },
        "structured_sections": {
            "paper_sections": [
                {
                    "section_name": "Introduction",
                    "section_type": "foundation",
                    "content": (
                        "HAGS is introduced to make AGS scalable to large urban search spaces. "
                        "Prior work developed non-myopic and cost-effective active search methods."
                    ),
                }
            ]
        },
    }

    lines = refiner._content_lines_for_fill("motivation", state, content, max_items=5)

    joined = "\n".join(lines)
    assert "HAGS is introduced" not in joined
    assert "Prior work developed" not in joined


def test_micro_layout_refiner_preserves_existing_concise_fill_lines():
    refiner = MicroLayoutRefiner()
    lines = [
        "AGS formulates outreach as budget-constrained sequential geospatial search.",
        "The policy combines a prediction module.",
        "Hierarchy is the main reason the method succeeds at urban scale.",
        "Overall empirical conclusion: HAGS is the strongest method across cost models, budgets, and target rates.",
    ]

    cleaned = refiner._clean_existing_fill_content("\n".join(lines))

    assert cleaned.splitlines() == lines


def test_micro_layout_refiner_repairs_existing_truncated_fill_tail():
    refiner = MicroLayoutRefiner()

    cleaned = refiner._clean_existing_fill_content(
        "Overall empirical conclusion: HAGS is the strongest method across cost."
    )

    assert cleaned == "Overall empirical conclusion: HAGS is the strongest method."


def test_micro_layout_refiner_uses_conservative_portrait_text_measurement():
    refiner = MicroLayoutRefiner()
    text = (
        "AGS uses prediction-and-search decomposition to select parcel queries under budget.\n"
        "The policy adapts online as new labels are discovered and balances exploration with exploitation."
    )

    landscape_height = refiner._estimate_text_height_fast(
        text,
        width_inches=12.0,
        font_size=44,
        line_spacing=1.0,
        template_layout={"template_name": "cluster_104_landscape", "orientation": "landscape"},
    )
    portrait_height = refiner._estimate_text_height_fast(
        text,
        width_inches=12.0,
        font_size=44,
        line_spacing=1.0,
        template_layout={"template_name": "cluster_22_portrait", "orientation": "portrait"},
    )

    assert portrait_height >= landscape_height


def test_micro_layout_refiner_defaults_to_fast_text_measurement_for_landscape(monkeypatch):
    refiner = MicroLayoutRefiner()

    def fail_slow_measurement(**_kwargs):
        raise AssertionError("slow python-pptx text measurement should not run by default")

    monkeypatch.setattr(
        "src.agents.micro_layout_refiner.measure_text_height",
        fail_slow_measurement,
    )

    measured = refiner._measure_text_height_for_refinement(
        text_content="A concise landscape block line.\nA second line used for sizing.",
        width_inches=12.0,
        font_name="Arial",
        font_size=44,
        line_spacing=1.0,
        template_layout={"template_name": "cluster_104_landscape", "orientation": "landscape"},
    )

    assert measured["optimal_height"] > 0


def test_micro_layout_force_fit_preserves_visual_footprint_contract():
    refiner = MicroLayoutRefiner()
    state = create_state("/tmp/paper.pdf", layout_template="cluster_104_landscape", width=54, height=27)
    state["template_fast_mode"] = True
    state["visual_assets"] = {"figure_2": {"asset_type": "figure", "aspect": 2.84}}
    lane = {"id": "slot_2", "x": 18.94, "y": 5.35, "w": 17.16, "h": 7.15}
    layout = [
        {
            "type": "visual",
            "id": "slot_2_method_visual_figure_2",
            "visual_id": "figure_2",
            "section_id": "slot_2_method_visual",
            "lane_id": "slot_2",
            "x": lane["x"] + 3.0,
            "y": lane["y"] + 0.9,
            "width": 11.17,
            "height": 3.93,
        },
        {
            "type": "text",
            "id": "slot_2_method_visual_text",
            "section_id": "slot_2_method_visual",
            "lane_id": "slot_2",
            "x": lane["x"] + 0.3,
            "y": lane["y"] + 6.8,
            "width": 16.5,
            "height": 1.0,
            "font_size": 44,
        },
    ]

    compressed = refiner._force_fit_lane(layout, lane, state, {"template_name": "cluster_104_landscape"})

    visual = next(element for element in compressed if element.get("type") == "visual")
    assert visual["width"] >= 10.63
    assert visual["visual_footprint"]["ok"] is True


def test_micro_layout_force_fit_rechecks_overflow_after_footprint_and_text_floors():
    refiner = MicroLayoutRefiner()
    state = create_state("/tmp/paper.pdf", layout_template="cluster_43_landscape", width=54, height=27)
    state["template_fast_mode"] = True
    state["visual_assets"] = {"generated_teaser_1": {"asset_type": "figure", "aspect": 1.9882}}
    lane = {"id": "slot_1", "x": 1.0, "y": 5.1, "w": 14.04, "h": 10.4}
    text = (
        "Eviction-prevention outreach faces a sequential search problem: canvassers have limited "
        "budget, uncertain parcel-level risk.\n"
        "• The paper formulates this as <color:#1e324d>Active Geospatial Search</color> (AGS): "
        "sequentially query parcels with geographic locations, features."
    )
    layout = [
        {
            "type": "section_container",
            "section_id": "section_1",
            "lane_id": "slot_1",
            "x": lane["x"],
            "y": lane["y"],
            "width": lane["w"],
            "height": 10.695,
        },
        {
            "type": "title_accent_block",
            "section_id": "section_1",
            "lane_id": "slot_1",
            "x": lane["x"],
            "y": lane["y"],
            "width": lane["w"],
            "height": 0.647,
        },
        {
            "type": "section_title",
            "id": "section_1_title",
            "section_id": "section_1",
            "lane_id": "slot_1",
            "x": lane["x"] + 0.28,
            "y": lane["y"] + 0.033,
            "width": lane["w"] - 0.52,
            "height": 0.689,
            "font_size": 46,
        },
        {
            "type": "visual",
            "id": "section_1_generated_teaser_1",
            "visual_id": "generated_teaser_1",
            "section_id": "section_1",
            "lane_id": "slot_1",
            "x": 2.208,
            "y": 5.963,
            "width": 11.624,
            "height": 5.847,
        },
        {
            "type": "text",
            "id": "section_1_text",
            "section_id": "section_1",
            "lane_id": "slot_1",
            "x": lane["x"] + 0.24,
            "y": 11.965,
            "width": lane["w"] - 0.48,
            "height": 3.809,
            "content": text,
            "font_family": "Arial",
            "font_size": 36,
        },
    ]

    compressed = refiner._force_fit_lane(
        layout,
        lane,
        state,
        {"template_name": "cluster_43_landscape", "orientation": "landscape"},
    )

    lane_bottom = lane["y"] + lane["h"]
    assert max(element["y"] + element["height"] for element in compressed) <= lane_bottom + 0.02
    visual = next(element for element in compressed if element.get("type") == "visual")
    assert visual["visual_footprint"]["ok"] is True


def test_micro_layout_refiner_validation_rejects_child_outside_container():
    state = create_state("/tmp/paper.pdf", layout_template="three_column_postergen")
    template = LayoutTemplates(54, 36, margin=1.0, col_gap=1.0).get_template(
        "three_column_postergen",
        header_height=(36 - 2) * 0.18,
    )
    lane_map = {lane["id"]: lane for lane in template["lanes"]}
    lane = lane_map["left"]
    layout = [
        {
            "type": "section_container",
            "section_id": "hags_core",
            "lane_id": "left",
            "x": lane["x"],
            "y": lane["y"],
            "width": lane["w"],
            "height": 3.0,
        },
        {
            "type": "text",
            "id": "hags_core_text",
            "x": lane["x"] + 0.3,
            "y": lane["y"] + 1.0,
            "width": lane["w"] - 0.6,
            "height": 3.5,
        },
    ]

    validation = MicroLayoutRefiner()._validate_refined_layout(layout, lane_map, state)

    assert any("child vertical overflow" in issue for issue in validation["issues"])


def test_micro_layout_refiner_validation_rejects_text_touching_container_bottom():
    state = create_state("/tmp/paper.pdf", layout_template="cluster_86_landscape", width=54, height=27)
    state["layout_template_metadata"] = {
        "template_name": "cluster_test",
        "layout_mode": "template_prior",
        "lanes": [{"id": "slot_1", "x": 1.0, "y": 4.0, "w": 12.0, "h": 8.0}],
    }
    lane_map = {lane["id"]: lane for lane in state["layout_template_metadata"]["lanes"]}
    layout = [
        {
            "type": "section_container",
            "section_id": "method",
            "lane_id": "slot_1",
            "x": 1.0,
            "y": 4.0,
            "width": 12.0,
            "height": 5.0,
        },
        {
            "type": "text",
            "id": "method_text",
            "section_id": "method",
            "lane_id": "slot_1",
            "x": 1.3,
            "y": 5.0,
            "width": 11.4,
            "height": 4.0,
            "font_size": 44,
            "line_spacing": 1.0,
            "content": "A fitted block line.\nAnother fitted block line.",
        },
    ]

    validation = MicroLayoutRefiner()._validate_refined_layout(layout, lane_map, state)

    assert any("text bottom padding too small" in issue for issue in validation["issues"])


def test_micro_layout_refiner_validation_rejects_template_section_overlap():
    state = create_state("/tmp/paper.pdf", layout_template="cluster_96_landscape", width=54, height=27)
    state["layout_template_metadata"] = {
        "template_name": "cluster_test",
        "layout_mode": "template_prior",
        "lanes": [
            {"id": "slot_1", "x": 1.0, "y": 4.0, "w": 12.0, "h": 5.0},
            {"id": "slot_2", "x": 1.0, "y": 8.8, "w": 12.0, "h": 5.0},
        ],
    }
    lane_map = {lane["id"]: lane for lane in state["layout_template_metadata"]["lanes"]}
    layout = [
        {
            "type": "section_container",
            "section_id": "first",
            "lane_id": "slot_1",
            "x": 1.0,
            "y": 4.0,
            "width": 12.0,
            "height": 5.1,
        },
        {
            "type": "section_container",
            "section_id": "second",
            "lane_id": "slot_2",
            "x": 1.0,
            "y": 8.8,
            "width": 12.0,
            "height": 5.0,
        },
    ]

    validation = MicroLayoutRefiner()._validate_refined_layout(layout, lane_map, state)

    assert any("section container overlap" in issue for issue in validation["issues"])


def test_micro_layout_refiner_skips_slot_width_resize_when_it_would_create_overlap(monkeypatch):
    state = create_state("/tmp/paper.pdf", layout_template="cluster_86_landscape", width=54, height=27)
    template = {
        "template_name": "cluster_86_landscape",
        "layout_mode": "template_prior",
        "slot_order": ["slot_1", "slot_3", "slot_4"],
        "lanes": [
            {"id": "slot_1", "x": 18.68, "y": 4.6, "w": 17.16, "h": 9.28},
            {"id": "slot_3", "x": 35.84, "y": 4.6, "w": 17.16, "h": 10.4},
            {"id": "slot_4", "x": 19.2, "y": 13.875, "w": 16.64, "h": 12.0},
        ],
        "adjacency_graph": {
            "slot_3": [{"slot_id": "slot_1", "orientation": "vertical"}],
        },
    }
    lane_map = {lane["id"]: dict(lane) for lane in template["lanes"]}
    section_containers = [
        {"section_id": "method", "lane_id": "slot_1", "container": {"height": 2.0}, "children": []},
        {"section_id": "details", "lane_id": "slot_3", "container": {"height": 2.0}, "children": []},
        {"section_id": "evaluation", "lane_id": "slot_4", "container": {"height": 2.0}, "children": []},
    ]
    demands = {"slot_1": 2.0, "slot_3": 11.0, "slot_4": 8.0}

    def fake_layout_section(group, lane, section_y, state_arg, params, template_arg):
        return [], demands[group["lane_id"]]

    refiner = MicroLayoutRefiner()
    monkeypatch.setattr(refiner, "_layout_section", fake_layout_section)

    updated = refiner._rebalance_template_block_slots(template, section_containers, lane_map, state)

    assert updated["slot_3"]["x"] == pytest.approx(35.84)
    assert updated["slot_3"]["w"] == pytest.approx(17.16)
    assert state["slot_pressure_report"]["slot_resize_applied"] is False
    assert state["slot_pressure_report"]["slot_resize_skipped"] == "resize_would_overlap_template_slots"


def test_vlm_layout_reviewer_syncs_container_after_patch():
    state = create_state("/tmp/paper.pdf", layout_template="three_column_postergen", enable_vlm_layout_review=True)
    template = LayoutTemplates(54, 36, margin=1.0, col_gap=1.0).get_template(
        "three_column_postergen",
        header_height=(36 - 2) * 0.18,
    )
    state["layout_template_metadata"] = template
    lane = template["lanes"][0]
    layout = [
        {
            "type": "section_container",
            "section_id": "hags_core",
            "lane_id": "left",
            "x": lane["x"],
            "y": lane["y"],
            "width": lane["w"],
            "height": 4.0,
        },
        {
            "type": "text",
            "id": "hags_core_text",
            "x": lane["x"] + 0.3,
            "y": lane["y"] + 1.0,
            "width": lane["w"] - 0.6,
            "height": 3.8,
            "font_size": 44,
        },
    ]
    patch = [{"target": "hags_core_text", "op": "increase_font_size", "value": 2}]

    patched = VLMLayoutReviewer()._apply_safe_patch(layout, patch, state)
    container = next(element for element in patched if element.get("type") == "section_container")
    text = next(element for element in patched if element.get("id") == "hags_core_text")

    assert patched is not None
    assert container["y"] + container["height"] >= text["y"] + text["height"]


def test_micro_layout_refiner_handles_two_plus_one_mixed_without_right_lane_overflow():
    state = create_state("/tmp/paper.pdf", layout_template="two_plus_one_mixed")
    state["layout_template_metadata"] = LayoutTemplates(54, 36, margin=1.0, col_gap=1.0).get_template(
        "two_plus_one_mixed",
        header_height=(36 - 2) * 0.18,
    )
    right_lane = state["layout_template_metadata"]["lanes"][2]
    state["styled_layout"] = [
        {
            "type": "section_container",
            "section_id": "results",
            "lane_id": "right",
            "x": right_lane["x"],
            "y": right_lane["y"],
            "width": right_lane["w"],
            "height": right_lane["h"] * 0.75,
            "importance_level": 1,
            "priority": 0.1,
        },
        {
            "type": "section_title",
            "id": "results_title",
            "x": right_lane["x"] + 0.2,
            "y": right_lane["y"] + 0.2,
            "width": right_lane["w"] - 0.4,
            "height": 1.0,
            "font_size": 64,
            "priority": 0.2,
        },
        {
            "type": "visual",
            "id": "results_visual",
            "slot_id": "results_visual",
            "visual_id": "table_1",
            "x": right_lane["x"] + 0.4,
            "y": right_lane["y"] + 1.6,
            "width": right_lane["w"] - 0.8,
            "height": 7.5,
            "priority": 0.4,
        },
        {
            "type": "text",
            "id": "results_text",
            "x": right_lane["x"] + 0.3,
            "y": right_lane["y"] + 9.6,
            "width": right_lane["w"] - 0.6,
            "height": 15.0,
            "content": "\n".join([f"Result point {i}" for i in range(1, 20)]),
            "font_family": "Arial",
            "font_size": 44,
            "font_color": "#000000",
            "priority": 0.5,
        },
    ]

    result = MicroLayoutRefiner()(state)
    right_sections = [
        element for element in result["styled_layout"]
        if element.get("type") == "section_container" and element.get("lane_id") == "right"
    ]

    assert right_sections
    assert max(section["y"] + section["height"] for section in right_sections) <= right_lane["y"] + right_lane["h"] + 0.05


def test_micro_layout_refiner_handles_one_plus_two_mixed_without_left_lane_overflow():
    state = create_state("/tmp/paper.pdf", layout_template="one_plus_two_mixed")
    state["layout_template_metadata"] = LayoutTemplates(54, 36, margin=1.0, col_gap=1.0).get_template(
        "one_plus_two_mixed",
        header_height=(36 - 2) * 0.18,
    )
    left_lane = state["layout_template_metadata"]["lanes"][0]
    state["styled_layout"] = [
        {
            "type": "section_container",
            "section_id": "intro",
            "lane_id": "left",
            "x": left_lane["x"],
            "y": left_lane["y"],
            "width": left_lane["w"],
            "height": left_lane["h"] * 0.82,
            "importance_level": 2,
            "priority": 0.1,
        },
        {
            "type": "section_title",
            "id": "intro_title",
            "x": left_lane["x"] + 0.2,
            "y": left_lane["y"] + 0.2,
            "width": left_lane["w"] - 0.4,
            "height": 1.0,
            "font_size": 64,
            "priority": 0.2,
        },
        {
            "type": "text",
            "id": "intro_text",
            "x": left_lane["x"] + 0.3,
            "y": left_lane["y"] + 1.4,
            "width": left_lane["w"] - 0.6,
            "height": 18.0,
            "content": "\n".join([f"Background point {i}" for i in range(1, 26)]),
            "font_family": "Arial",
            "font_size": 44,
            "font_color": "#000000",
            "priority": 0.5,
        },
        {
            "type": "visual",
            "id": "intro_visual",
            "slot_id": "intro_visual",
            "visual_id": "figure_1",
            "x": left_lane["x"] + 0.5,
            "y": left_lane["y"] + 20.0,
            "width": left_lane["w"] - 1.0,
            "height": 6.5,
            "priority": 0.4,
        },
    ]

    result = MicroLayoutRefiner()(state)
    left_sections = [
        element for element in result["styled_layout"]
        if element.get("type") == "section_container" and element.get("lane_id") == "left"
    ]

    assert left_sections
    assert max(section["y"] + section["height"] for section in left_sections) <= left_lane["y"] + left_lane["h"] + 0.05


def test_template_selector_prefers_balanced_three_column_layout():
    selector = TemplateSelector(load_config())
    state = create_state("/tmp/paper.pdf", layout_template="adaptive_auto")
    state["visual_assets"] = {
        "figure_1": {"aspect": 2.0, "asset_type": "figure"},
        "figure_2": {"aspect": 2.1, "asset_type": "figure"},
        "table_1": {"aspect": 3.0, "asset_type": "table"},
        "table_2": {"aspect": 3.2, "asset_type": "table"},
    }
    structured_sections = {
        "paper_sections": [
            {"section_name": "Introduction", "section_type": "foundation", "key_points": ["A", "B"], "contains_figures": ["figure_1"], "contains_tables": []},
            {"section_name": "Method", "section_type": "method", "key_points": ["A", "B"], "contains_figures": ["figure_2"], "contains_tables": []},
            {"section_name": "Results", "section_type": "evaluation", "key_points": ["A", "B"], "contains_figures": [], "contains_tables": ["table_1", "table_2"]},
        ]
    }
    classified_visuals = {
        "key_visual": "figure_2",
        "problem_illustration": ["figure_1"],
        "method_workflow": ["figure_2"],
        "main_results": ["table_1"],
        "comparative_results": ["table_2"],
        "supporting": [],
    }

    result = selector.select(state, structured_sections, classified_visuals, state["visual_assets"])

    assert result["selection_mode"] == "adaptive_auto"
    assert result["selected_template"] in {"three_column_postergen", "two_plus_one_mixed", "one_plus_two_mixed"}
    assert result["selected_template"] == "three_column_postergen"


def test_adaptive_auto_layout_uses_dense_column_gaps():
    cfg = load_config()
    state = create_state("/tmp/paper.pdf", layout_template="adaptive_auto", width=54, height=36)
    state["resolved_layout_template"] = "three_column_postergen"
    layout = LayoutAgent()._resolve_template_layout(state)
    lanes = layout["lanes"]
    gaps = [
        round(float(lanes[index + 1]["x"]) - (float(lanes[index]["x"]) + float(lanes[index]["w"])), 4)
        for index in range(len(lanes) - 1)
    ]

    assert layout["template_name"] == "three_column_postergen"
    assert max(gaps) <= cfg["adaptive_auto_dense_layout"]["max_column_gap_inches"]
    assert lanes[0]["x"] == pytest.approx(cfg["adaptive_auto_dense_layout"]["poster_margin"])


def test_template_selector_exposes_right_heavy_lane_preference_in_auto_mode():
    selector = TemplateSelector(load_config())
    state = create_state("/tmp/paper.pdf", layout_template="adaptive_auto")
    state["visual_assets"] = {
        "figure_1": {"aspect": 2.4, "asset_type": "figure"},
        "figure_2": {"aspect": 2.4, "asset_type": "figure"},
        "table_1": {"aspect": 1.6, "asset_type": "table"},
        "table_2": {"aspect": 1.7, "asset_type": "table"},
    }
    structured_sections = {
        "paper_sections": [
            {"section_name": "Introduction", "section_type": "foundation", "key_points": ["A"], "contains_figures": ["figure_1"], "contains_tables": []},
            {"section_name": "Method", "section_type": "method", "key_points": ["A"], "contains_figures": ["figure_2"], "contains_tables": []},
            {
                "section_name": "Results",
                "section_type": "evaluation",
                "key_points": ["A", "B", "C", "D", "E", "F", "G"],
                "contains_figures": [],
                "contains_tables": ["table_1", "table_2"],
            },
        ]
    }
    classified_visuals = {
        "key_visual": "figure_2",
        "problem_illustration": ["figure_1"],
        "method_workflow": ["figure_2"],
        "main_results": ["table_1"],
        "comparative_results": ["table_2"],
        "supporting": [],
    }

    result = selector.select(state, structured_sections, classified_visuals, state["visual_assets"])

    assert result["preferred_template"] == "two_plus_one_mixed"
    assert result["selected_template"] in {"three_column_postergen", "two_plus_one_mixed", "one_plus_two_mixed"}


def test_template_selector_respects_manual_template_request():
    selector = TemplateSelector(load_config())
    state = create_state("/tmp/paper.pdf", layout_template="one_plus_two_mixed")

    result = selector.select(
        state,
        structured_sections={"paper_sections": []},
        classified_visuals={},
        visual_assets={},
    )

    assert result["selection_mode"] == "manual"
    assert result["selected_template"] == "one_plus_two_mixed"


def test_template_selector_auto_uses_standard_template_whitelist():
    selector = TemplateSelector(load_config())
    state = create_state("/tmp/paper.pdf", layout_template="auto")
    state["paper_poster_keypoints"] = [
        {"id": index, "key_point": f"Keypoint {index}", "section": "Method"}
        for index in range(1, 11)
    ]

    result = selector.select(
        state,
        structured_sections={"paper_sections": []},
        classified_visuals={},
        visual_assets={},
    )

    assert result["selection_mode"] == "standard_auto"
    expected_templates = set(load_config()["standard_template_policy"]["auto_templates"])
    assert result["selected_template"] in expected_templates
    selected = next(candidate for candidate in result["candidates"] if candidate["template_name"] == result["selected_template"])
    assert selected["template_name"] in expected_templates


def test_template_selector_standard_auto_prefers_default_standard_landscape_template():
    selector = TemplateSelector(load_config())
    state = create_state("/tmp/paper.pdf", layout_template="auto", width=54, height=27)
    visual_assets = {
        "figure_1": {"asset_type": "figure", "aspect": 2.0},
        "figure_2": {"asset_type": "figure", "aspect": 1.5},
        "table_1": {"asset_type": "table", "aspect": 2.4},
    }
    structured_sections = {
        "paper_sections": [
            {"section_name": f"Section {index}", "section_type": "method", "key_points": ["A"]}
            for index in range(6)
        ]
    }

    result = selector.select(
        state,
        structured_sections=structured_sections,
        classified_visuals={"main_results": ["table_1"]},
        visual_assets=visual_assets,
    )

    assert result["selection_mode"] == "standard_auto"
    assert result["selected_template"] == "cluster_43_landscape"


def _block_refinement_state(tmp_path, utilization=0.45):
    state = create_state(
        str(tmp_path / "paper.pdf"),
        layout_template="cluster_3_portrait",
        enable_block_vlm_review=True,
    )
    state["output_dir"] = str(tmp_path / "output")
    lane = {"id": "slot_1", "x": 1.0, "y": 2.0, "w": 8.0, "h": 30.0}
    used_height = lane["h"] * utilization
    state["layout_template_metadata"] = {
        "template_name": "cluster_3_portrait",
        "layout_mode": "template_prior",
        "lanes": [lane],
    }
    state["story_board"] = {
        "spatial_content_plan": {
            "sections": [
                {
                    "section_id": "method",
                    "section_title": "Method",
                    "column_assignment": "slot_1",
                    "slot_id": "slot_1",
                    "vertical_priority": "top",
                    "text_content": ["Existing factual bullet about the method."],
                    "visual_assets": [{"visual_id": "figure_1"}],
                }
            ]
        }
    }
    state["styled_layout"] = [
        {
            "type": "section_container",
            "section_id": "method",
            "lane_id": "slot_1",
            "slot_id": "slot_1",
            "x": lane["x"],
            "y": lane["y"],
            "width": lane["w"],
            "height": used_height,
        },
        {
            "type": "text",
            "id": "method_text",
            "section_id": "method",
            "lane_id": "slot_1",
            "slot_id": "slot_1",
            "x": lane["x"] + 0.3,
            "y": lane["y"] + 1.0,
            "width": lane["w"] - 0.6,
            "height": max(used_height - 1.0, 0.5),
            "font_size": 44,
            "line_spacing": 1.0,
            "content": "Existing factual bullet about the method.",
        },
    ]
    return state


def test_block_occupancy_analyzer_formula_actions(tmp_path):
    analyzer = BlockOccupancyAnalyzer()

    def analyze_with_visible_extent(utilization):
        state = _block_refinement_state(tmp_path, utilization=utilization)
        lane = state["layout_template_metadata"]["lanes"][0]
        state["styled_layout"].append(
            {
                "type": "visual",
                "id": "occupancy_extent_marker",
                "section_id": "method",
                "lane_id": "slot_1",
                "slot_id": "slot_1",
                "x": lane["x"] + 0.3,
                "y": lane["y"],
                "width": lane["w"] - 0.6,
                "height": lane["h"] * utilization,
            }
        )
        return analyzer.analyze(state)["blocks"][0]

    low = analyze_with_visible_extent(0.45)
    moderate = analyze_with_visible_extent(0.92)
    near_target = analyze_with_visible_extent(0.986)
    crowded = analyze_with_visible_extent(0.997)

    assert low["action"] == "expand"
    assert low["target_extra_chars"] > moderate["target_extra_chars"]
    assert moderate["action"] == "expand"
    assert 0 < moderate["target_extra_chars"] < low["target_extra_chars"]
    assert near_target["action"] == "keep"
    assert near_target["target_extra_chars"] == 0
    assert crowded["action"] == "reduce"


def test_block_occupancy_analyzer_uses_real_child_content_not_container_background(tmp_path):
    state = _block_refinement_state(tmp_path, utilization=1.0)
    state["styled_layout"][0]["height"] = 30.0
    state["styled_layout"][1]["y"] = 3.0
    state["styled_layout"][1]["height"] = 8.0

    block = BlockOccupancyAnalyzer().analyze(state)["blocks"][0]

    assert block["container_bbox"]["h"] == pytest.approx(30.0)
    assert block["used_height"] < 9.0
    assert block["visible_content_height"] < 8.0
    assert block["bottom_whitespace"] > 21.0
    assert block["action"] == "expand"
    assert block["target_extra_chars"] > 0


def test_final_quality_gate_rejects_block_below_hard_minimum(tmp_path):
    state = _block_refinement_state(tmp_path, utilization=0.50)
    state["output_dir"] = str(tmp_path / "output")
    state["template_layout_mode"] = "template_prior"
    state["final_poster_accepted"] = True
    content_dir = Path(state["output_dir"]) / "content"
    content_dir.mkdir(parents=True, exist_ok=True)
    (content_dir / "micro_layout_report.json").write_text(
        json.dumps({"validation": {"issues": []}}),
        encoding="utf-8",
    )

    result = _run_final_quality_gate(state)

    assert result["final_poster_accepted"] is False
    assert result["final_quality_gate"]["accepted"] is False
    assert result["final_quality_gate"]["failures"][0]["category"] == "occupancy"
    assert (content_dir / "final_quality_gate.json").exists()


def test_final_quality_gate_rejects_mean_below_98_percent(tmp_path):
    state = _block_refinement_state(tmp_path, utilization=0.95)
    state["output_dir"] = str(tmp_path / "output")
    state["template_layout_mode"] = "template_prior"
    state["final_poster_accepted"] = True
    content_dir = Path(state["output_dir"]) / "content"
    content_dir.mkdir(parents=True, exist_ok=True)
    (content_dir / "micro_layout_report.json").write_text(
        json.dumps({"validation": {"issues": []}}),
        encoding="utf-8",
    )

    result = _run_final_quality_gate(state)

    assert result["final_poster_accepted"] is False
    assert result["final_quality_gate"]["accepted"] is False
    assert any(
        failure["category"] == "occupancy_mean"
        for failure in result["final_quality_gate"]["failures"]
    )


def test_final_quality_gate_rejects_excessive_bottom_whitespace(tmp_path):
    state = _block_refinement_state(tmp_path, utilization=0.96)
    state["output_dir"] = str(tmp_path / "output")
    state["template_layout_mode"] = "template_prior"
    state["final_poster_accepted"] = True
    content_dir = Path(state["output_dir"]) / "content"
    content_dir.mkdir(parents=True, exist_ok=True)
    (content_dir / "micro_layout_report.json").write_text(
        json.dumps({"validation": {"issues": []}}),
        encoding="utf-8",
    )

    result = _run_final_quality_gate(state)

    assert result["final_poster_accepted"] is False
    assert result["final_quality_gate"]["accepted"] is False
    assert any(
        failure["category"] == "bottom_whitespace"
        for failure in result["final_quality_gate"]["failures"]
    )


def test_final_quality_gate_allows_subline_bottom_whitespace_when_within_line_height(tmp_path, monkeypatch):
    class FakeAnalyzer:
        def analyze(self, _state):
            return {
                "summary": {"mean_utilization": 0.975},
                "blocks": [
                    {
                        "slot_id": "slot_1",
                        "section_id": "method",
                        "section_title": "Method",
                        "available_height": 10.0,
                        "utilization": 0.975,
                        "bottom_whitespace": 0.55,
                        "line_height": 0.7,
                        "action": "keep",
                        "visual_count": 0,
                    }
                ],
            }

    monkeypatch.setattr(
        "src.agents.block_occupancy_analyzer.BlockOccupancyAnalyzer",
        lambda: FakeAnalyzer(),
    )

    state = create_state(str(tmp_path / "paper.pdf"), layout_template="cluster_3_portrait")
    state["output_dir"] = str(tmp_path / "output")
    state["template_layout_mode"] = "template_prior"
    state["final_poster_accepted"] = True
    state["layout_template_metadata"] = {"lanes": []}
    state["styled_layout"] = []
    content_dir = Path(state["output_dir"]) / "content"
    content_dir.mkdir(parents=True, exist_ok=True)
    (content_dir / "micro_layout_report.json").write_text(
        json.dumps({"validation": {"issues": []}}),
        encoding="utf-8",
    )

    result = _run_final_quality_gate(state)

    assert result["final_quality_gate"]["accepted"] is True


def test_final_quality_gate_dedupes_degraded_states_and_removes_stale_repair_report(tmp_path, monkeypatch):
    class FakeAnalyzer:
        def analyze(self, _state):
            return {
                "summary": {"mean_utilization": 0.975},
                "blocks": [
                    {
                        "slot_id": "slot_1",
                        "section_id": "method",
                        "section_title": "Method",
                        "available_height": 10.0,
                        "utilization": 0.975,
                        "bottom_whitespace": 0.02,
                        "line_height": 0.7,
                        "action": "keep",
                        "visual_count": 0,
                    }
                ],
            }

    monkeypatch.setattr(
        "src.agents.block_occupancy_analyzer.BlockOccupancyAnalyzer",
        lambda: FakeAnalyzer(),
    )

    state = create_state(str(tmp_path / "paper.pdf"), layout_template="cluster_43_landscape", width=54, height=27)
    state["output_dir"] = str(tmp_path / "output")
    state["template_layout_mode"] = "template_prior"
    state["final_poster_accepted"] = True
    state["layout_template_metadata"] = {"lanes": []}
    state["styled_layout"] = []
    state["degraded_quality_states"] = [
        {"component": "vlm_layout_reviewer", "category": "vlm_layout_review", "fallback": "deterministic_acceptance", "reason": "request 1"},
        {"component": "vlm_layout_reviewer", "category": "vlm_layout_review", "fallback": "deterministic_acceptance", "reason": "request 2"},
    ]
    content_dir = Path(state["output_dir"]) / "content"
    content_dir.mkdir(parents=True, exist_ok=True)
    stale_report = content_dir / "final_quality_repair_report.json"
    stale_report.write_text(json.dumps({"stale": True}), encoding="utf-8")
    (content_dir / "micro_layout_report.json").write_text(
        json.dumps({"validation": {"issues": []}}),
        encoding="utf-8",
    )

    result = _run_final_quality_gate(state)

    assert result["final_quality_gate"]["accepted"] is True
    assert len(result["final_quality_gate"]["degraded_quality_states"]) == 1
    assert not stale_report.exists()


def test_final_quality_gate_rejects_degraded_required_vlm_review(tmp_path, monkeypatch):
    state = _block_refinement_state(tmp_path, utilization=0.98)
    state["output_dir"] = str(tmp_path / "output")
    state["template_layout_mode"] = "template_prior"
    state["final_poster_accepted"] = True
    state["enable_vlm_layout_review"] = True
    state["vlm_layout_review"] = {
        "source": "fallback",
        "review_available": False,
        "degraded": True,
        "accept": False,
        "warnings": ["response.failed"],
    }
    content_dir = Path(state["output_dir"]) / "content"
    content_dir.mkdir(parents=True, exist_ok=True)
    (content_dir / "micro_layout_report.json").write_text(
        json.dumps({"validation": {"issues": []}}),
        encoding="utf-8",
    )

    result = _run_final_quality_gate(state)

    assert result["final_quality_gate"]["accepted"] is False
    assert any(
        failure.get("category") == "quality_review_unavailable"
        and failure.get("component") == "vlm_layout_reviewer"
        for failure in result["final_quality_gate"]["failures"]
    )


def test_final_quality_gate_records_affiliation_logo_degraded_state(tmp_path, monkeypatch):
    class FakeAnalyzer:
        def analyze(self, _state):
            return {
                "summary": {"mean_utilization": 0.975},
                "blocks": [
                    {
                        "slot_id": "slot_1",
                        "section_id": "method",
                        "section_title": "Method",
                        "available_height": 10.0,
                        "utilization": 0.975,
                        "bottom_whitespace": 0.02,
                        "line_height": 0.7,
                        "action": "keep",
                        "visual_count": 0,
                    }
                ],
            }

    monkeypatch.setattr(
        "src.agents.block_occupancy_analyzer.BlockOccupancyAnalyzer",
        lambda: FakeAnalyzer(),
    )

    state = create_state(
        str(tmp_path / "paper.pdf"),
        layout_template="cluster_43_landscape",
        width=54,
        height=27,
        enable_affiliation_logos=True,
    )
    state["output_dir"] = str(tmp_path / "output")
    state["template_layout_mode"] = "template_prior"
    state["final_poster_accepted"] = True
    state["layout_template_metadata"] = {"lanes": []}
    state["styled_layout"] = []
    state["affiliations"] = ["Example Research University"]
    state["affiliation_logos"] = []
    content_dir = Path(state["output_dir"]) / "content"
    content_dir.mkdir(parents=True, exist_ok=True)
    (content_dir / "micro_layout_report.json").write_text(
        json.dumps({"validation": {"issues": []}}),
        encoding="utf-8",
    )

    result = _run_final_quality_gate(state)

    assert result["final_quality_gate"]["accepted"] is True
    degraded = result["final_quality_gate"]["degraded_quality_states"]
    assert degraded[-1]["component"] == "affiliation_logo_agent"
    assert degraded[-1]["category"] == "affiliation_logo_resolution"
    assert degraded[-1]["fallback"] == "no_affiliation_logo"


def test_final_quality_gate_rejects_micro_layout_lane_overflow(tmp_path, monkeypatch):
    class FakeAnalyzer:
        def analyze(self, _state):
            return {
                "summary": {"mean_utilization": 0.99},
                "blocks": [
                    {
                        "slot_id": "slot_1",
                        "section_id": "method",
                        "section_title": "Method",
                        "available_height": 10.0,
                        "utilization": 0.99,
                        "bottom_whitespace": 0.02,
                        "line_height": 0.7,
                        "action": "keep",
                        "visual_count": 0,
                    }
                ],
            }

    monkeypatch.setattr(
        "src.agents.block_occupancy_analyzer.BlockOccupancyAnalyzer",
        lambda: FakeAnalyzer(),
    )

    state = create_state(str(tmp_path / "paper.pdf"), layout_template="cluster_43_landscape", width=54, height=27)
    state["output_dir"] = str(tmp_path / "output")
    state["template_layout_mode"] = "template_prior"
    state["final_poster_accepted"] = True
    state["layout_template_metadata"] = {"lanes": []}
    state["styled_layout"] = []
    content_dir = Path(state["output_dir"]) / "content"
    content_dir.mkdir(parents=True, exist_ok=True)
    (content_dir / "micro_layout_report.json").write_text(
        json.dumps(
            {
                "validation": {"issues": []},
                "lanes": [
                    {
                        "lane_id": "slot_1",
                        "force_fit_used": True,
                        "final_overflow": 0.12,
                    }
                ],
            }
        ),
        encoding="utf-8",
    )

    result = _run_final_quality_gate(state)

    assert result["final_poster_accepted"] is False
    assert result["final_quality_gate"]["accepted"] is False
    assert any(
        failure["category"] == "micro_layout_lane_overflow"
        for failure in result["final_quality_gate"]["failures"]
    )


def test_final_quality_gate_rejects_oversized_body_font(tmp_path, monkeypatch):
    class FakeAnalyzer:
        def analyze(self, _state):
            return {
                "summary": {"mean_utilization": 0.99},
                "blocks": [
                    {
                        "slot_id": "slot_3",
                        "section_id": "sec_application",
                        "section_title": "Application",
                        "available_height": 10.0,
                        "utilization": 0.99,
                        "bottom_whitespace": 0.02,
                        "line_height": 0.7,
                        "action": "keep",
                        "visual_count": 1,
                    }
                ],
            }

    monkeypatch.setattr(
        "src.agents.block_occupancy_analyzer.BlockOccupancyAnalyzer",
        lambda: FakeAnalyzer(),
    )

    state = create_state(str(tmp_path / "paper.pdf"), layout_template="cluster_43_landscape", width=54, height=27)
    state["output_dir"] = str(tmp_path / "output")
    state["template_layout_mode"] = "template_prior"
    state["final_poster_accepted"] = True
    state["layout_template_metadata"] = {"lanes": []}
    state["styled_layout"] = [
        {
            "type": "text",
            "id": "sec_application_text",
            "section_id": "sec_application",
            "lane_id": "slot_3",
            "slot_id": "slot_3",
            "font_size": 54,
            "content": "Short application takeaway.",
        }
    ]
    content_dir = Path(state["output_dir"]) / "content"
    content_dir.mkdir(parents=True, exist_ok=True)
    (content_dir / "micro_layout_report.json").write_text(
        json.dumps({"validation": {"issues": []}}),
        encoding="utf-8",
    )

    result = _run_final_quality_gate(state)

    assert result["final_poster_accepted"] is False
    assert result["final_quality_gate"]["accepted"] is False
    assert any(
        failure["category"] == "body_font_scale"
        for failure in result["final_quality_gate"]["failures"]
    )


def test_final_gate_whitespace_repair_feeds_block_refiner_expand_input(tmp_path):
    state = _block_refinement_state(tmp_path, utilization=0.96)
    state["output_dir"] = str(tmp_path / "output")
    state["template_layout_mode"] = "template_prior"
    state["final_poster_accepted"] = True
    content_dir = Path(state["output_dir"]) / "content"
    content_dir.mkdir(parents=True, exist_ok=True)
    (content_dir / "micro_layout_report.json").write_text(
        json.dumps({"validation": {"issues": []}}),
        encoding="utf-8",
    )

    result = _run_final_quality_gate(state)
    repair_occupancy = _build_final_gate_refinement_occupancy(result)

    assert repair_occupancy["source"] == "final_quality_gate_repair"
    assert len(repair_occupancy["blocks"]) == 1
    block = repair_occupancy["blocks"][0]
    assert block["section_id"] == "method"
    assert block["action"] == "expand"
    assert block["final_gate_repair"] is True
    assert block["target_extra_chars"] >= load_config()["block_refinement"]["near_line_rewrite_extra_chars"]


def test_block_content_refiner_preserves_final_gate_repair_extra_budget():
    state = create_state("/tmp/paper.pdf", layout_template="cluster_43_landscape", width=54, height=27)
    state["template_fast_mode"] = True
    state["block_occupancy_report"] = {
        "settings": {
            "acceptable_min": 0.96,
            "hard_max": 0.995,
        },
        "blocks": [
            {
                "slot_id": "slot_3",
                "section_id": "sec_application",
                "action": "expand",
                "target_extra_chars": 80,
                "utilization": 0.8399,
                "available_height": 12.4,
                "visible_content_height": 10.415,
                "used_height": 10.415,
                "bottom_whitespace": 1.985,
                "line_height": 1.058,
                "chars_per_line": 46,
                "visual_count": 1,
                "final_gate_repair": True,
                "reason": "final quality gate requested full block text rewrite for bottom whitespace",
            }
        ],
    }
    state["block_vlm_review"] = {"blocks": []}

    actions = BlockContentRefiner()._decide_actions(state)

    assert len(actions) == 1
    assert actions[0]["target_extra_chars"] >= load_config()["block_refinement"]["near_line_rewrite_extra_chars"]


def test_final_quality_gate_allows_readable_single_line_title_too_small_vlm_status(tmp_path):
    state = _block_refinement_state(tmp_path, utilization=0.98)
    state["output_dir"] = str(tmp_path / "output")
    state["template_layout_mode"] = "template_prior"
    state["final_poster_accepted"] = True
    state["layout_template_metadata"]["orientation"] = "portrait"
    state["styled_layout"][1].update(
        {
            "y": 2.1,
            "height": 29.9,
            "content": "\n".join(
                f"Filled factual line {index} with method detail and evidence."
                for index in range(1, 30)
            ),
        }
    )
    state["header_plan"] = {"title": {"single_line": True, "font_size": 45}}
    state["vlm_layout_review"] = {"global_assessment": {"title_readability": "too_small"}}
    content_dir = Path(state["output_dir"]) / "content"
    content_dir.mkdir(parents=True, exist_ok=True)
    (content_dir / "micro_layout_report.json").write_text(
        json.dumps({"validation": {"issues": []}}),
        encoding="utf-8",
    )

    result = _run_final_quality_gate(state)

    assert result["final_quality_gate"]["accepted"] is True
    assert result["final_quality_gate"]["overrides"][0]["reason"] == "single_line_title_policy"


def test_final_quality_gate_rejects_visual_below_footprint_contract(tmp_path):
    state = _block_refinement_state(tmp_path, utilization=0.96)
    state["output_dir"] = str(tmp_path / "output")
    state["template_layout_mode"] = "template_prior"
    state["final_poster_accepted"] = True
    state["visual_assets"] = {"figure_1": {"asset_type": "figure", "aspect": 2.8}}
    state["styled_layout"].append(
        {
            "type": "visual",
            "id": "method_visual_figure_1",
            "visual_id": "figure_1",
            "section_id": "method",
            "lane_id": "slot_1",
            "slot_id": "method_visual_figure_1",
            "x": 2.0,
            "y": 4.0,
            "width": 6.0,
            "height": 2.1,
        }
    )
    content_dir = Path(state["output_dir"]) / "content"
    content_dir.mkdir(parents=True, exist_ok=True)
    (content_dir / "micro_layout_report.json").write_text(
        json.dumps({"validation": {"issues": []}}),
        encoding="utf-8",
    )

    result = _run_final_quality_gate(state)

    assert result["final_poster_accepted"] is False
    assert result["final_quality_gate"]["accepted"] is False
    assert any(
        failure["category"] == "visual_footprint"
        for failure in result["final_quality_gate"]["failures"]
    )


def test_final_quality_gate_rejects_template_section_geometry_overlap(tmp_path):
    state = create_state(str(tmp_path / "paper.pdf"), layout_template="cluster_96_landscape", width=54, height=27)
    state["output_dir"] = str(tmp_path / "output")
    state["template_layout_mode"] = "template_prior"
    state["final_poster_accepted"] = True
    state["layout_template_metadata"] = {
        "template_name": "cluster_test",
        "layout_mode": "template_prior",
        "lanes": [
            {"id": "slot_1", "x": 1.0, "y": 4.0, "w": 10.0, "h": 4.0},
            {"id": "slot_2", "x": 1.0, "y": 7.8, "w": 10.0, "h": 4.0},
        ],
    }
    state["story_board"] = {
        "spatial_content_plan": {
            "sections": [
                {"section_id": "first", "section_title": "First", "slot_id": "slot_1", "text_content": ["First block."]},
                {"section_id": "second", "section_title": "Second", "slot_id": "slot_2", "text_content": ["Second block."]},
            ]
        }
    }
    state["styled_layout"] = [
        {
            "type": "section_container",
            "section_id": "first",
            "lane_id": "slot_1",
            "slot_id": "slot_1",
            "x": 1.0,
            "y": 4.0,
            "width": 10.0,
            "height": 4.05,
        },
        {
            "type": "text",
            "id": "first_text",
            "section_id": "first",
            "lane_id": "slot_1",
            "x": 1.2,
            "y": 4.3,
            "width": 9.6,
            "height": 3.45,
            "content": "First block.",
            "font_size": 44,
        },
        {
            "type": "section_container",
            "section_id": "second",
            "lane_id": "slot_2",
            "slot_id": "slot_2",
            "x": 1.0,
            "y": 7.8,
            "width": 10.0,
            "height": 3.9,
        },
        {
            "type": "text",
            "id": "second_text",
            "section_id": "second",
            "lane_id": "slot_2",
            "x": 1.2,
            "y": 8.1,
            "width": 9.6,
            "height": 3.45,
            "content": "Second block.",
            "font_size": 44,
        },
    ]
    content_dir = Path(state["output_dir"]) / "content"
    content_dir.mkdir(parents=True, exist_ok=True)
    (content_dir / "micro_layout_report.json").write_text(
        json.dumps({"validation": {"issues": []}}),
        encoding="utf-8",
    )

    result = _run_final_quality_gate(state)

    assert any("overlap" in issue for issue in _section_geometry_issues(state))
    assert result["final_poster_accepted"] is False
    assert any(
        failure["category"] == "section_geometry"
        for failure in result["final_quality_gate"]["failures"]
    )


def test_block_content_refiner_expands_underfilled_block_without_changing_refs(tmp_path, monkeypatch):
    state = _block_refinement_state(tmp_path, utilization=0.45)
    state["block_occupancy_report"] = BlockOccupancyAnalyzer().analyze(state)
    state["block_vlm_review"] = {
        "blocks": [
            {
                "slot_id": "slot_1",
                "section_id": "method",
                "status": "underfilled",
                "severity": "medium",
                "description": "large whitespace remains",
            }
        ]
    }
    before_section = deepcopy(state["story_board"]["spatial_content_plan"]["sections"][0])

    def fake_expansion(self, state, actions, section_by_id):
        return {
            "method": {
                "rewritten_bullets": [
                    "The method rewrites the original factual point so the block reads as one continuous section.",
                    "The method uses paper-grounded evidence to add detail while preserving the original section assignment."
                ]
            }
        }

    monkeypatch.setattr(BlockContentRefiner, "_generate_expansion_patches", fake_expansion)
    patch = BlockContentRefiner().refine(state)
    after_section = state["story_board"]["spatial_content_plan"]["sections"][0]

    assert patch["applied"] is True
    assert len(after_section["text_content"]) == 2
    assert after_section["text_content"][0] != before_section["text_content"][0]
    assert after_section["section_id"] == before_section["section_id"]
    assert after_section["slot_id"] == before_section["slot_id"]
    assert after_section["visual_assets"] == before_section["visual_assets"]


def test_block_content_refiner_rewrites_near_line_bottom_gap_instead_of_skipping(tmp_path, monkeypatch):
    state = _block_refinement_state(tmp_path, utilization=0.956)
    state["block_occupancy_report"] = {
        "settings": {"acceptable_min": 0.96, "hard_max": 0.995},
        "blocks": [
            {
                "slot_id": "slot_1",
                "section_id": "method",
                "section_title": "Method",
                "available_height": 12.125,
                "visible_content_height": 11.5969,
                "used_height": 11.5969,
                "bottom_whitespace": 0.5281,
                "line_height": 0.6389,
                "chars_per_line": 74,
                "utilization": 0.9564,
                "action": "expand",
                "target_extra_chars": 60,
                "visual_count": 1,
                "reason": "utilization below target",
            }
        ],
    }
    state["block_vlm_review"] = {
        "blocks": [
            {
                "slot_id": "slot_1",
                "section_id": "method",
                "status": "ok",
                "severity": "low",
                "description": "content is readable but bottom whitespace remains",
            }
        ]
    }

    actions = BlockContentRefiner()._decide_actions(state)

    assert actions
    assert actions[0]["action"] == "expand"
    assert actions[0]["target_extra_chars"] > 0


def test_block_content_refiner_fallback_expands_vlm_underfilled_small_gap(tmp_path, monkeypatch):
    state = _block_refinement_state(tmp_path, utilization=0.943)
    state["template_fast_mode"] = True
    state["raw_text"] = (
        "The method updates predictions after each queried parcel label arrives online. "
        "The adaptive search policy balances exploration and exploitation under a limited outreach budget."
    )
    state["block_occupancy_report"] = {
        "settings": {"acceptable_min": 0.96, "hard_max": 0.995},
        "blocks": [
            {
                "slot_id": "slot_1",
                "section_id": "method",
                "section_title": "Method",
                "utilization": 0.943,
                "action": "expand",
                "target_extra_chars": 27,
                "visual_count": 0,
                "reason": "utilization below target",
            }
        ],
    }
    state["block_vlm_review"] = {
        "blocks": [
            {
                "slot_id": "slot_1",
                "section_id": "method",
                "status": "underfilled",
                "severity": "medium",
                "description": "visible bottom whitespace remains",
            }
        ]
    }

    class EmptyPatchResponse:
        content = json.dumps({"patches": []})
        input_tokens = 1
        output_tokens = 1

    class EmptyPatchAgent:
        def __init__(self, *args, **kwargs):
            pass

        def step(self, prompt):
            return EmptyPatchResponse()

    monkeypatch.setattr("src.agents.block_content_refiner.LangGraphAgent", EmptyPatchAgent)

    patch = BlockContentRefiner().refine(state)
    section = state["story_board"]["spatial_content_plan"]["sections"][0]

    assert patch["applied"] is True
    assert patch["actions_considered"][0]["target_extra_chars"] >= 70
    assert len(section["text_content"]) == 2
    assert "queried parcel label" in section["text_content"][1]


def test_block_content_refiner_expands_geometry_underfill_despite_medium_crowded_vlm(tmp_path, monkeypatch):
    state = _block_refinement_state(tmp_path, utilization=0.955)
    state["block_occupancy_report"] = {
        "settings": {"acceptable_min": 0.96, "hard_max": 0.995},
        "blocks": [
            {
                "slot_id": "slot_1",
                "section_id": "method",
                "section_title": "Method",
                "utilization": 0.955,
                "action": "expand",
                "target_extra_chars": 70,
                "visual_count": 1,
                "reason": "bottom whitespace above final threshold",
            }
        ],
    }
    state["block_vlm_review"] = {
        "blocks": [
            {
                "slot_id": "slot_1",
                "section_id": "method",
                "status": "crowded",
                "severity": "medium",
                "description": "dense but still has bottom whitespace",
            }
        ]
    }

    def fake_expansion(self, state, actions, section_by_id):
        return {"method": {"new_bullets": ["Adds a concise paper-grounded result detail to close the remaining whitespace."]}}

    monkeypatch.setattr(BlockContentRefiner, "_generate_expansion_patches", fake_expansion)

    patch = BlockContentRefiner().refine(state)

    assert patch["applied"] is True
    assert patch["actions_considered"][0]["action"] == "expand"


def test_block_content_refiner_caps_second_round_expand_by_remaining_geometry(tmp_path, monkeypatch):
    state = _block_refinement_state(tmp_path, utilization=0.905)
    state["block_refinement_count"] = 1
    state["block_occupancy_report"] = {
        "settings": {"acceptable_min": 0.96, "hard_max": 0.995},
        "blocks": [
            {
                "slot_id": "slot_1",
                "section_id": "method",
                "section_title": "Method",
                "available_height": 12.39,
                "visible_content_height": 11.2168,
                "used_height": 11.2168,
                "line_height": 0.6069,
                "chars_per_line": 42,
                "utilization": 0.905,
                "action": "expand",
                "target_extra_chars": 34,
                "visual_count": 0,
                "reason": "utilization below target",
            }
        ],
    }
    state["block_vlm_review"] = {
        "blocks": [
            {
                "slot_id": "slot_1",
                "section_id": "method",
                "status": "underfilled",
                "severity": "medium",
                "description": "visible bottom whitespace remains",
            }
        ]
    }

    def fake_expansion(self, state, actions, section_by_id):
        assert actions[0]["target_extra_chars"] == 42
        return {"method": {"new_bullets": ["Online labels update parcel risk during search."]}}

    monkeypatch.setattr(BlockContentRefiner, "_generate_expansion_patches", fake_expansion)

    patch = BlockContentRefiner().refine(state)

    assert patch["applied"] is True
    assert patch["actions_considered"][0]["target_extra_chars"] == 42


def test_block_content_refiner_allows_near_one_line_bottom_whitespace():
    refiner = BlockContentRefiner()
    safe_extra_chars = refiner._safe_extra_chars_for_block(
        {
            "available_height": 18.98,
            "visible_content_height": 18.1317,
            "used_height": 18.1317,
            "bottom_whitespace": 0.8483,
            "line_height": 0.8306,
            "chars_per_line": 86,
        }
    )

    assert safe_extra_chars == 86


def test_block_content_refiner_keeps_capacity_bullet_count_during_expansion():
    refiner = BlockContentRefiner()
    current = ["Primary evaluation fact.", "Secondary robustness fact."]
    rewritten = [
        "Primary evaluation fact with clearer protocol context.",
        "Secondary robustness fact with a stronger takeaway.",
        "Additional implementation detail that should be merged rather than added as a third paragraph.",
    ]

    result = refiner._apply_rewrite(
        current,
        rewritten,
        {"target_extra_chars": 90, "max_final_bullets": 2},
    )

    assert len(result) == 2
    assert sum(len(item) for item in result) <= sum(len(item) for item in current) + 104


def test_block_content_refiner_caps_expansion_with_whole_sentences():
    refiner = BlockContentRefiner()
    current = [
        "Temporal features capture neighboring-frame motion for efficient coarse restoration.",
        "Identity features combine the current frame with audio and landmark cues.",
    ]
    rewritten = [
        "Temporal features capture neighboring-frame motion for efficient coarse restoration across multiple frame offsets.",
        "Identity features combine the current frame with audio and landmark cues to recover facial detail.",
        "The reconstruction module fuses both feature streams into a high-quality restored frame.",
    ]

    result = refiner._apply_rewrite(
        current,
        rewritten,
        {
            "target_extra_chars": 80,
            "max_final_bullets": 3,
            "max_final_chars": 250,
        },
    )

    assert sum(len(item) for item in result) <= 250
    assert all(item in rewritten for item in result)
    assert all(item.endswith(".") for item in result)
    assert refiner._clean_bullets(["1, both identity and audio cues improve facial restoration."]) == [
        "Both identity and audio cues improve facial restoration."
    ]


def test_block_content_refiner_allows_short_callout_for_final_whitespace_repair():
    state = _block_refinement_state(Path("/tmp"), utilization=0.91)
    section = state["story_board"]["spatial_content_plan"]["sections"][0]
    section["text_content"] = ["A complete visual explanation that already occupies several lines."]
    section["target_bullets"] = 1
    section["capacity_budget"] = {"target_bullets": 1}
    state["template_fast_mode"] = True
    state["block_occupancy_report"] = {
        "settings": {"acceptable_min": 0.96, "hard_max": 0.995},
        "blocks": [{
            "slot_id": "slot_1",
            "section_id": "method",
            "utilization": 0.91,
            "visual_count": 1,
            "action": "expand",
            "target_extra_chars": 80,
            "bottom_whitespace": 1.0,
            "available_height": 10.0,
            "used_height": 9.0,
            "visible_content_height": 9.0,
            "line_height": 0.8,
            "chars_per_line": 50,
            "final_gate_repair": True,
        }],
    }
    state["block_vlm_review"] = {"blocks": []}

    actions = BlockContentRefiner()._decide_actions(state)

    assert actions[0]["action"] == "expand"
    assert actions[0]["max_final_bullets"] == 2


def test_curator_allows_one_caption_for_single_bullet_visual_block():
    curator = StoryBoardCurator()

    minimum = curator._minimum_text_items_for_section(
        {
            "text_content": ["Concise visual interpretation."],
            "capacity_budget": {
                "target_bullets": 1,
                "visual_policy": "figure_caption",
            },
        }
    )

    assert minimum == 1


def test_block_content_refiner_resets_header_review_before_relayout():
    state = create_state("/tmp/paper.pdf", layout_template="cluster_13_portrait", width=36, height=50.88)
    state["header_block_review"] = {"accepted": True}
    state["header_block_patch_applied"] = True

    BlockContentRefiner()._reset_downstream_state(state)

    assert state["header_block_review"] is None
    assert state["header_block_patch_applied"] is False


def test_block_content_refiner_reduces_crowded_block_without_changing_refs(tmp_path):
    state = _block_refinement_state(tmp_path, utilization=0.997)
    section = state["story_board"]["spatial_content_plan"]["sections"][0]
    section["text_content"] = [
        "A long factual bullet that is useful but can be shortened when the block is crowded by too much text.",
        "A second factual bullet describing a supporting detail from the paper.",
        "A third lower-priority factual bullet that can be removed first.",
        "A fourth low-priority factual bullet that should be removed in a crowded block.",
    ]
    state["styled_layout"][1]["content"] = "\n".join(section["text_content"])
    state["styled_layout"][1]["height"] = 29.0
    state["block_occupancy_report"] = {
        "settings": {"acceptable_min": 0.96, "hard_max": 0.995},
        "blocks": [
            {
                "slot_id": "slot_1",
                "section_id": "method",
                "section_title": "Method",
                "utilization": 0.997,
                "action": "reduce",
                "target_extra_chars": 0,
                "visual_count": 0,
                "reason": "utilization exceeds hard max",
            }
        ],
    }
    state["block_vlm_review"] = {
        "blocks": [
            {
                "slot_id": "slot_1",
                "section_id": "method",
                "status": "crowded",
                "severity": "medium",
                "description": "block looks too dense",
            }
        ]
    }
    before_section = deepcopy(section)

    patch = BlockContentRefiner().refine(state)
    after_section = state["story_board"]["spatial_content_plan"]["sections"][0]

    assert patch["applied"] is True
    assert len(after_section["text_content"]) < len(before_section["text_content"])
    assert after_section["section_id"] == before_section["section_id"]
    assert after_section["slot_id"] == before_section["slot_id"]
    assert after_section["visual_assets"] == before_section["visual_assets"]


def test_capacity_shortening_never_turns_a_single_sentence_into_a_fragment():
    planner = TemplateBlockPlanner()
    refiner = BlockContentRefiner()

    planner_sentence = "GAVN is built from three modules: inter-frame temporal modeling, intra-frame identity modeling, and a reconstruction module that fuses both cues."
    refiner_sentence = "The largest gains appear in lip-sync metrics and mouth-region restoration, showing that audio guidance improves reconstruction."
    assert planner._truncate_on_word_boundary(planner_sentence, 97) == planner_sentence
    assert refiner._truncate_on_word_boundary(refiner_sentence, 94) == refiner_sentence
    assert planner._truncate_on_word_boundary(
        "The first sentence is complete. The second sentence contains supporting implementation details.",
        38,
    ) == "The first sentence is complete."
    assert planner._truncate_on_word_boundary(
        "The first complete sentence is longer than the limit. A short second sentence follows.",
        24,
    ) == "The first complete sentence is longer than the limit."
    assert normalize_text_for_poster("Each visit consumes scarce outreach budget and may also.").endswith("budget.")
    assert normalize_text_for_poster("HAGS scales city-wide search by first selecting a.").endswith("search.")
    assert normalize_text_for_poster("The policy then selects the next property within the chosen region using local.").endswith("region.")
    repaired_called_tail = normalize_text_for_poster(
        "The paper proposes reinforcement learning for AGS, then addresses scaling challenges using a hierarchical approach called."
    )
    assert repaired_called_tail.endswith("challenges.")
    assert "called" not in repaired_called_tail
    assert normalize_text_for_poster("Updates labels instead of relying only on stale.").endswith("labels.")
    assert normalize_text_for_poster("The predictor is updated with binary cross-entropy, letting.").endswith("entropy.")
    assert normalize_text_for_poster("Teams work with limited staff, limited.").endswith("staff.")
    assert normalize_text_for_poster("judges differ in quality, consistency,.").endswith("consistency.")
    assert normalize_text_for_poster("the method improves generalization through.").endswith("generalization.")
    assert normalize_text_for_poster("Parcels are queried within a budget, while costs may.").endswith("budget.")
    assert normalize_text_for_poster("The search policy is trained with R.").endswith("trained.")
    assert normalize_text_for_poster("Data uses tabular features and imagery; targets are.").endswith("imagery.")
    assert normalize_text_for_poster("Formalizes eviction outreach as a sequential.").endswith("outreach.")
    assert normalize_text_for_poster(
        "The framework has two stages: (1) Primitive Patch Compression, which uses a 3D."
    ).endswith("Primitive Patch Compression.")
    assert normalize_text_for_poster("HAGS scales this process by splitting.").endswith("process.")
    assert normalize_text_for_poster("HAGS achieves the best target discovery across all budgets, typically.").endswith("budgets.")
    assert normalize_text_for_poster("choose a parcel inside that region..").endswith("region.")
    assert normalize_text_for_poster("Unknown current risk creates a sequential.").endswith("risk.")
    assert normalize_text_for_poster("HAGS shares the policy across regions, injecting.").endswith("regions.")
    assert normalize_text_for_poster("The policy identifies target parcels while ac.").endswith("parcels.")
    assert normalize_text_for_poster("Evaluation reports average number of targets fou.").endswith("targets.")
    assert normalize_text_for_poster("The formulation maximizes discoveries withi.").endswith("discoveries.")
    assert normalize_text_for_poster("The policy chooses parcels from predictions, past ou.").endswith("predictions.")
    assert normalize_text_for_poster("Average targets found under unifor.").endswith("found.")
    assert normalize_text_for_poster("The non-hierarchical AGS model; gains ar.").endswith("gains.")
    assert normalize_text_for_poster("Targets found under travel or query-cos.").endswith("travel.")
    assert normalize_text_for_poster("Symmetry acts as inductive bias for thousan.").endswith("bias.")
    assert normalize_text_for_poster("HAGS shares geospatial locality and reducing.").endswith("locality.")
    assert normalize_text_for_poster("HAGS wins in large-area search, with especially strong.").endswith("search.")
    assert normalize_text_for_poster("first choose a region, then choose a parcel within that re.").endswith("parcel.")
    assert normalize_text_for_poster("injecting geospatial locality and symmetry as inductive bias for large-area.").endswith("bias.")
    assert normalize_text_for_poster("first choose a region, then a parcel wi.").endswith("parcel.")
    assert normalize_text_for_poster("reducing action complexity for thousands of parc.").endswith("complexity.")
    assert normalize_text_for_poster("outperforming random, active search, greedy heuristics, and the non-hierarchi.").endswith("heuristics.")
    assert normalize_text_for_poster("must balance exploitation with exploration to update.").endswith("exploration.")
    assert normalize_text_for_poster("features from tabular records and overhead imagery; performance.").endswith("imagery.")
    assert normalize_text_for_poster("features plus overhead imagery, measuring average number.").endswith("imagery.")
    assert normalize_text_for_poster("features plus NAIP overhead imagery, and the average number.").endswith("imagery.")
    assert normalize_text_for_poster("strong gains under tighter budgets and lower eviction pr.").endswith("budgets.")
    assert normalize_text_for_poster("true risk is only revealed after.").endswith("revealed.")
    assert normalize_text_for_poster("outperforming random, greedy baselines, and the non.").endswith("baselines.")
    assert normalize_text_for_poster("injecting geospatial locality and symmetry as inductive bias for larg.").endswith("bias.")
    assert normalize_text_for_poster("Outreach operates under tight.").endswith("operates.")
    assert normalize_text_for_poster("The MDP enables princ.").endswith("MDP.")
    assert normalize_text_for_poster("Targets improve by approxima.").endswith("improve.")
    assert normalize_text_for_poster("HAGS improves discovery and substant.").endswith("discovery.")
    assert normalize_text_for_poster("The policy is optimized with a co.").endswith("optimized.")
    assert normalize_text_for_poster("The policy chooses next parcels as actions wit.").endswith("actions.")
    assert normalize_text_for_poster("The method learns a region-selection policy, and a share.").endswith("policy.")
    assert normalize_text_for_poster("Features pass through a shared mul.").endswith("shared.")
    assert normalize_text_for_poster("HAGS improves over the stronges.").endswith("improves.")
    assert normalize_text_for_poster("Outreach works despite limited.").endswith("works.")
    assert normalize_text_for_poster("Search uses either tabular or visual features with either unif.").endswith("features.")
    assert normalize_text_for_poster("The modules are jointly optimized vi.").endswith("optimized.")
    assert normalize_text_for_poster("Extend HAGS to scale learning and search to large.").endswith("search.")
    assert normalize_text_for_poster("HAGS finds more targets found, especially.").endswith("found.")
    assert normalize_text_for_poster("Multi-modal representations combining tabular and imagery features further.").endswith("features.")
    assert normalize_text_for_poster("Construct multimodal parcel features by combining rich tabular data (eviction histories, ownership, property.").endswith("data.")
    assert normalize_text_for_poster("HAGS shares locality as inductive bias to handle tens.").endswith("bias.")
    assert "Multimodal parcel" not in normalize_text_for_poster("Overall findings: - HAGS is strongest. - Multimodal parcel.")
    assert normalize_text_for_poster("The hierarchy matters with thousands of parcels, where flat action spaces become computationally and statistically.").endswith("parcels.")
    assert normalize_text_for_poster("The policy balances exploration with exploitation to visit.").endswith("exploitation.")
    assert normalize_text_for_poster("Parameter sharing acts as inductive bias for large urban.").endswith("bias.")
    assert normalize_text_for_poster("Performance is measured by average numbe.").endswith("measured.")
    assert normalize_text_for_poster("Outreach chooses visits under limited budgets, even though true near-term eviction risk.").endswith("budgets.")
    assert normalize_text_for_poster("Exploration improves predictions as new.").endswith("predictions.")
    assert normalize_text_for_poster("Canvassers choose properties to visit to find and support.").endswith("visit.")
    assert normalize_text_for_poster("Geospatial locality reduces complexity and improve.").endswith("complexity.")
    assert normalize_text_for_poster("Targets found within a searc.").endswith("found.")
    assert normalize_text_for_poster("Next, we describe in detail the data we use, as well as the baseline methods, before presenting our results.") == ""
    assert normalize_text_for_poster("HAGS finds 5-17% more at‑risk properties.") == "HAGS finds 5-17% more at-risk properties."
    assert normalize_text_for_poster("Our focus here is on large-area search; we defer most results involving small-area search to the Supplement.") == ""
    assert normalize_text_for_poster("We also provide results of random queries for calibration purposes.") == ""
    assert normalize_text_for_poster("Louis, USA Brown School at Washington University in St.") == ""
    assert normalize_text_for_poster(
        "After the episode, the search policy and initial predictor are updated jointly using reinforcement learning plus supervised loss: L_AGS = L_RL."
    ) == "After the episode, the search policy and initial predictor are updated jointly using reinforcement learning plus supervised loss."
    assert normalize_text_for_poster(
        "Overall, the gap identified by the paper is the lack of a framework that jointly handles."
    ) == "Overall, the gap identified by the paper is the lack of a framework."
    teaser_summary = GeneratedTeaserAgent()._truncate_on_word_boundary(
        "Eviction-prevention outreach must decide which rental properties to canvass under tight budgets, while true near-term eviction risk is only partially known and can change online.",
        150,
    )
    assert len(teaser_summary) > 80
    assert "cha." not in teaser_summary


def test_block_content_refiner_fallback_uses_a_minimum_complete_sentence_budget():
    refiner = BlockContentRefiner()

    weak = refiner._fallback_new_bullets(
        "Hierarchy is the main reason the method succeeds at urban scale.",
        [],
        {"target_extra_chars": 46},
    )
    assert weak == ["Hierarchy is the main reason the method succeeds at urban scale."]

    bullets = refiner._fallback_new_bullets(
        "Hierarchy is the main reason the method succeeds at urban scale. The paper introduces Active Geospatial Search to select parcels under limited outreach budgets while updating predictions after each observed label.",
        [],
        {"target_extra_chars": 46},
    )

    assert bullets == ["Hierarchy is the main reason the method succeeds at urban scale."]
    assert all(item.endswith(".") for item in bullets)
    assert sum(len(item) for item in bullets) <= 80


def test_template_block_planner_filters_reference_text_and_truncation_fragments():
    planner = TemplateBlockPlanner()
    state = create_state("/tmp/paper.pdf")
    state["raw_text"] = """
    Introduction
    The method prioritizes outreach cases with high eviction risk across neighborhoods.
    It updates regional scores after each field visit.

    References
    Desmond, M. 2016. Evicted: Poverty and Profit in the American City. Proceedings Press.
    Smith, A. 2020. Unrelated citation in Journal of Housing.
    """
    section = {
        "section_title": "Method",
        "section_id": "method",
        "content_role": "method",
    }

    source_sentences = planner._source_sentences_for_section(section, state)

    assert source_sentences
    assert all("Evicted" not in sentence for sentence in source_sentences)
    complete_sentence = "The planner tries to maximize discovery across neighborhoods using policy gradients."
    assert planner._truncate_on_word_boundary(complete_sentence, 37) == complete_sentence
    complete_sentence = "Risk targets for outreach are selected from high-priority geospatial regions."
    assert planner._truncate_on_word_boundary(complete_sentence, 38) == complete_sentence
    assert planner._is_clean_poster_bullet("risk targets fo.") is False
    assert planner._is_clean_poster_bullet("maximize dis.") is False
    assert planner._is_clean_poster_bullet("large-area se.") is False


def test_template_capacity_fallback_preserves_an_overlong_complete_sentence():
    planner = TemplateBlockPlanner()
    state = create_state("/tmp/paper.pdf")
    sentence = (
        "GAVN is built from three modules: inter-frame temporal modeling, intra-frame identity modeling, "
        "and a reconstruction module that fuses both cues."
    )

    bullets, warning = planner._fit_bullets_to_budget(
        [sentence],
        {"min_chars": 90, "target_chars": 90, "max_chars": 97, "target_bullets": 1},
        {"section_id": "method", "section_title": "GAVN Overview"},
        state,
        allow_expand=False,
    )

    assert bullets == [sentence]
    assert warning == "complete_sentence_exceeds_capacity"


def test_template_capacity_uses_actual_wide_visual_height_in_fast_mode():
    planner = TemplateBlockPlanner()
    state = create_state("/tmp/paper.pdf", width=54, height=27)
    state["template_fast_mode"] = True
    state["fast_block_contract"] = {
        "by_slot": {
            "slot_3": {
                "target_chars": 90,
                "min_chars": 90,
                "max_chars": 97,
                "target_bullets": 1,
                "visual_policy": "figure_caption",
            }
        }
    }
    state["visual_assets"] = {"figure_3": {"asset_type": "figure", "aspect": 3.64}}
    section = {
        "section_id": "pipeline",
        "region_id": "slot_3",
        "content_role": "method",
        "visual_assets": [{"visual_id": "figure_3"}],
    }
    region = {"region_id": "slot_3", "x": 36.88, "y": 5.1, "w": 16.12, "h": 12.4}

    budget = planner._capacity_budget_for_section(
        section,
        region,
        state,
        planner._capacity_settings(),
    )

    assert budget["reserved_visual_height"] < 5.0
    assert budget["target_chars"] > 250
    assert budget["target_bullets"] >= 3
    assert budget["source"] == "fast_template_actual_visual_contract"


def test_template_block_planner_does_not_expand_method_blocks_with_result_summaries():
    planner = TemplateBlockPlanner()
    state = create_state("/tmp/paper.pdf")
    state["raw_text"] = """
    HAGS is introduced to make AGS scalable to large urban search spaces.
    Overall empirical conclusion: HAGS is the strongest method across cost models, budgets, and target rates.
    In HAGS, the area is divided into regions.
    """
    section = {
        "section_title": "HAGS Overview",
        "section_id": "keypoint_group_02_hags_overview",
        "content_role": "method",
        "source_section": "Introduction",
        "source_sections": ["Introduction"],
        "source_keypoints": ["HAGS is introduced to make AGS scalable to large urban search spaces."],
        "text_content": ["HAGS is introduced to make AGS scalable to large urban search spaces."],
    }

    source_sentences = planner._source_sentences_for_section(section, state)

    assert any("area is divided into regions" in sentence for sentence in source_sentences)
    assert all("strongest method" not in sentence for sentence in source_sentences)
    expanded = planner._expand_bullets_from_source(
        ["HAGS is introduced to make AGS scalable to large urban search spaces."],
        {"min_chars": 160, "max_chars": 260, "target_bullets": 2},
        section,
        state,
    )
    assert all("strongest method" not in sentence for sentence in expanded)


def test_block_content_refiner_forces_min_budget_for_vlm_underfilled_keep(tmp_path):
    state = _block_refinement_state(tmp_path, utilization=0.87)
    state["block_occupancy_report"] = BlockOccupancyAnalyzer().analyze(state)
    block = state["block_occupancy_report"]["blocks"][0]
    block["action"] = "keep"
    block["target_extra_chars"] = 0
    state["block_vlm_review"] = {
        "blocks": [
            {
                "slot_id": "slot_1",
                "section_id": "method",
                "status": "underfilled",
                "severity": "medium",
                "description": "visible whitespace remains",
            }
        ]
    }

    actions = BlockContentRefiner()._decide_actions(state)

    assert actions
    assert actions[0]["action"] == "expand"
    assert actions[0]["target_extra_chars"] > 0


def test_block_content_refiner_compresses_caption_for_visual_too_small(tmp_path):
    state = _block_refinement_state(tmp_path, utilization=0.99)
    state["styled_layout"].append({
        "type": "visual",
        "id": "method_figure",
        "section_id": "method",
        "lane_id": "slot_1",
        "slot_id": "slot_1",
        "x": 1.3,
        "y": 3.0,
        "width": 6.0,
        "height": 3.0,
        "visual_id": "figure_1",
    })
    state["block_occupancy_report"] = {
        "settings": {"acceptable_min": 0.96, "hard_max": 0.995},
        "blocks": [
            {
                "slot_id": "slot_1",
                "section_id": "method",
                "section_title": "Method",
                "utilization": 0.99,
                "action": "keep",
                "target_extra_chars": 0,
                "visual_count": 1,
                "reason": "near target",
            }
        ],
    }
    state["block_vlm_review"] = {
        "blocks": [
            {
                "slot_id": "slot_1",
                "section_id": "method",
                "status": "visual_too_small",
                "severity": "high",
                "description": "figure labels are unreadable",
            }
        ]
    }

    actions = BlockContentRefiner()._decide_actions(state)

    assert actions
    assert actions[0]["action"] == "reduce"
    assert actions[0]["reason"] == "compress text to prioritize visual scale for unreadable figure/table labels"


def test_block_content_refiner_expands_underfilled_visual_too_small_block(tmp_path):
    state = _block_refinement_state(tmp_path, utilization=0.74)
    state["styled_layout"].append({
        "type": "visual",
        "id": "result_table",
        "section_id": "method",
        "lane_id": "slot_1",
        "slot_id": "slot_1",
        "x": 1.3,
        "y": 3.0,
        "width": 6.0,
        "height": 2.0,
        "visual_id": "table_1",
    })
    state["block_occupancy_report"] = BlockOccupancyAnalyzer().analyze(state)
    state["block_vlm_review"] = {
        "blocks": [
            {
                "slot_id": "slot_1",
                "section_id": "method",
                "status": "visual_too_small",
                "severity": "high",
                "description": "table labels are small, but the block has empty lower space",
            }
        ]
    }
    state["block_refinement_count"] = 1

    actions = BlockContentRefiner()._decide_actions(state)

    assert actions
    assert actions[0]["action"] == "expand"
    assert actions[0]["target_extra_chars"] >= 120


def test_block_content_refiner_keeps_medium_crowded_block_below_hard_max(tmp_path):
    state = _block_refinement_state(tmp_path, utilization=0.973)
    state["block_occupancy_report"] = BlockOccupancyAnalyzer().analyze(state)
    block = state["block_occupancy_report"]["blocks"][0]
    block["action"] = "keep"
    block["utilization"] = 0.973
    state["block_vlm_review"] = {
        "blocks": [
            {
                "slot_id": "slot_1",
                "section_id": "method",
                "status": "crowded",
                "severity": "medium",
                "description": "dense but not overflowing",
            }
        ]
    }

    actions = BlockContentRefiner()._decide_actions(state)

    assert actions == []


def test_block_content_refiner_fast_reduce_preserves_min_chars(tmp_path):
    refiner = BlockContentRefiner()
    bullets = [
        "PhishAgent combines online and offline knowledge retrieval with multimodal webpage evidence to address stale brand knowledge, local-brand ambiguity, and delayed phishing indicators.",
        "The agent keeps latency low by using a single-iteration workflow that retrieves candidate brand evidence before making the final phishing decision.",
        "This design targets realistic phishing pages where screenshots, logos, HTML, and domain cues must be interpreted together.",
    ]
    section = {
        "min_chars": 450,
        "max_chars": 550,
        "capacity_budget": {"min_chars": 450, "max_chars": 550},
    }
    action = {"vlm_status": "overflow", "utilization": 0.99}

    reduced = refiner._reduce_bullets_fast(bullets, action, section)

    assert sum(len(item) for item in reduced) >= 450
    assert len(reduced) == len(bullets)


def test_block_content_refiner_fast_skips_medium_crowded_reduce(tmp_path):
    state = create_state(str(tmp_path / "paper.pdf"), enable_block_vlm_review=True)
    state["template_fast_mode"] = True
    state["block_occupancy_report"] = {
        "settings": {"acceptable_min": 0.90, "hard_max": 0.98},
        "blocks": [
            {
                "slot_id": "slot_1",
                "section_id": "motivation",
                "utilization": 0.989,
                "action": "reduce",
                "target_extra_chars": 0,
            }
        ],
    }
    state["block_vlm_review"] = {
        "blocks": [
            {
                "slot_id": "slot_1",
                "section_id": "motivation",
                "status": "crowded",
                "severity": "medium",
            }
        ]
    }

    actions = BlockContentRefiner()._decide_actions(state)

    assert actions == []


def test_block_content_refiner_fast_allows_light_text_fill_repair(tmp_path):
    state = create_state(str(tmp_path / "paper.pdf"), enable_block_vlm_review=True)
    state["template_fast_mode"] = True
    state["block_occupancy_report"] = {
        "settings": {"acceptable_min": 0.90, "hard_max": 0.98},
        "blocks": [
            {
                "slot_id": "slot_6",
                "section_id": "results",
                "utilization": 0.89,
                "action": "expand",
                "target_extra_chars": 520,
                "visual_count": 0,
                "reason": "real content leaves bottom whitespace",
            }
        ],
    }
    state["block_vlm_review"] = {"blocks": [{"slot_id": "slot_6", "section_id": "results", "status": "underfilled", "severity": "medium"}]}

    actions = BlockContentRefiner()._decide_actions(state)

    assert actions
    assert actions[0]["action"] == "expand"
    assert actions[0]["target_extra_chars"] <= 420


def test_block_vlm_reviewer_falls_back_when_request_fails(tmp_path, monkeypatch):
    preview_path = tmp_path / "preview.png"
    Image.new("RGB", (900, 500), "white").save(preview_path)
    state = create_state(str(tmp_path / "paper.pdf"), enable_block_vlm_review=True)
    state["output_dir"] = str(tmp_path / "output")
    state["poster_width"] = 9.0
    state["poster_height"] = 5.0
    state["poster_preview_path"] = str(preview_path)
    state["block_occupancy_report"] = {
        "blocks": [
            {
                "slot_id": "slot_1",
                "section_id": "method",
                "section_title": "Method",
                "bbox": {"x": 1.0, "y": 1.0, "w": 4.0, "h": 2.0},
                "utilization": 0.99,
                "action": "reduce",
                "visual_count": 1,
                "reason": "above hard max",
            }
        ]
    }
    monkeypatch.setenv("VLM_BASE_URL", "https://example.invalid/v1")
    monkeypatch.setenv("VLM_API_KEY", "test-key")
    monkeypatch.setenv("VLM_MODEL", "gpt-5.4")
    reviewer = BlockVLMReviewer()
    monkeypatch.setattr(reviewer, "_post_vlm_request", lambda *args, **kwargs: (_ for _ in ()).throw(RuntimeError("bad gateway")))

    result = reviewer(state)

    assert result["errors"] == []
    review = result["block_vlm_review"]
    assert review["source"] == "fallback"
    assert review["blocks"][0]["status"] == "crowded"
    assert review["contact_sheet_path"].endswith("block_contact_sheet.png")
    assert "bad gateway" in review["warnings"][0]
    assert result["degraded_quality_states"][-1]["component"] == "block_vlm_reviewer"
    assert result["degraded_quality_states"][-1]["category"] == "block_vlm_review"
    assert result["degraded_quality_states"][-1]["fallback"] == "occupancy_only_block_review"


def test_vlm_layout_reviewer_rejects_unresolved_template_whitespace(tmp_path):
    state = create_state(str(tmp_path / "paper.pdf"), enable_vlm_layout_review=True)
    state["output_dir"] = str(tmp_path / "output")
    state["template_layout_mode"] = "template_prior"
    content_dir = Path(state["output_dir"]) / "content"
    content_dir.mkdir(parents=True)
    (content_dir / "micro_layout_report.json").write_text(
        json.dumps(
            {
                "lanes": [
                    {"lane_id": "slot_4", "final_utilization": 0.72},
                    {"lane_id": "slot_1", "final_utilization": 0.94},
                ]
            }
        )
    )
    state["styled_layout"] = [
        {
            "type": "section_container",
            "section_id": "main_results",
            "slot_id": "slot_4",
            "lane_id": "slot_4",
        }
    ]
    review = {
        "overall_score": 90,
        "accept": True,
        "issues": [
            {
                "severity": "medium",
                "category": "whitespace",
                "target": "main_results",
                "description": "large empty area remains",
            }
        ],
        "patch": [],
    }

    gated = VLMLayoutReviewer()._enforce_template_acceptance_gate(review, state)

    assert gated["accept"] is False
    assert any("unresolved_whitespace=True" in warning for warning in gated["warnings"])


def test_vlm_layout_reviewer_accepts_clean_draft_after_max_template_repair(tmp_path):
    state = create_state(str(tmp_path / "paper.pdf"), enable_vlm_layout_review=True)
    state["output_dir"] = str(tmp_path / "output")
    state["template_layout_mode"] = "template_prior"
    state["template_repair_count"] = 1
    content_dir = Path(state["output_dir"]) / "content"
    content_dir.mkdir(parents=True)
    (content_dir / "micro_layout_report.json").write_text(
        json.dumps({"validation": {"issues": []}, "lanes": []})
    )
    review = {
        "overall_score": 72,
        "accept": False,
        "issues": [{"severity": "high", "category": "overflow", "target": "sec_10"}],
        "warnings": [],
    }

    accepted = VLMLayoutReviewer()._accept_after_max_template_repair(review, state)

    assert accepted["accept"] is True
    assert any("accepted after max repair count" in warning for warning in accepted["warnings"])
