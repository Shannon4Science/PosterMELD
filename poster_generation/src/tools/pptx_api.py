import math
from typing import Optional, Tuple, Dict, Any, Union
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn


class PPTXDirector:
    """
    通用且强大的 Python-PPTX 封装类，专为构建 Poster 排版工作流而设计。
    提供了针对文本框、形状、线条与图像操作的标准自动化能力。
    """
    def __init__(self, template_path: Optional[str] = None):
        if template_path and Path(template_path).exists():
            self.prs = Presentation(template_path)
        else:
            self.prs = Presentation()
            
        # 确保至少存在一个空白幻灯片 (索引 6 通常为空白布局)
        layout = self.prs.slide_layouts[6] 
        self.slide = self.prs.slides.add_slide(layout)

    def set_slide_dimensions(self, width_inches: float, height_inches: float):
        """
        设置幻灯片（画板）的长宽，单位为英寸。
        """
        self.prs.slide_width = Inches(width_inches)
        self.prs.slide_height = Inches(height_inches)

    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        """
        内部辅助方法：将十六进制颜色字符串转换为 RGBColor 对象。
        """
        hex_color = hex_color.lstrip('#')
        if len(hex_color) == 3:
            hex_color = ''.join(c + c for c in hex_color)
        r, g, b = (int(hex_color[i:i+2], 16) for i in (0, 2, 4))
        return RGBColor(r, g, b)

    def add_text_box(
        self,
        x: float, y: float, w: float, h: float,
        text: str,
        font_family: str = "Arial",
        font_size: int = 14,
        font_color: str = "#000000",
        bold: bool = False,
        italic: bool = False,
        alignment: str = "left",
        line_spacing: float = 1.0,
        kerning_pt: float = 0.0,
    ):
        """
        在指定位置添加并高度定制一个文本框。
        kerning_pt: 自定义字间距（字符间距）。
        所有的尺寸参数均以英寸(Inches)为单位。
        """
        textbox = self.slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = textbox.text_frame
        tf.word_wrap = True
        
        # 清除默认格式
        tf.clear()
        
        p = tf.paragraphs[0]
        p.text = text
        
        if alignment == "center":
            p.alignment = PP_ALIGN.CENTER
        elif alignment == "right":
            p.alignment = PP_ALIGN.RIGHT
        elif alignment == "justify":
            p.alignment = PP_ALIGN.JUSTIFY
        else:
            p.alignment = PP_ALIGN.LEFT
            
        p.line_spacing = line_spacing

        # 将样式精确应用到文本运行 (run) 级别
        for run in p.runs:
            run.font.name = font_family
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = self._hex_to_rgb(font_color)
            
            # 使用底层 XML 调整字间距 (Character Spacing)
            if kerning_pt != 0.0:
                # `spc` 标签的单位是百分之一磅 (1/100 point)
                spacing_val = int(kerning_pt * 100)
                # 确保 rPr 节点存在
                rPr = run._r.get_or_add_rPr()
                rPr.set('spc', str(spacing_val))
                
        return textbox

    def add_shape(
        self,
        shape_type: int, # 取值参考 MSO_SHAPE 枚举
        x: float, y: float, w: float, h: float,
        fill_color: Optional[str] = None,
        fill_opacity: Optional[float] = None,
        border_color: Optional[str] = None,
        border_opacity: Optional[float] = None,
        border_width: float = 1.0,
        border_style: str = "solid",
        shadow: Optional[Dict[str, Any]] = None,
    ):
        """
        添加形状（如矩形、圆形等）。支持设置填充色和边框色。
        """
        shape = self.slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
        
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = self._hex_to_rgb(fill_color)
            self._apply_color_alpha(shape.fill.fore_color._color, fill_opacity)
        else:
            shape.fill.background()
            
        if border_color:
            shape.line.color.rgb = self._hex_to_rgb(border_color)
            self._apply_color_alpha(shape.line.color._color, border_opacity)
            shape.line.width = Pt(border_width)
            if str(border_style).lower() in {"dash", "dashed"}:
                line = shape.line
                ln = line._ln
                prstDash = parse_xml('<a:prstDash xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" val="dash"/>')
                ln.append(prstDash)
        else:
            shape.line.fill.background()

        if shadow and shadow.get("enabled", True):
            self._apply_outer_shadow(shape, shadow)
            
        return shape

    def _apply_color_alpha(self, color_obj: Any, opacity: Optional[float]) -> None:
        if opacity is None:
            return
        try:
            opacity_value = max(0.0, min(float(opacity), 1.0))
        except (TypeError, ValueError):
            return
        alpha = int(round(opacity_value * 100000))
        srgb = color_obj._xClr
        for child in list(srgb):
            if child.tag == qn("a:alpha"):
                srgb.remove(child)
        srgb.append(parse_xml(f'<a:alpha xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" val="{alpha}"/>'))

    def _apply_outer_shadow(self, shape, shadow: Dict[str, Any]) -> None:
        """Apply a DrawingML outer shadow to a shape."""
        color = str(shadow.get("color", "#000000")).lstrip("#")[:6].upper()
        if len(color) != 6:
            color = "000000"
        alpha_raw = float(shadow.get("alpha", 0.16))
        alpha = int(alpha_raw * 100000) if alpha_raw <= 1 else int(alpha_raw)
        alpha = max(0, min(alpha, 100000))
        blur_pt = max(float(shadow.get("blur_pt", shadow.get("blur", 4.0))), 0.0)
        distance_pt = max(float(shadow.get("distance_pt", shadow.get("distance", 2.0))), 0.0)
        angle = float(shadow.get("angle", 45.0))
        blur_emu = int(blur_pt * 12700)
        distance_emu = int(distance_pt * 12700)
        direction = int(angle * 60000)

        sp_pr = shape._element.spPr
        for child in list(sp_pr):
            if child.tag == qn("a:effectLst"):
                sp_pr.remove(child)
        effect_xml = (
            f'<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            f'<a:outerShdw blurRad="{blur_emu}" dist="{distance_emu}" dir="{direction}" '
            f'algn="ctr" rotWithShape="0">'
            f'<a:srgbClr val="{color}"><a:alpha val="{alpha}"/></a:srgbClr>'
            f'</a:outerShdw></a:effectLst>'
        )
        sp_pr.append(parse_xml(effect_xml))

    def add_connector(
        self,
        start_x: float, start_y: float,
        end_x: float, end_y: float,
        color: str = "#000000",
        width_pt: float = 1.0,
        is_dashed: bool = False
    ):
        """
        添加连接线（直线）。由于 python-pptx 没有直接的 'add_line' 函数，
        我们使用 add_connector 来实现。
        """
        connector = self.slide.shapes.add_connector(
            1, # MSO_CONNECTOR.STRAIGHT == 1
            Inches(start_x), Inches(start_y),
            Inches(end_x), Inches(end_y)
        )
        connector.line.color.rgb = self._hex_to_rgb(color)
        connector.line.width = Pt(width_pt)
        if is_dashed:
            # 通过底层 OXML 设置虚线 (prstDash)
            line = connector.line
            ln = line._ln
            prstDash = parse_xml('<a:prstDash xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" val="dash"/>')
            ln.append(prstDash)
            
        return connector

    def add_image(
        self,
        image_path: str,
        x: float, y: float,
        w: float, h: float,
        keep_aspect_ratio: bool = True
    ):
        """
        渲染并插入图片。如果 keep_aspect_ratio 为真，我们会自动计算图片在目标宽高 (w, h) 内的自适应布局。
        """
        if not keep_aspect_ratio:
            return self.slide.shapes.add_picture(image_path, Inches(x), Inches(y), Inches(w), Inches(h))
            
        from PIL import Image
        with Image.open(image_path) as img:
            orig_w, orig_h = img.size
            aspect = orig_w / orig_h
        
        target_aspect = w / h
        if aspect > target_aspect:
            # 原图比目标框更宽，此时宽度达到上限
            final_w = w
            final_h = w / aspect
        else:
            # 原图比目标框更高，此时高度达到上限
            final_h = h
            final_w = h * aspect
            
        # 居中对齐到提供的框内
        centered_x = x + (w - final_w) / 2
        centered_y = y + (h - final_h) / 2
        
        return self.slide.shapes.add_picture(
            image_path, Inches(centered_x), Inches(centered_y),
            Inches(final_w), Inches(final_h)
        )

    def save(self, output_path: str):
        """保存幻灯片到指定路径"""
        self.prs.save(output_path)
