"""
powerpoint rendering using python-pptx
"""

import re
import qrcode
import shutil
import subprocess
import hashlib
from pathlib import Path
from typing import Dict, Any, Optional
import json

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from src.state.poster_state import PosterState
from utils.src.logging_utils import log_agent_info, log_agent_success, log_agent_error
from src.config.poster_config import load_config
from src.utils.text_cleanup import repair_possessive_title_apostrophe
from src.utils.style_options import resolve_poster_visual_style, resolve_typography_config


class Renderer:
    """powerpoint rendering with styling support"""
    
    def __init__(self):
        self.name = "renderer"
        self.styling_interfaces = None
        
        # load configuration
        self.config = load_config()
        self.layout_constants = self.config["layout_constants"]
        self.powerpoint_config = self.config["powerpoint"]
        self.indentation_config = self.config["indentation"]
        self.typography_config = self.config["typography"]
        self.visual_style_config = self.config.get("poster_visual_style", {})

    def __call__(self, state: PosterState) -> PosterState:
        log_agent_info(self.name, "Starting Rendering Process")
        
        try:
            self.visual_style_config = resolve_poster_visual_style(state, self.config)
            self.typography_config = resolve_typography_config(state, self.config)
            self.styling_interfaces = self._load_styling_interfaces(state)
            render_stage = state.get("render_stage", "final")
            suffix = "_draft" if render_stage == "draft" else ""
            output_path = Path(state["output_dir"]) / f"{state['poster_name']}{suffix}.pptx"
            self._render_presentation(state, output_path)
            state["pptx_output_path"] = str(output_path)
            state["current_agent"] = self.name
            
            # convert to png if possible
            png_path = self._convert_to_png(output_path)
            state["poster_preview_path"] = png_path
            state["final_poster_accepted"] = render_stage == "final" and state.get("draft_status") == "accepted"
            
            log_agent_success(self.name, f"rendered poster: {output_path}")
            if png_path:
                log_agent_success(self.name, f"generated preview: {png_path}")
                
        except Exception as e:
            log_agent_error(self.name, f"rendering failed: {e}")
            state["errors"].append(f"{self.name}: {e}")
            
        return state

    def _load_styling_interfaces(self, state: PosterState) -> Dict[str, Any]:
        """load styling interfaces from font agent output file"""
        if state.get("styling_interfaces"):
            interfaces = dict(state["styling_interfaces"])
            interfaces["line_spacing"] = 1.0
            return interfaces

        styling_path = Path(state["output_dir"]) / "content" / "styling_interfaces.json"
        if styling_path.exists():
            with open(styling_path, 'r', encoding='utf-8') as f:
                interfaces = json.load(f)
            interfaces["line_spacing"] = 1.0
            return interfaces
        else:
            # fallback to defaults with 1.0 line spacing
            return {
                "bullet_point_marker": "•",
                "bold_start_tag": "**",
                "bold_end_tag": "**",
                "italic_start_tag": "*",
                "italic_end_tag": "*",
                "color_start_tag": "<color:",
                "color_end_tag": "</color>",
                "line_spacing": 1.0,
                "paragraph_spacing": 0.1
            }

    def _render_presentation(self, state: PosterState, output_path: Path):
        """render complete presentation"""
        from src.tools.pptx_api import PPTXDirector
        self.director = PPTXDirector()
        self.director.set_slide_dimensions(state["poster_width"], state["poster_height"])
        slide = self.director.slide

        # TODO: generate QR code if needed
        qr_code_path = None
        if state.get("url"):
            qr_code_path = self._generate_qr_code(state["url"], state["output_dir"])
        self._render_output_dir = Path(state["output_dir"])

        self._render_background_image(slide, state)

        # use styled_layout if available, fallback to design_layout
        layout_data = state.get("styled_layout", state.get("design_layout", []))
        if not layout_data:
            raise ValueError("no styled_layout or design_layout found")
        
        # sort elements by priority to ensure proper rendering order
        sorted_elements = sorted(layout_data, key=lambda x: x.get("priority", 0.5))
        
        for element in sorted_elements:
            self._render_element(slide, element, state, qr_code_path)
        
        self.director.save(str(output_path))

    def _render_background_image(self, slide, state: PosterState):
        background_path = state.get("background_image_path")
        if not background_path or not Path(background_path).exists():
            return
        self.director.add_image(
            background_path,
            0,
            0,
            state["poster_width"],
            state["poster_height"],
            keep_aspect_ratio=False,
        )

    def _render_element(self, slide, element: Dict, state: PosterState, qr_code_path: Optional[str]):
        """render individual element based on type"""
        element_type = element.get("type")
        
        # handle QR code elements
        if element_type == "qr_code" and qr_code_path:
            self._render_qr_code(slide, element, qr_code_path)
            return
        
        # get appropriate renderer
        renderer_map = {
            "template_background": self._render_template_shape,
            "template_header_background": self._render_template_shape,
            "template_footer_background": self._render_template_shape,
            "title": self._render_title,
            "section_title": self._render_section_title,
            "title_accent_block": self._render_title_accent_block,
            "title_accent_line": self._render_title_accent_line,
            "conf_logo": self._render_conf_logo,
            "aff_logo": self._render_aff_logo,
            "institution_logo": self._render_institution_logo,
            "logo_divider": self._render_logo_divider,
            "section_container": self._render_section_container,
            "text": self._render_text,
            "visual": self._render_visual,
            "mixed": self._render_mixed,
        }
        
        renderer = renderer_map.get(element_type)
        if renderer:
            renderer(slide, element, state)
        else:
            log_agent_error(self.name, f"unknown element type: {element_type}")

    def _render_template_shape(self, slide, element: Dict, state: PosterState):
        fill_color = element.get("fill_color", "#FFFFFF")
        border_color = element.get("border_color")
        border_width = element.get("border_width", 0.0)
        self.director.add_shape(
            MSO_SHAPE.RECTANGLE,
            element["x"],
            element["y"],
            element["width"],
            element["height"],
            fill_color=fill_color,
            border_color=border_color,
            border_width=border_width,
            shadow=element.get("shadow"),
        )

    def _render_title(self, slide, element: Dict, state: PosterState):
        """render poster title with authors"""
        x, y, w, h = (Inches(element[k]) for k in ["x", "y", "width", "height"])
        
        log_agent_info(self.name, f"rendering title at ({x.inches:.1f}, {y.inches:.1f})")

        content = element.get("content", "Title\nAuthors")
        lines = content.split("\n")
        
        # separate title and authors
        title_lines = lines[:-1] if len(lines) > 1 else lines
        authors_text = lines[-1] if len(lines) > 1 else ""
        title_text = "\n".join(line.strip() for line in title_lines if line.strip())
        title_font_size = self.styling_interfaces.get("font_sizes", {}).get("title", 100)
        author_font_size = self.styling_interfaces.get("font_sizes", {}).get("authors", 72)
        title_font_size = element.get("font_size", title_font_size)
        author_font_size = element.get("author_font_size", author_font_size)
        style_enabled = self.visual_style_config.get("enabled", False)
        title_style = self.visual_style_config.get("main_title", {}) if style_enabled else {}
        title_style_override = element.get("main_title_style_override") or {}
        if isinstance(title_style_override, dict):
            title_style = {**title_style, **title_style_override}
        title_font_family = title_style.get("font_family") or element.get("font_family", "Georgia")
        author_font_family = title_style.get("author_font_family") or element.get("author_font_family", "Arial")
        title_color = self._parse_color(title_style.get("font_color") or element.get("font_color", "#07164A"))
        author_color = self._parse_color(title_style.get("author_font_color", "#333333"))
        title_shadow_cfg = title_style.get("shadow") or {}
        title_shadow = title_shadow_cfg if title_shadow_cfg.get("enabled", False) else None
        alignment = element.get("alignment", "left")
        title_single_line = bool(element.get("title_single_line", True))
        title_wrap_policy = str(element.get("title_wrap_policy") or "single_line")
        title_word_wrap = not (title_single_line or title_wrap_policy == "two_line")
        subtitle_single_line = bool(element.get("subtitle_single_line", True))
        if title_single_line:
            title_text = self._single_line_text(title_text)
            title_font_size = self._fit_single_line_font_size(
                title_text,
                w.inches,
                float(title_font_size),
                min_key="portrait_title_single_line_min_font_size" if h.inches > w.inches else "title_single_line_min_font_size",
            )

        if element.get("title_text") or element.get("subtitle_text"):
            title_text = str(element.get("title_text") or title_text).strip()
            subtitle_text = str(element.get("subtitle_text") or "").strip()
            authors_text = str(element.get("authors_text") or authors_text).strip()
            subtitle_font_size = float(element.get("subtitle_font_size", max(float(title_font_size) * 0.58, 24)))
            subtitle_color = self._parse_color(element.get("subtitle_font_color", title_style.get("subtitle_font_color", "#374151")))
            if title_single_line:
                title_text = self._single_line_text(title_text)
                title_font_size = self._fit_single_line_font_size(
                    title_text,
                    w.inches,
                    float(title_font_size),
                    min_key="portrait_title_single_line_min_font_size" if h.inches > w.inches else "title_single_line_min_font_size",
                )
            if subtitle_text and subtitle_single_line:
                subtitle_text = self._single_line_text(subtitle_text)
                subtitle_font_size = self._fit_single_line_font_size(
                    subtitle_text,
                    w.inches,
                    subtitle_font_size,
                    min_key="subtitle_single_line_min_font_size",
                )

            author_box_height = float(element.get("author_box_height") or max((float(author_font_size) / 72) * 1.15, 0.55))
            author_box_height = min(author_box_height, max(h.inches * 0.32, 0.55))
            subtitle_box_height = float(element.get("subtitle_box_height") or (max((subtitle_font_size / 72) * 1.12, 0.36) if subtitle_text else 0.0))
            subtitle_gap_inches = float(element.get("title_to_subtitle_gap_inches", 0.08 if subtitle_text else 0.0))
            author_gap_inches = float(
                element.get(
                    "author_top_gap_inches",
                    self.typography_config.get("title_author_gap_points", 16) / 72,
                )
            )
            author_x = float(element.get("author_x")) if element.get("author_x") is not None else x.inches
            author_width = float(element.get("author_width")) if element.get("author_width") is not None else w.inches
            author_word_wrap = bool(element.get("author_word_wrap", False))
            title_box_height = float(
                element.get("title_box_height")
                or max(h.inches - subtitle_box_height - subtitle_gap_inches - author_gap_inches - author_box_height, h.inches * 0.45)
            )
            if title_box_height + subtitle_gap_inches + subtitle_box_height + author_gap_inches + author_box_height > h.inches:
                title_box_height = max(h.inches - subtitle_gap_inches - subtitle_box_height - author_gap_inches - author_box_height, h.inches * 0.38)

            cursor_y = y.inches
            self._add_title_textbox(
                slide,
                title_text,
                element,
                x.inches,
                cursor_y,
                w.inches,
                title_box_height,
                font_size=title_font_size,
                font_family=title_font_family,
                bold=True,
                color=title_color,
                line_spacing=self.typography_config["line_spacing"],
                alignment=alignment,
                shadow=title_shadow,
                word_wrap=title_word_wrap,
            )
            cursor_y += title_box_height
            if subtitle_text:
                cursor_y += subtitle_gap_inches
                self._add_title_textbox(
                    slide,
                    subtitle_text,
                    element,
                    x.inches,
                    cursor_y,
                    w.inches,
                    subtitle_box_height,
                    font_size=subtitle_font_size,
                    font_family=element.get("subtitle_font_family", author_font_family),
                    bold=False,
                    color=subtitle_color,
                    line_spacing=self.typography_config["line_spacing"],
                    alignment=alignment,
                    word_wrap=not subtitle_single_line,
                )
                cursor_y += subtitle_box_height

            if authors_text:
                cursor_y += author_gap_inches
                self._add_title_textbox(
                    slide,
                    authors_text,
                    element,
                    author_x,
                    cursor_y,
                    author_width,
                    min(author_box_height, max(y.inches + h.inches - cursor_y, 0.35)),
                    font_size=author_font_size,
                    font_family=author_font_family,
                    bold=False,
                    color=author_color,
                    line_spacing=self.typography_config["line_spacing"] + 0.1,
                    alignment=alignment,
                    word_wrap=author_word_wrap,
                )
            return

        if authors_text:
            author_gap_inches = float(
                element.get(
                    "author_top_gap_inches",
                    self.typography_config.get("title_author_gap_points", 16) / 72,
                )
            )
            author_box_height = max((float(author_font_size) / 72) * 1.15, 0.55)
            author_box_height = min(author_box_height, max(h.inches * 0.32, 0.55))
            author_x = float(element.get("author_x")) if element.get("author_x") is not None else x.inches
            author_width = float(element.get("author_width")) if element.get("author_width") is not None else w.inches
            author_word_wrap = bool(element.get("author_word_wrap", False))
            title_box_height = max(h.inches - author_box_height - author_gap_inches, h.inches * 0.55)
            if title_box_height + author_gap_inches + author_box_height > h.inches:
                title_box_height = max(h.inches - author_gap_inches - author_box_height, h.inches * 0.45)
            author_y = y.inches + title_box_height + author_gap_inches

            self._add_title_textbox(
                slide,
                title_text,
                element,
                x.inches,
                y.inches,
                w.inches,
                title_box_height,
                font_size=title_font_size,
                font_family=title_font_family,
                bold=True,
                color=title_color,
                line_spacing=self.typography_config["line_spacing"],
                alignment=alignment,
                shadow=title_shadow,
                word_wrap=title_word_wrap,
            )
            self._add_title_textbox(
                slide,
                authors_text.strip(),
                element,
                author_x,
                author_y,
                author_width,
                min(author_box_height, max(y.inches + h.inches - author_y, 0.4)),
                font_size=author_font_size,
                font_family=author_font_family,
                bold=False,
                color=author_color,
                line_spacing=self.typography_config["line_spacing"] + 0.1,
                alignment=alignment,
                word_wrap=author_word_wrap,
            )
            return

        self._add_title_textbox(
            slide,
            "\n".join(title_line.strip() for title_line in title_lines if title_line.strip()),
            element,
            x.inches,
            y.inches,
            w.inches,
            h.inches,
            font_size=title_font_size,
            font_family=title_font_family,
            bold=True,
            color=title_color,
            line_spacing=self.typography_config["line_spacing"],
            alignment=alignment,
            shadow=title_shadow,
            word_wrap=title_word_wrap,
        )

    def _add_title_textbox(
        self,
        slide,
        text: str,
        element: Dict,
        x: float,
        y: float,
        width: float,
        height: float,
        font_size: float,
        font_family: str,
        bold: bool,
        color: RGBColor,
        line_spacing: float,
        alignment: str = "left",
        shadow: Optional[Dict[str, Any]] = None,
        word_wrap: bool = True,
    ):
        if shadow and shadow.get("enabled", True):
            self._add_title_textbox(
                slide,
                text,
                element,
                x + float(shadow.get("offset_x_inches", 0.04)),
                y + float(shadow.get("offset_y_inches", 0.04)),
                width,
                height,
                font_size=font_size,
                font_family=font_family,
                bold=bold,
                color=self._parse_color(shadow.get("color", "#8E96B2")),
                line_spacing=line_spacing,
                alignment=alignment,
                word_wrap=word_wrap,
            )
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
        tf = tb.text_frame
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.word_wrap = bool(word_wrap)
        tf.margin_left = Inches(0)
        tf.margin_right = Inches(0)
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = font_family
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        if alignment.lower() == "center":
            p.alignment = PP_ALIGN.CENTER
        elif alignment.lower() == "right":
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        return tb

    def _single_line_text(self, text: str) -> str:
        return re.sub(r"\s+", " ", str(text or "")).strip()

    def _fit_single_line_font_size(
        self,
        text: str,
        width_inches: float,
        desired_size: float,
        *,
        min_key: str,
    ) -> float:
        clean_text = self._single_line_text(text)
        if not clean_text or width_inches <= 0:
            return desired_size
        header_config = self.config.get("header_planner", {})
        if min_key == "subtitle_single_line_min_font_size":
            avg_char_width = float(
                header_config.get(
                    "portrait_subtitle_fit_avg_char_width_em",
                    header_config.get("subtitle_fit_avg_char_width_em", header_config.get("title_fit_avg_char_width_em", 0.56)),
                )
            )
        else:
            avg_char_width = float(header_config.get("title_fit_avg_char_width_em", 0.56))
        width_safety = float(header_config.get("title_fit_width_safety", 0.94))
        usable_width = max(width_inches * width_safety, 0.1)
        estimated_size = (usable_width * 72) / max(len(clean_text) * avg_char_width, 1)
        min_size = float(header_config.get(min_key, 34))
        return max(min_size, min(float(desired_size), estimated_size))

    def _render_section_title(self, slide, element: Dict, state: PosterState):
        """render section title with enhanced styling"""
        x, y, w, h = (float(element[k]) for k in ["x", "y", "width", "height"])
        
        section_title = repair_possessive_title_apostrophe(element.get("section_title", "").strip())
        if not section_title:
            return
        
        log_agent_info(self.name, f"rendering section title: '{section_title}'")

        style_enabled = self.visual_style_config.get("enabled", False)
        section_style = self.visual_style_config.get("section_title", {}) if style_enabled else {}
        section_title_font_size = self.styling_interfaces.get("font_sizes", {}).get("section_title", 48)
        font_size = self._fit_section_title_font_size(section_title, element, section_title_font_size)
        font_family = element.get("font_family") or section_style.get("font_family", "Georgia")
        font_weight = element.get("font_weight", section_style.get("font_weight", "bold"))
        font_color = element.get("font_color") or section_style.get("font_color", "#FFFFFF")
        alignment = element.get("alignment", section_style.get("alignment", "center")).lower()
        wordart_style = element.get("wordart_style") or {}
        shadow = wordart_style.get("shadow") or section_style.get("shadow")
        number_text = str(element.get("section_number") or "").strip()
        numbering_mode = str(element.get("section_numbering_mode") or "off").strip().lower()
        title_x, title_w = x, w
        number_box = None
        if numbering_mode == "small" and number_text:
            number_width = max(float(element.get("section_number_width", 0.46) or 0.46), 0.18)
            number_gap = max(float(element.get("section_number_gap", 0.12) or 0.12), 0.02)
            number_width = min(number_width, max(w * 0.24, 0.18))
            if w - number_width - number_gap >= 0.4:
                title_x = x + number_width + number_gap
                title_w = w - number_width - number_gap
                number_box = {
                    "x": x,
                    "y": y,
                    "w": number_width,
                    "h": h,
                    "font_size": max(22, int(font_size * float(element.get("section_number_font_scale", 0.62) or 0.62))),
                }
        font_size = self._fit_section_title_font_size(
            section_title,
            {**element, "width": title_w},
            font_size,
        )

        if shadow and shadow.get("enabled", True):
            if number_box:
                self._add_section_title_textbox(
                    slide,
                    number_text,
                    number_box["x"] + float(shadow.get("offset_x_inches", 0.025)),
                    number_box["y"] + float(shadow.get("offset_y_inches", 0.025)),
                    number_box["w"],
                    number_box["h"],
                    font_family,
                    number_box["font_size"],
                    font_weight,
                    shadow.get("color", "#AEB6D6"),
                    "right",
                )
            self._add_section_title_textbox(
                slide,
                section_title,
                title_x + float(shadow.get("offset_x_inches", 0.025)),
                y + float(shadow.get("offset_y_inches", 0.025)),
                title_w,
                h,
                font_family,
                font_size,
                font_weight,
                shadow.get("color", "#AEB6D6"),
                alignment,
            )

        if number_box:
            self._add_section_title_textbox(
                slide,
                number_text,
                number_box["x"],
                number_box["y"],
                number_box["w"],
                number_box["h"],
                font_family,
                number_box["font_size"],
                font_weight,
                font_color,
                "right",
            )
        self._add_section_title_textbox(
            slide,
            section_title,
            title_x,
            y,
            title_w,
            h,
            font_family,
            font_size,
            font_weight,
            font_color,
            alignment,
        )

    def _add_section_title_textbox(
        self,
        slide,
        section_title: str,
        x: float,
        y: float,
        w: float,
        h: float,
        font_family: str,
        font_size: int,
        font_weight: str,
        font_color: str,
        alignment: str,
    ):
        textbox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = textbox.text_frame
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.word_wrap = False
        tf.clear()
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        
        if len(tf.paragraphs) > 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = section_title
        p.font.name = font_family
        p.font.size = Pt(font_size)
        p.font.bold = font_weight == "bold"
        p.font.color.rgb = self._parse_color(font_color)
        if alignment == "center":
            p.alignment = PP_ALIGN.CENTER
        elif alignment == "right":
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        return textbox

    def _fit_section_title_font_size(self, title: str, element: Dict, fallback_size: int) -> int:
        font_size = int(element.get("font_size", fallback_size))
        width_inches = max(float(element.get("width", 0.0) or 0.0), 0.1)
        title_chars = max(len(title.strip()), 1)
        chars_per_inch = float(
            self.config.get("rendering", {}).get("section_title_chars_per_inch_at_48pt", 3.1)
        )
        capacity = width_inches * chars_per_inch * (48 / max(font_size, 1))
        if title_chars <= capacity:
            return font_size
        fitted = int(font_size * capacity / title_chars)
        return max(34, min(font_size, fitted))

    def _render_title_accent_block(self, slide, element: Dict, state: PosterState):
        """render color block accent for section titles"""
        x, y, w, h = (Inches(element[k]) for k in ["x", "y", "width", "height"])
        
        # use 'color' field from layout agent
        fill_color = element.get("color", element.get("fill_color", "#1E3A8A"))
        
        log_agent_info(self.name, f"rendering title accent block: {fill_color} at ({x.inches:.2f}, {y.inches:.2f})")
        
        self.director.add_shape(
            MSO_SHAPE.RECTANGLE,
            element["x"],
            element["y"],
            element["width"],
            element["height"],
            fill_color=fill_color,
            shadow=element.get("shadow"),
        )

    def _render_title_accent_line(self, slide, element: Dict, state: PosterState):
        """render underline accent for section titles"""
        x, y, w, h = (Inches(element[k]) for k in ["x", "y", "width", "height"])
        
        # use 'color' field from layout agent
        fill_color = element.get("color", element.get("fill_color", "#E8E8E8"))
        
        log_agent_info(self.name, f"rendering title accent line: {fill_color} at ({x.inches:.2f}, {y.inches:.2f})")
        
        self.director.add_shape(
            MSO_SHAPE.RECTANGLE,
            element["x"],
            element["y"],
            element["width"],
            element["height"],
            fill_color=fill_color,
            shadow=element.get("shadow"),
        )

    def _render_section_container(self, slide, element: Dict, state: PosterState):
        """render section container with optional debug border and mono_light background for critical sections"""
        x, y, w, h = (Inches(element[k]) for k in ["x", "y", "width", "height"])
        
        is_debug = element.get("debug_border", False)
        importance_level = element.get("importance_level", 2)
        template_meta = state.get("layout_template_metadata") or {}
        is_extracted_template = bool(template_meta.get("extracted_template"))
        
        # apply background fill based on importance level
        fill_color = element.get("fill_color")
        if not fill_color and importance_level == 1:
            # critical section gets mono_light background color
            color_scheme = state.get("color_scheme", {})
            fill_color = color_scheme.get("mono_light", "#e6eaef")
            log_agent_info(self.name, f"applied mono_light background ({fill_color}) to critical section")

        border_color = element.get("border_color")
        border_width = element.get("border_width", 1.0)
        border_style = element.get("border_style", "solid")
        if not is_debug and not is_extracted_template and not fill_color and not border_color:
            return

        # apply border based on debug mode
        if is_debug:
            border_color = "#FF0000"
            border_width = 2.0
            log_agent_info(self.name, f"added debug section border")

        if not fill_color and not border_color:
            return

        self.director.add_shape(
            MSO_SHAPE.RECTANGLE,
            element["x"],
            element["y"],
            element["width"],
            element["height"],
            fill_color=fill_color,
            fill_opacity=element.get("fill_opacity"),
            border_color=border_color,
            border_opacity=element.get("border_opacity"),
            border_width=border_width,
            border_style=border_style,
            shadow=element.get("shadow"),
        )

    def _render_text(self, slide, element: Dict, state: PosterState):
        """render text elements with enhanced formatting"""
        x, y, w, h = (Inches(element[k]) for k in ["x", "y", "width", "height"])
        
        content = element.get("content", "").strip()
        if not content:
            return
        
        log_agent_info(self.name, f"rendering text element: {element.get('id', 'unknown')}")
        
        # add text with margins
        margin = self.layout_constants["text_margin_renderer"]  # reduced margin for better space utilization
        self._add_enhanced_text(
            slide, content, 
            x + Inches(margin), y, 
            w - Inches(2 * margin), h,
            element
        )

    def _render_visual(self, slide, element: Dict, state: PosterState):
        """Render visual elements using only resolved visual assets."""
        visual_id = element.get("visual_id")
        slot_id = element.get("slot_id") or element.get("id")
        
        if visual_id and slot_id:
            self._add_resolved_visual(slot_id, visual_id, element, state)

    def _render_mixed(self, slide, element: Dict, state: PosterState):
        """render mixed elements (text and visual)"""
        # for now, treat as text element
        self._render_text(slide, element, state)

    def _render_conf_logo(self, slide, element: Dict, state: PosterState):
        """render conference logo"""
        logo_path = state.get("logo_path")
        if logo_path and Path(logo_path).exists():
            self._render_logo_with_aspect_ratio(slide, element, logo_path)

    def _render_aff_logo(self, slide, element: Dict, state: PosterState):
        """render affiliation logo"""
        aff_logo_path = state.get("aff_logo_path")
        if aff_logo_path and Path(aff_logo_path).exists():
            self._render_logo_with_aspect_ratio(slide, element, aff_logo_path)

    def _render_institution_logo(self, slide, element: Dict, state: PosterState):
        """render an auto-resolved institution logo."""
        logo_path = element.get("image_path")
        if logo_path and Path(logo_path).exists():
            self._render_logo_with_aspect_ratio(slide, element, logo_path)

    def _render_logo_divider(self, slide, element: Dict, state: PosterState):
        """Render a thin vertical rule between affiliation logos and conference logo."""
        color_scheme = state.get("color_scheme") or {}
        line_color = color_scheme.get("theme", "#AAAAAA")
        self.director.add_connector(
            element["x"], element["y"],
            element["x"], element["y"] + element["height"],
            color=line_color,
            width_pt=1.2,
        )

    def _add_enhanced_text(self, slide, text: str, left, top, width, height, element: Dict):
        """add text with enhanced formatting support"""
        if not text.strip():
            return
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.word_wrap = True
        tf.clear()
        
        # enforce height constraints to prevent text overflow beyond textbox bounds
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        # add small margins to ensure text stays within bounds
        tf.margin_top = Inches(0.05)
        tf.margin_bottom = Inches(0.05)
        
        # get font properties from element
        font_family = element.get("font_family", "Arial")
        font_size = element.get("font_size", 40)
        font_color = element.get("font_color", "#000000")
        line_spacing = element.get("line_spacing", self.styling_interfaces["line_spacing"])
        
        self._format_enhanced_text(tf, text, font_family, font_size, font_color, line_spacing)
        
        # debug info for formatting
        total_runs = sum(len(p.runs) for p in tf.paragraphs)
        log_agent_info(self.name, f"created {len(tf.paragraphs)} paragraphs with {total_runs} formatted runs")

    def _format_enhanced_text(self, text_frame, text: str, font_family: str, font_size: int, font_color: str, line_spacing: float):
        """format text with enhanced bullet point and bold support using 1.0 line spacing"""
        text_frame.clear()
        
        body_text_font_size = self.styling_interfaces.get("font_sizes", {}).get("body_text", 40)
        effective_font_size = font_size if font_size != 40 else body_text_font_size
        base_font_size = Pt(max(effective_font_size, 1))
        base_color = self._parse_color(font_color)
        
        # split by single newlines only (treat as simple line breaks)
        lines = text.split('\n')
        
        for line_idx, line in enumerate(lines):
            original_line = line  # keep original line for indentation detection
            line = line.strip()
            if not line:
                continue
            
            # create paragraph for each line
            if line_idx == 0 and len(text_frame.paragraphs) > 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            # handle indentation by checking if line starts with ◦ (sub-bullet)
            if line.strip().startswith(self.indentation_config["secondary_bullet_char"]):  # secondary bullet character
                # set paragraph level for indentation
                p.level = self.indentation_config["secondary_level"]
            else:
                p.level = self.indentation_config["primary_level"]
            
            # add formatted text content (don't clear p.text)
            self._add_formatted_runs(p, line, font_family, base_font_size, base_color)
            
            # set paragraph properties - force 1.0 line spacing
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = line_spacing
            p.space_before = Pt(0)
            p.space_after = Pt(0)

    def _add_formatted_runs(self, paragraph, text: str, font_family: str, 
                          base_font_size, base_color):
        """add text with all formatting as separate runs - following pptx best practices"""
        self._parse_and_add_runs(paragraph, text, font_family, base_font_size, base_color)

    def _parse_and_add_runs(self, paragraph, text: str, font_family: str, 
                          base_font_size, base_color):
        """parse text and create separate runs for each format type"""
        # tokenize the text into formatting segments
        segments = self._tokenize_formatting(text)
        
        # create runs for each segment
        for segment in segments:
            run = paragraph.add_run()
            run.text = segment['text']
            run.font.name = font_family
            run.font.size = base_font_size
            
            # apply formatting based on segment type
            if segment['color']:
                run.font.color.rgb = self._parse_color(segment['color'])
            else:
                run.font.color.rgb = base_color
            
            if segment['bold']:
                run.font.bold = True
            
            if segment['italic']:
                run.font.italic = True
    
    def _tokenize_formatting(self, text: str) -> list:
        """tokenize text into formatting segments with precise position tracking"""
        segments = []
        i = 0
        
        while i < len(text):
            # check for color markup: <color:#RRGGBB>text</color>
            color_match = re.match(r'<color:(#[0-9A-Fa-f]{6})>', text[i:])
            if color_match:
                color_hex = color_match.group(1)
                opening_tag_end = i + color_match.end()
                
                # find closing </color> tag using absolute position
                closing_tag_pattern = r'</color>'
                color_content_start = opening_tag_end
                closing_match = re.search(closing_tag_pattern, text[color_content_start:])
                
                if closing_match:
                    # calculate absolute positions
                    color_content_end = color_content_start + closing_match.start()
                    closing_tag_end = color_content_start + closing_match.end()
                    
                    # extract content between color tags
                    colored_text = text[color_content_start:color_content_end]
                    
                    # process colored text with automatic bold
                    if colored_text.strip():  # only process non-empty content
                        segments.append({
                            'text': colored_text,
                            'bold': True,  # all colored text is bold
                            'italic': False,
                            'color': color_hex
                        })
                    
                    # move past the entire color block
                    i = closing_tag_end
                    continue
                else:
                    # malformed color tag, treat as regular text
                    segments.append({
                        'text': text[i],
                        'bold': False,
                        'italic': False,
                        'color': None
                    })
                    i += 1
                    continue
            
            # check for bold: **text**
            bold_match = re.match(r'\*\*(.*?)\*\*', text[i:])
            if bold_match:
                bold_text = bold_match.group(1)
                segments.append({
                    'text': bold_text,
                    'bold': True,
                    'italic': False,
                    'color': None
                })
                i += bold_match.end()
                continue
            
            # check for italic: *text*
            italic_match = re.match(r'\*(.*?)\*', text[i:])
            if italic_match:
                italic_text = italic_match.group(1)
                segments.append({
                    'text': italic_text,
                    'bold': False,
                    'italic': True,
                    'color': None
                })
                i += italic_match.end()
                continue
            
            # regular text - find next formatting marker
            next_format = re.search(r'(\*\*|\*|<color:)', text[i:])
            if next_format:
                regular_text = text[i:i + next_format.start()]
            else:
                regular_text = text[i:]
            
            if regular_text:
                segments.append({
                    'text': regular_text,
                    'bold': False,
                    'italic': False,
                    'color': None
                })
            
            if next_format:
                if next_format.start() == 0:
                    marker = next_format.group(1)
                    segments.append({
                        'text': marker,
                        'bold': False,
                        'italic': False,
                        'color': None
                    })
                    i += len(marker)
                else:
                    i += next_format.start()
            else:
                break
        
        return segments
    
    def _parse_bold_italic(self, text: str, color: str) -> list:
        """simplified bold/italic parser - only used for nested formatting"""
        segments = []
        i = 0
        
        while i < len(text):
            # check for bold
            bold_match = re.match(r'\*\*(.*?)\*\*', text[i:])
            if bold_match:
                bold_text = bold_match.group(1)
                segments.append({
                    'text': bold_text,
                    'bold': True,
                    'italic': False,
                    'color': color
                })
                i += bold_match.end()
                continue
            
            # check for italic
            italic_match = re.match(r'\*(.*?)\*', text[i:])
            if italic_match:
                italic_text = italic_match.group(1)
                segments.append({
                    'text': italic_text,
                    'bold': bool(color),  # force bold if color is present
                    'italic': True,
                    'color': color
                })
                i += italic_match.end()
                continue
            
            # regular text
            next_format = re.search(r'(\*\*|\*)', text[i:])
            if next_format:
                regular_text = text[i:i + next_format.start()]
            else:
                regular_text = text[i:]
            
            if regular_text:
                segments.append({
                    'text': regular_text,
                    'bold': bool(color),  # force bold if color is present
                    'italic': False,
                    'color': color
                })
            
            if next_format:
                i += next_format.start()
            else:
                break
        
        return segments

    def _add_resolved_visual(self, slot_id: str, visual_id: str, element: Dict[str, Any], state: PosterState):
        """Render a resolved visual slot without any image decision logic."""
        resolved_entry = self._get_resolved_visual_entry(slot_id, visual_id, state)
        if not resolved_entry:
            raise ValueError(f"missing resolved visual asset for slot '{slot_id}'")

        visual_path = resolved_entry.get("resolved_path")
        if not visual_path or not Path(visual_path).exists():
            raise ValueError(f"resolved visual path not found for slot '{slot_id}': {visual_path}")

        try:
            self.director.add_image(
                visual_path,
                element["x"],
                element["y"],
                element["width"],
                element["height"],
                keep_aspect_ratio=False,
            )
        except Exception as e:
            log_agent_error(self.name, f"failed to render visual slot {slot_id}: {e}")
            state.setdefault("errors", []).append(f"{self.name}: failed to render visual slot {slot_id}: {e}")

    def _render_logo_with_aspect_ratio(self, slide, element: Dict, image_path: str):
        """render logo with proper aspect ratio preservation"""
        x, y, w, h = (Inches(element[k]) for k in ["x", "y", "width", "height"])
        
        try:
            render_path = self._trim_logo_whitespace(image_path)
            self.director.add_image(render_path, element["x"], element["y"], element["width"], element["height"], keep_aspect_ratio=True)
                                   
        except Exception as e:
            log_agent_error(self.name, f"failed to render logo: {e}")

    def _trim_logo_whitespace(self, image_path: str) -> str:
        try:
            from PIL import Image, ImageChops

            source = Path(image_path)
            with Image.open(source).convert("RGBA") as image:
                white = Image.new("RGBA", image.size, (255, 255, 255, 255))
                diff = ImageChops.difference(image, white).convert("L")
                mask = diff.point(lambda value: 255 if value > 12 else 0)
                bbox = mask.getbbox()
                if not bbox:
                    return image_path
                content_area = (bbox[2] - bbox[0]) * (bbox[3] - bbox[1])
                total_area = image.size[0] * image.size[1]
                if total_area <= 0 or content_area / total_area > 0.86:
                    return image_path

                pad_x = max(int((bbox[2] - bbox[0]) * 0.04), 8)
                pad_y = max(int((bbox[3] - bbox[1]) * 0.04), 8)
                crop_box = (
                    max(bbox[0] - pad_x, 0),
                    max(bbox[1] - pad_y, 0),
                    min(bbox[2] + pad_x, image.size[0]),
                    min(bbox[3] + pad_y, image.size[1]),
                )
                output_dir = getattr(self, "_render_output_dir", source.parent) / "assets"
                output_dir.mkdir(parents=True, exist_ok=True)
                digest = hashlib.sha1(f"{source.resolve()}:{crop_box}".encode("utf-8")).hexdigest()[:10]
                output_path = output_dir / f"{source.stem}_trimmed_{digest}.png"
                if not output_path.exists():
                    image.crop(crop_box).save(output_path)
                return str(output_path)
        except Exception:
            return image_path
        return image_path

    def _get_resolved_visual_entry(self, slot_id: str, visual_id: str, state: PosterState) -> Optional[Dict[str, Any]]:
        """Resolve a slot-level visual entry from the visual asset agent output."""
        resolved_visual_assets = state.get("resolved_visual_assets") or {}
        if slot_id in resolved_visual_assets:
            return resolved_visual_assets[slot_id]
        return resolved_visual_assets.get(visual_id)

    def _parse_color(self, color_str: str) -> RGBColor:
        """parse color string to RGBColor"""
        hex_color = color_str.lstrip('#')
        r, g, b = (int(hex_color[i:i+2], 16) for i in (0, 2, 4))
        return RGBColor(r, g, b)

    def _generate_qr_code(self, url: str, output_dir: str) -> str:
        """generate QR code for URL"""
        qr = qrcode.QRCode(
            version=1,
            error_correction=qrcode.constants.ERROR_CORRECT_L,
            box_size=10,
            border=2,
        )
        qr.add_data(url)
        qr.make(fit=True)
        
        img = qr.make_image(fill_color="black", back_color="white")
        qr_path = Path(output_dir) / "qr_code.png"
        img.save(qr_path)
        
        return str(qr_path)

    def _render_qr_code(self, slide, element: Dict, qr_code_path: str):
        """render QR code element"""
        x, y, w, h = (Inches(element[k]) for k in ["x", "y", "width", "height"])
        slide.shapes.add_picture(qr_code_path, x, y, w, h)

    def _convert_to_png(self, pptx_path: Path) -> Optional[str]:
        """Convert PPTX to PNG using LibreOffice, with macOS QuickLook as fallback."""
        try:
            output_dir = pptx_path.parent
            
            import platform
            system = platform.system().lower()
            
            if system == "windows":
                libreoffice_paths = [
                    r"C:\Program Files\LibreOffice\program\soffice.exe",
                    r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
                    r"C:\Users\%USERNAME%\AppData\Local\Programs\LibreOffice\program\soffice.exe",
                    "soffice.exe",
                    "libreoffice.exe"
                ]
            elif system == "linux":
                libreoffice_paths = [
                    "/usr/bin/libreoffice",
                    "/usr/local/bin/libreoffice",
                    "/snap/bin/libreoffice",
                    "/usr/bin/soffice",
                    "libreoffice",
                    "soffice"
                ]
            elif system == "darwin":  # macOS
                libreoffice_paths = [
                    "/Applications/LibreOffice.app/Contents/MacOS/soffice",
                    "/usr/local/bin/libreoffice",
                    "libreoffice",
                    "soffice"
                ]
            else:
                libreoffice_paths = [
                    "libreoffice",
                    "soffice"
                ]
            
            for lo_path in libreoffice_paths:
                try:
                    cmd = [
                        lo_path, "--headless", "--convert-to", "png",
                        "--outdir", str(output_dir), str(pptx_path)
                    ]
                    
                    result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
                    
                    if result.returncode == 0:
                        png_name = pptx_path.stem + ".png"
                        png_path = output_dir / png_name
                        if png_path.exists():
                            return str(png_path)
                            
                except (subprocess.SubprocessError, FileNotFoundError):
                    continue
            
            log_agent_error(self.name, "LibreOffice not found - install for PNG conversion")
            
        except Exception as e:
            log_agent_error(self.name, f"PNG conversion failed: {e}")

        return self._convert_to_png_with_quicklook(pptx_path)

    def _convert_to_png_with_quicklook(self, pptx_path: Path) -> Optional[str]:
        """Generate a preview image on macOS when LibreOffice is unavailable."""
        if shutil.which("qlmanage") is None:
            return None

        preview_dir = pptx_path.parent / "preview"
        preview_dir.mkdir(parents=True, exist_ok=True)
        try:
            result = subprocess.run(
                [
                    "qlmanage",
                    "-t",
                    "-s",
                    "2000",
                    "-o",
                    str(preview_dir),
                    str(pptx_path),
                ],
                capture_output=True,
                text=True,
                timeout=60,
            )
            if result.returncode != 0:
                log_agent_error(self.name, f"QuickLook preview failed: {result.stderr.strip()}")
                return None

            candidates = sorted(preview_dir.glob(f"{pptx_path.name}*.png"))
            if not candidates:
                candidates = sorted(preview_dir.glob("*.png"))
            if not candidates:
                return None

            final_path = pptx_path.with_suffix(".png")
            candidates[0].replace(final_path)
            return str(final_path)
        except Exception as e:
            log_agent_error(self.name, f"QuickLook preview failed: {e}")
            return None


def renderer_node(state: PosterState) -> Dict[str, Any]:
    result = Renderer()(state)
    return {
        **state,
        "pptx_output_path": result.get("pptx_output_path"),
        "poster_preview_path": result.get("poster_preview_path"),
        "render_stage": result.get("render_stage", state.get("render_stage")),
        "draft_status": result.get("draft_status", state.get("draft_status")),
        "final_poster_accepted": result.get("final_poster_accepted", state.get("final_poster_accepted", False)),
        "tokens": result["tokens"],
        "current_agent": result["current_agent"],
        "errors": result["errors"]
    }
